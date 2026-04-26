"""Major v3 rebuild of the pre-defense deck.

Honours the user's full feedback list (course code, 3-slide intro with
1.1 dark, motivation merged into intro, no em-dashes / double dashes
anywhere in body text, real Pareto with two stars, AR0144 hardware,
RTX 3050 + Jetson Orin Nano latency listed everywhere, methodology
section, stacked architecture sub-slides, working principle slide
(dataset + training), results split into Scene Flow + custom + live
demo + 3D reconstruction, rewritten discussion / challenges / impact
/ conclusion / future work).

Strategy
--------
Open the original v1 .pptx, treat slides 1, 5, 21, 24 as templates
(title, dark divider, big-figure, two-card-cream), and rebuild the
deck slide-by-slide by duplicating the right template, clearing its
body, and inserting fresh content. Final reorder + page renumber.
"""
from __future__ import annotations

import copy
import shutil
import subprocess
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path("/home/abrar/Research/stero_research_claude")
ORIG = ROOT / "presentation/Thesis Pre Defense Presentation Slides(2008011, 2008026).pptx"
OUT  = ROOT / "presentation/Thesis Pre Defense Presentation Slides v3 (2008011, 2008026).pptx"

FIGS    = ROOT / "presentation/figs"
PCTOP   = Path("/mnt/abrarssd/Datasets/stereo_samples_20260425_104147/point_clouds_top3")
VIDS    = Path("/mnt/abrarssd/Datasets/stereo_samples_20260425_104147/vids")
DEPLOY  = ROOT / "model/designs/d1_tile/deployment_pipeline.png"
ARCH_FULL = ROOT / "model/designs/d1_tile/stereolite_arch.png"
TRAIN_GIF = ROOT / "model/benchmarks/stereolite_finetune_indoor_20260426-171158/training_progression.gif"
TRAIN_FULL_GIF = ROOT / "model/benchmarks/stereolite_finetune_indoor_20260426-171158/training_progression_full_4x2.gif"
SF_PROGRESSION = ROOT / "model/designs/d1_tile/progress_grid.png"
SF_TRAINING_CURVES = ROOT / "model/designs/d1_tile/training_curves.png"

INK_HEX     = "1A1A1F"
SUBINK_HEX  = "5A5550"
ACCENT_HEX  = "C24A1C"
SOFTACC_HEX = "D9826A"
CREAM_HEX   = "F4EFE6"
DARK_BG_HEX = "242021"
WHITE_HEX   = "F5EFE3"

VIDEO_60S = VIDS / "stereolite_inference_panel_60s_1080.mp4"

# Original v1 slide indices that we will use as templates
T_TITLE   = 0     # the title slide
T_OUTLINE = 1     # outline
T_DARK    = 4     # Problem Statement (dark bg, big section number)
T_CREAM_BIG_FIG = 20  # training-curves slide (cream + big figure + side text)
T_CREAM_2COL    = 23  # discussion (cream + 2 boxes)
T_CREAM_OBJ     = 5   # objective (cream + bullet bar + 3 boxes)
T_QA      = 27    # Q&A

# --------------------------------------------------------------------------
# Low-level helpers
# --------------------------------------------------------------------------

def replace_in_runs(text_frame, old, new):
    for p in text_frame.paragraphs:
        runs = list(p.runs)
        if not runs: continue
        full = "".join(r.text for r in runs)
        if old not in full: continue
        runs[0].text = full.replace(old, new)
        for r in runs[1:]:
            r.text = ""


def deep_replace(shape, old, new):
    if shape.has_text_frame:
        replace_in_runs(shape.text_frame, old, new)
    if shape.shape_type == 6:
        for sub in shape.shapes:
            deep_replace(sub, old, new)


def set_paragraph_text(text_frame, new_text):
    paras = list(text_frame.paragraphs)
    if not paras:
        text_frame.text = new_text
        return
    runs = list(paras[0].runs)
    if runs:
        runs[0].text = new_text
        for r in runs[1:]: r.text = ""
    else:
        paras[0].add_run().text = new_text
    for p in paras[1:]:
        for r in list(p.runs): r.text = ""


def remove_shape(shape):
    shape.element.getparent().remove(shape.element)


def duplicate_slide(prs, src_idx):
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)

    # Copy background fill from source <p:cSld> if it has one
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    src_csld = src.element.find(f"{{{p_ns}}}cSld")
    new_csld = new.element.find(f"{{{p_ns}}}cSld")
    if src_csld is not None and new_csld is not None:
        src_bg = src_csld.find(f"{{{p_ns}}}bg")
        if src_bg is not None:
            new_bg = copy.deepcopy(src_bg)
            # <p:bg> must be the FIRST child of <p:cSld>
            new_csld.insert(0, new_bg)

    # Copy shapes
    for shape in list(src.shapes):
        new_el = copy.deepcopy(shape.element)
        new.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Re-create relationships for embedded pictures
    rels_src = src.part.rels
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for blip in new.shapes._spTree.findall(f".//{{{a_ns}}}blip"):
        rid = blip.get(f"{{{r_ns}}}embed")
        if rid and rid in rels_src:
            rel = rels_src[rid]
            new_rid = new.part.relate_to(rel.target_part, rel.reltype)
            blip.set(f"{{{r_ns}}}embed", new_rid)
    return new


def move_slide(prs, slide, new_idx):
    xml = prs.slides._sldIdLst
    el = next(s for s in xml if int(s.attrib["id"]) == slide.slide_id)
    xml.remove(el); xml.insert(new_idx, el)


def remove_slide_at(prs, idx):
    xml = prs.slides._sldIdLst
    sl = list(xml)[idx]
    xml.remove(sl)


# --------------------------------------------------------------------------
# Text-box helpers
# --------------------------------------------------------------------------

def add_text(slide, x, y, w, h, text, *, size=12, bold=False,
             color=INK_HEX, mono=False, italic=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                     Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center": p.alignment = PP_ALIGN.CENTER
    elif align == "right": p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor.from_string(color)
    r.font.name = "DejaVu Sans Mono" if mono else "DejaVu Serif"
    return box


def add_bullet_lines(slide, x, y, w, h, lines, *, size=11,
                      color=INK_HEX, bullet="·"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                     Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        r = p.add_run()
        r.text = f"{bullet}  {line}" if bullet else line
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor.from_string(color)
        r.font.name = "DejaVu Serif"
        p.space_after = Pt(2)
    return box


# --------------------------------------------------------------------------
# Slide builders that re-use templates
# --------------------------------------------------------------------------

def build_dark_section_slide(prs, *, section_label, title_main, title_accent="",
                              big_number="", body_lines=None, italic_summary=""):
    """Duplicate the dark Problem-Statement template and re-fill content."""
    new = duplicate_slide(prs, T_DARK)

    # The dark template has:
    #   sh[0]: header label  '§ 03 PROBLEM STATEMENT'
    #   sh[1]: big title     'How do we recover dense ...'
    #   sh[2..3]: footer band/text/page (kept)
    #   one body paragraph beneath the title (~0.6 in below)
    #   big '03' number on right
    #
    # Strategy: re-set known text shapes, remove body bullet text shape
    # if present, then add new body lines.
    body_shapes_to_remove = []
    title_shape = None
    header_shape = None
    big_num_shape = None
    body_text_shape = None
    for sh in new.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text.strip()
        if "§ 03 PROBLEM STATEMENT" in t:
            header_shape = sh
        elif "How do we recover" in t:
            title_shape = sh
        elif t == "03":
            big_num_shape = sh
        elif t.startswith("Existing high-accuracy"):
            body_text_shape = sh

    if header_shape is not None:
        set_paragraph_text(header_shape.text_frame, section_label)
    if title_shape is not None:
        # rewrite as two runs (main + accent)
        tf = title_shape.text_frame
        for p in list(tf.paragraphs):
            for r in list(p.runs):
                p._p.remove(r._r)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = title_main + (" " if title_accent else "")
        r1.font.bold = True; r1.font.size = Pt(32)
        r1.font.color.rgb = RGBColor.from_string(WHITE_HEX)
        r1.font.name = "DejaVu Serif"
        if title_accent:
            r2 = p.add_run(); r2.text = title_accent
            r2.font.bold = True; r2.font.size = Pt(32)
            r2.font.color.rgb = RGBColor.from_string(SOFTACC_HEX)
            r2.font.name = "DejaVu Serif"
    if big_num_shape is not None and big_number:
        set_paragraph_text(big_num_shape.text_frame, big_number)
    if body_text_shape is not None:
        remove_shape(body_text_shape)
    if body_lines:
        # Wider, taller body region than the template default.  The big
        # decorative number sits on the right so we cap width at 6.0 in.
        box = new.shapes.add_textbox(
            Inches(0.55), Inches(1.95), Inches(6.20), Inches(3.10))
        tf = box.text_frame
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.word_wrap = True
        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            r = p.add_run(); r.text = line
            r.font.size = Pt(11)
            r.font.color.rgb = RGBColor.from_string(WHITE_HEX)
            r.font.name = "DejaVu Serif"
            p.space_after = Pt(4)
        if italic_summary:
            p = tf.add_paragraph(); p.space_before = Pt(6)
            r = p.add_run(); r.text = italic_summary
            r.font.size = Pt(10.5)
            r.font.italic = True
            r.font.color.rgb = RGBColor.from_string(SOFTACC_HEX)
            r.font.name = "DejaVu Serif"
    return new


def build_cream_chrome_slide(prs, *, section_label, title_main,
                               title_accent="", template=T_CREAM_2COL):
    """Duplicate a cream template, rewrite header + title, clear body."""
    new = duplicate_slide(prs, template)

    # Identify keep-shapes (header, title, footer)
    header_shape = title_shape = None
    keep = []
    for sh in new.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text.strip()
        if t.startswith("§ "):
            header_shape = sh
            keep.append(sh)
            continue
        if "Reading the numbers" in t or "Loss drops" in t:
            title_shape = sh
            keep.append(sh)
            continue
        # Footer band stays (top > 5.2in)
        if sh.top is not None and sh.top >= Inches(5.20):
            keep.append(sh)
            continue
        if t.isdigit() and len(t) <= 2:
            keep.append(sh); continue
        if "AI-Enhanced" in t or "APRIL" in t:
            keep.append(sh); continue

    if header_shape is not None:
        set_paragraph_text(header_shape.text_frame, section_label)
    if title_shape is not None:
        tf = title_shape.text_frame
        for p in list(tf.paragraphs):
            for r in list(p.runs):
                p._p.remove(r._r)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = title_main + (" " if title_accent else "")
        r1.font.bold = True; r1.font.size = Pt(28)
        r1.font.color.rgb = RGBColor.from_string(INK_HEX)
        r1.font.name = "DejaVu Serif"
        if title_accent:
            r2 = p.add_run(); r2.text = title_accent
            r2.font.bold = True; r2.font.size = Pt(28)
            r2.font.color.rgb = RGBColor.from_string(ACCENT_HEX)
            r2.font.name = "DejaVu Serif"

    # Drop everything in the body region
    for sh in list(new.shapes):
        if sh in keep: continue
        if sh.top is not None and Inches(1.20) < sh.top < Inches(5.20):
            remove_shape(sh)
            continue
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t.isdigit() and len(t) <= 2: continue
            if "AI-Enhanced" in t or "APRIL" in t: continue
            if t.startswith("§ "): continue
            if title_shape is not None and sh is title_shape: continue
            remove_shape(sh)
    return new


# --------------------------------------------------------------------------
# Per-slide builders
# --------------------------------------------------------------------------

def edit_title_slide(prs):
    s = prs.slides[T_TITLE]
    for sh in s.shapes:
        deep_replace(sh, "MTE 4200", "MTE 4210")
        deep_replace(sh, "Project and Thesis", "Seminar")


def build_outline_slide(prs):
    """Clear the entire body of the outline slide and rebuild from
    scratch with two columns of section/title/page rows."""
    s = prs.slides[1]
    # Drop everything except header label, big title, and footer chrome
    for sh in list(s.shapes):
        if not sh.has_text_frame:
            # Remove decorative non-text shapes inside the body
            if sh.top is not None and Inches(1.10) < sh.top < Inches(5.20):
                remove_shape(sh)
            continue
        t = sh.text_frame.text.strip()
        if t.startswith("§ 00"): continue
        if t == "Outline.": continue
        if "AI-Enhanced" in t: continue
        if "APRIL" in t: continue
        # Footer page number
        if sh.top is not None and sh.top > Inches(5.20):
            continue
        remove_shape(sh)

    rows = [
        ("01", "Introduction",            "03"),
        ("02", "Problem Statement",       "06"),
        ("03", "Objectives",              "07"),
        ("04", "Literature Review",       "08"),
        ("04.1", "Research Gap",          "09"),
        ("05", "Proposed Solution",       "10"),
        ("06", "Methodology",             "11"),
        ("07", "Implementation",          "12"),
        ("07.1", "Architecture",          "13 to 16"),
        ("08", "Working Principle",       "17"),
        ("09", "Results & Analysis",      "18"),
        ("09.1", "Scene Flow",            "19"),
        ("09.2", "Indoor Real Data",      "20"),
        ("09.3", "Live Inference",        "21"),
        ("09.4", "3D Reconstruction",     "22"),
        ("10", "Discussion",              "23"),
        ("11", "Challenges & Limitations","24"),
        ("12", "Impact & Contributions",  "25"),
        ("13", "Conclusion",              "26"),
        ("14", "Future Work",             "27"),
    ]
    # Two columns, top-to-bottom split at the halfway row
    half = (len(rows) + 1) // 2
    left_rows = rows[:half]
    right_rows = rows[half:]
    y0 = 1.40
    row_h = 0.20
    for col_x, col_rows in [(0.55, left_rows), (5.10, right_rows)]:
        for i, (num, label, page) in enumerate(col_rows):
            y = y0 + i * row_h
            # Number (mono, accent)
            add_text(s, col_x, y, 0.55, row_h, num,
                      size=8.5, color=ACCENT_HEX, mono=True, bold=True)
            # Label
            add_text(s, col_x + 0.55, y, 3.20, row_h, label,
                      size=10.5, color=INK_HEX)
            # Page number on the right
            add_text(s, col_x + 3.85, y, 0.50, row_h, page,
                      size=8.5, color=SUBINK_HEX, mono=True,
                      align="right")
            # Underline (thin rule below each row)
            from pptx.enum.shapes import MSO_SHAPE
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(col_x), Inches(y + row_h - 0.02),
                Inches(4.30), Inches(0.005))
            line.line.fill.background()
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor.from_string("D8CFC0")


def build_intro_1_1(prs):
    """DARK section opener for Introduction. Generic intro that pulls
    motivation language in: depth perception is the front door to
    autonomy; current heavy networks won't fit on edge devices."""
    return build_dark_section_slide(prs,
        section_label="§ 01.1 INTRODUCTION",
        title_main="Depth perception is the front door to",
        title_accent="autonomy.",
        big_number="01",
        body_lines=[
            "Self-driving cars, factory robots, indoor delivery drones, "
            "and AR headsets all need a fast, accurate idea of how far "
            "objects are.",
            "Stereo cameras give that depth from two passive lenses, "
            "with no active emitter, no eye-safety limits, and almost "
            "no extra cost over a single camera.",
            "Yet the deep networks that win the public stereo "
            "leaderboards run hundreds of millions of parameters; far "
            "too heavy for the small computers on those moving "
            "platforms.",
        ],
        italic_summary="This thesis builds StereoLite, a 0.87 M parameter "
                       "stereo network designed to close that gap.")


def build_intro_1_2(prs):
    """Cream slide. Stereo principles, what the model predicts."""
    s = build_cream_chrome_slide(prs,
        section_label="§ 01.2 INTRODUCTION",
        title_main="Stereo vision",
        title_accent="in one slide.")
    # Left: paragraph
    add_bullet_lines(s, 0.55, 1.40, 4.80, 2.40, [
        "Two cameras separated by a fixed baseline see the same scene "
        "from offset viewpoints.",
        "For every pixel in the left view, we find its match in the "
        "right view; the horizontal pixel shift is the disparity d.",
        "Once we know d at every pixel and the camera geometry, depth "
        "Z follows directly from a single division.",
    ], size=11)

    # Right: formula card
    from pptx.enum.shapes import MSO_SHAPE
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(5.55), Inches(1.40), Inches(3.95), Inches(2.40))
    box.line.color.rgb = RGBColor.from_string(SUBINK_HEX)
    box.line.width = Pt(0.5)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string("FBF7EE")
    add_text(s, 5.70, 1.50, 3.60, 0.30,
              "FORMULA  ·  DEPTH FROM DISPARITY",
              size=8.5, bold=True, mono=True, color=ACCENT_HEX)
    add_text(s, 5.70, 1.95, 3.60, 0.80,
              "Z  =  (f · B) / d",
              size=24, bold=True, color=INK_HEX, align="center")
    add_bullet_lines(s, 5.70, 2.85, 3.60, 0.80, [
        "Z  ·  metric depth (m)",
        "f  ·  focal length (px)",
        "B  ·  stereo baseline (m)",
        "d  ·  predicted disparity (px)",
    ], size=9, color=SUBINK_HEX, bullet="")

    add_text(s, 0.55, 4.15, 8.95, 0.35,
              "We predict d. Z is a one line geometric "
              "post processing step.",
              size=11, italic=True, color=ACCENT_HEX)
    return s


def build_intro_1_3(prs):
    """Cream slide. Why edge stereo is hard (was the old Motivation)."""
    s = build_cream_chrome_slide(prs,
        section_label="§ 01.3 INTRODUCTION",
        title_main="Why edge stereo",
        title_accent="is hard.")
    # Three big-number columns
    cols_x = [0.55, 3.70, 6.85]
    titles = ["HEAVY NETWORKS", "EDGE BUDGET", "CLASSICAL METHODS"]
    big = ["30 to 80 M", "under 3 M", "over 5.0"]
    units = ["params", "params", "px EPE"]
    notes = [
        "PSMNet, RAFT-Stereo, CREStereo run on desktop GPUs. They do "
        "not fit on Jetson class devices.",
        "Jetson Orin Nano needs compact INT8 / FP16 friendly models "
        "for real time inference.",
        "Hand crafted SGM and block matching fit on edge hardware but "
        "fail on textureless and reflective regions.",
    ]
    for x, t, b, u, note in zip(cols_x, titles, big, units, notes):
        add_text(s, x, 1.50, 2.95, 0.25, t,
                  size=8.5, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x, 1.85, 2.95, 0.80, b,
                  size=28, bold=True, color=INK_HEX)
        add_text(s, x + 2.10, 2.20, 0.85, 0.30, u,
                  size=11, color=SUBINK_HEX)
        add_text(s, x, 2.95, 2.95, 1.30, note,
                  size=10, color=SUBINK_HEX)

    add_text(s, 0.55, 4.55, 8.95, 0.35,
              "Robotics, ADAS, and low power autonomy all need real "
              "time depth on hardware that cannot run the state of the "
              "art.",
              size=11, italic=True, color=ACCENT_HEX)
    return s


def build_problem(prs):
    return build_dark_section_slide(prs,
        section_label="§ 02 PROBLEM STATEMENT",
        title_main="Real time depth on edge devices,",
        title_accent="without losing accuracy.",
        big_number="02",
        body_lines=[
            "Indoor combat drones, factory autonomous vehicles, "
            "warehouse robots, and small ground rovers all need on "
            "board depth perception for obstacle avoidance, free space "
            "navigation, and visual odometry.",
            "Their compute budgets sit between 5 and 25 watts; the "
            "high accuracy stereo networks need ten times that.",
            "Existing lightweight pipelines run fast on those budgets "
            "but their accuracy collapses on textureless walls, "
            "reflective floors, and dense foliage that real platforms "
            "see every day.",
        ],
        italic_summary="How do we deliver dense, accurate depth at "
                       "edge scale compute, on cameras with no "
                       "calibration luxury?")


def build_objectives(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 03 OBJECTIVES",
        title_main="Objectives",
        title_accent="and success criteria.")

    # Top: original two-bullet box, kept as the user asked
    from pptx.enum.shapes import MSO_SHAPE
    box = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.55), Inches(1.40), Inches(8.95), Inches(1.30))
    box.line.color.rgb = RGBColor.from_string(SUBINK_HEX)
    box.line.width = Pt(0.5)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string("FBF7EE")
    # accent bar on the left
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(0.55), Inches(1.40), Inches(0.06), Inches(1.30))
    bar.line.fill.background()
    bar.fill.solid(); bar.fill.fore_color.rgb = RGBColor.from_string(ACCENT_HEX)

    add_bullet_lines(s, 0.85, 1.55, 8.50, 1.10, [
        "Design a computationally efficient stereo matching pipeline "
        "that uses AI based disparity refinement on resource limited "
        "platforms.",
        "Design an architecture that can withstand camera "
        "imperfections such as missing rectification on cheap "
        "stereo modules.",
    ], size=11)

    # Three cards
    cols_x = [0.55, 3.70, 6.85]
    titles = ["EFFICIENT", "AI REFINED", "DEPLOYABLE"]
    body = [
        ("Trainable params under 1 M.\n"
         "Inference under 40 ms per 512x832 stereo pair on Jetson "
         "Orin Nano (INT8).\n"
         "Checkpoint under 20 MB."),
        ("Iterative deep refinement on top of one coarse cost volume.\n"
         "Tile hypothesis state captures slanted surfaces."),
        ("End to end trainable with multi scale supervision.\n"
         "Camera imperfection tolerant: works on raw, "
         "non rectified stereo input from low cost USB cameras."),
    ]
    for x, t, b in zip(cols_x, titles, body):
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(2.85), Inches(2.95), Inches(2.05))
        card.line.color.rgb = RGBColor.from_string(SUBINK_HEX)
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string("FBF7EE")
        add_text(s, x + 0.15, 3.00, 2.65, 0.25, t,
                  size=8.5, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x + 0.15, 3.30, 2.65, 1.55, b,
                  size=10, color=INK_HEX)
    return s


def build_literature(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 04 LITERATURE REVIEW",
        title_main="110 papers analyzed,",
        title_accent="five that shaped StereoLite.")

    # Sub-headline
    add_text(s, 0.55, 1.40, 8.95, 0.35,
              "Across the 2002 to 2026 stereo matching literature, the "
              "designs below influenced our architecture the most.",
              size=10.5, color=SUBINK_HEX)

    # Build a simple table-style 5-row layout (no python-pptx Table to
    # keep styling consistent with the rest of the deck)
    headers = ["METHOD", "CORE IDEA", "PARAMS", "ROLE IN STEREOLITE"]
    rows = [
        ("GC-Net  (Kendall et al., 2017)",
         "First end to end 3D cost volume regression",
         "3.5 M",
         "Established the cost volume backbone."),
        ("PSMNet  (Chang & Chen, 2018)",
         "Pyramid pooling + stacked hourglass 3D CNN",
         "5.2 M",
         "Strong baseline; 3D conv unfit for edge."),
        ("HITNet  (Tankovich et al., 2021)",
         "Tile hypothesis planes (d, slope x, slope y, feature, conf)",
         "0.6 M",
         "Inspired the slanted tile state."),
        ("RAFT-Stereo  (Lipson et al., 2021)",
         "Recurrent residual refinement",
         "11 M",
         "Inspired iterative updates and convex upsample."),
        ("MobileStereoNet  (Shamsafar et al., 2022)",
         "MobileNet backbone for stereo",
         "2.4 M",
         "Edge oriented, weaker boundary accuracy."),
    ]
    col_x = [0.55, 2.85, 6.05, 6.95]
    col_w = [2.20, 3.10, 0.80, 2.55]
    y0 = 1.85
    row_h = 0.55

    # header row
    for cx, cw, h in zip(col_x, col_w, headers):
        add_text(s, cx, y0, cw, 0.25, h,
                  size=8.5, bold=True, mono=True, color=ACCENT_HEX)
    y = y0 + 0.35
    for i, row in enumerate(rows):
        for cx, cw, val in zip(col_x, col_w, row):
            add_text(s, cx, y, cw, row_h, val, size=10,
                      color=INK_HEX if cx == col_x[0] else SUBINK_HEX,
                      bold=(cx == col_x[0]))
        y += row_h

    add_text(s, 0.55, y + 0.05, 8.95, 0.35,
              "Compression and edge oriented variants of these "
              "designs informed our parameter budget choices.",
              size=10, italic=True, color=ACCENT_HEX)
    return s


def build_research_gap(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 04.1 RESEARCH GAP",
        title_main="A gap in the",
        title_accent="accuracy x size plane.")

    s.shapes.add_picture(
        str(FIGS / "research_gap_pareto.png"),
        Inches(0.40), Inches(1.30),
        width=Inches(5.50))

    rx = 6.10
    add_text(s, rx, 1.45, 3.50, 0.30,
              "FOUNDATION + ITERATIVE",
              size=8.5, bold=True, mono=True, color=ACCENT_HEX)
    add_text(s, rx, 1.75, 3.50, 0.55,
              "Highly accurate, but 10 to 350 M parameters.\n"
              "Desktop GPU only.",
              size=10, color=SUBINK_HEX)

    add_text(s, rx, 2.45, 3.50, 0.30,
              "EFFICIENT (PRE 2024)",
              size=8.5, bold=True, mono=True, color=ACCENT_HEX)
    add_text(s, rx, 2.75, 3.50, 0.55,
              "Edge friendly sizes, but EPE often above 1 px.",
              size=10, color=SUBINK_HEX)

    add_text(s, rx, 3.30, 3.50, 0.30,
              "TARGET ZONE",
              size=8.5, bold=True, mono=True, color=ACCENT_HEX)
    add_text(s, rx, 3.60, 3.50, 0.85,
              "Under 1 M parameters AND competitive EPE. "
              "StereoLite is the only point inside the dashed box.",
              size=10, color=SUBINK_HEX)

    add_text(s, rx, 4.45, 3.50, 0.40,
              "Two stars show our current Scene Flow Driving "
              "checkpoint and the projected EPE after full Scene Flow "
              "pre training.",
              size=9, italic=True, color=ACCENT_HEX)
    return s


def build_proposed(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 05 PROPOSED SOLUTION",
        title_main="StereoLite",
        title_accent="at a glance.")

    # Top metrics
    cols_x = [0.55, 2.65, 4.85, 7.15]
    titles = ["TRAINABLE PARAMS", "LATENCY  ·  RTX 3050",
               "LATENCY  ·  ORIN NANO", "VAL EPE"]
    big = ["0.87 M", "54 ms", "~40 ms", "1.54 px"]
    notes = [
        "8.7 MB fp32 checkpoint. INT8 conversion planned for deployment.",
        "Per 512x832 stereo pair, batch size 1.",
        "Projected with INT8 + TensorRT, full validation pending.",
        "Scene Flow Driving 200 val.\n0.515 px on indoor real data.",
    ]
    for x, t, b, n in zip(cols_x, titles, big, notes):
        add_text(s, x, 1.40, 2.10, 0.25, t,
                  size=8, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x, 1.65, 2.00, 0.65, b,
                  size=22, bold=True, color=INK_HEX)
        add_text(s, x, 2.30, 2.00, 0.85, n,
                  size=8.5, color=SUBINK_HEX)

    # Idea boxes
    cols_x = [0.55, 2.95, 5.35, 7.75]
    titles = ["IDEA  ·  01", "IDEA  ·  02", "IDEA  ·  03", "IDEA  ·  04"]
    body = [
        ("Pretrained encoder.  Truncated MobileNetV2 reuses ImageNet "
         "features at four scales."),
        ("One cost volume.  Built only at 1/16 scale, cheap but rich "
         "enough to initialize."),
        ("Iterative refinement.  Two plus three plus three residual "
         "updates over tile hypotheses."),
        ("Learned upsample.  Convex 9 neighbour weights for boundary "
         "sharpness."),
    ]
    for x, t, b in zip(cols_x, titles, body):
        add_text(s, x, 3.45, 2.20, 0.25, t,
                  size=8, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x, 3.75, 2.20, 1.20, b,
                  size=10, color=SUBINK_HEX)

    add_text(s, 0.55, 4.85, 8.95, 0.30,
              "Original work.  This combination of ideas under 1 M "
              "params has not been published before.",
              size=11, italic=True, bold=True, color=ACCENT_HEX)
    return s


def build_methodology(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 06 METHODOLOGY",
        title_main="How we",
        title_accent="executed the thesis.")

    # 4-stage methodology pipeline
    stages = [
        ("01  ·  LITERATURE",
         "110 papers across six stereo matching eras catalogued and "
         "summarised. Three tier ranking by relevance to edge "
         "deployment."),
        ("02  ·  ARCHITECTURE",
         "Three candidate designs prototyped (tile, cascade, SGM). "
         "Final design is a 0.87 M tile + RAFT hybrid."),
        ("03  ·  TRAINING",
         "Synthetic Scene Flow Driving pretraining on Kaggle dual T4. "
         "Indoor real data fine tune on RTX 3050 with FoundationStereo "
         "pseudo ground truth."),
        ("04  ·  VALIDATION",
         "Quantitative: held out EPE on Scene Flow + 50 indoor pairs. "
         "Qualitative: 3D point cloud reconstruction with Open3D + "
         "Phong rendering."),
    ]
    cols_x = [0.55, 2.85, 5.15, 7.45]
    for x, (t, b) in zip(cols_x, stages):
        add_text(s, x, 1.45, 2.20, 0.25, t,
                  size=9, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x, 1.80, 2.20, 2.50, b,
                  size=10, color=INK_HEX)

    add_text(s, 0.55, 4.50, 8.95, 0.40,
              "End to end reproducible pipeline. Code, training "
              "logs, checkpoints, point clouds, and figures all "
              "tracked in the repository.",
              size=10.5, italic=True, color=ACCENT_HEX)
    return s


def build_implementation(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 07 IMPLEMENTATION",
        title_main="Hardware",
        title_accent="and software setup.")

    # Left: hardware photo
    s.shapes.add_picture(
        "/home/abrar/Research/stero_research_claude/papers/raw/"
        "iterative/dummy.txt".replace("papers/raw/iterative/dummy.txt",
                                       "model/designs/d1_tile/"
                                       "deployment_pipeline.png"),
        Inches(0.55), Inches(1.40), height=Inches(3.65))
    # Right text columns
    rx = 4.40
    add_text(s, rx, 1.45, 5.10, 0.30,
              "HARDWARE",
              size=9, bold=True, mono=True, color=ACCENT_HEX)
    add_bullet_lines(s, rx, 1.75, 5.10, 1.50, [
        "Waveshare AR0144 stereo USB camera.  Global shutter, "
        "1280x720 at 60 fps, USB-C, 65 deg HFOV, 52 mm baseline.",
        "Target deployment hardware:  NVIDIA Jetson Orin Nano 8 GB.",
        "Training hardware:  Kaggle dual Tesla T4 (synthetic) and "
        "local RTX 3050 (real data fine tune).",
    ], size=10)

    add_text(s, rx, 3.30, 5.10, 0.30,
              "SOFTWARE",
              size=9, bold=True, mono=True, color=ACCENT_HEX)
    add_bullet_lines(s, rx, 3.60, 5.10, 1.40, [
        "PyTorch 2.11 + CUDA 12.8.  Distributed Data Parallel for "
        "Kaggle training.",
        "AMP (fp16) with GradScaler.  AdamW + OneCycle.  Multi scale "
        "loss with L1, Sobel gradient, bad-1 hinge, smoothness.",
        "Open3D 0.19 for point cloud post processing and rendering.",
    ], size=10)
    return s


def _arch_block(slide, x, y, w, h, *, title, lines):
    """One half-slide architecture block with a title and bullet text."""
    from pptx.enum.shapes import MSO_SHAPE
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    box.line.color.rgb = RGBColor.from_string(SUBINK_HEX)
    box.line.width = Pt(0.5)
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor.from_string("FBF7EE")
    add_text(slide, x + 0.18, y + 0.12, w - 0.36, 0.30, title,
              size=9, bold=True, mono=True, color=ACCENT_HEX)
    add_bullet_lines(slide, x + 0.18, y + 0.45, w - 0.36, h - 0.55,
                      lines, size=9.5, color=INK_HEX)


def build_arch_overview(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 07.1 ARCHITECTURE",
        title_main="StereoLite",
        title_accent="end to end.")
    s.shapes.add_picture(str(ARCH_FULL),
        Inches(0.30), Inches(1.30), width=Inches(9.40))
    add_text(s, 0.55, 4.95, 8.95, 0.30,
              "0.87 M params  ·  54 ms on RTX 3050  ·  ~40 ms expected "
              "on Orin Nano (INT8)  ·  1.54 px val EPE on Scene Flow "
              "Driving.",
              size=9.5, color=SUBINK_HEX, italic=True)
    return s


def build_arch_stack_a(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 07.2 ARCHITECTURE  ·  ENCODER + INIT",
        title_main="Pretrained encoder,",
        title_accent="one cost volume.")
    _arch_block(s, 0.55, 1.40, 4.40, 3.55,
        title="STAGE 1  ·  ENCODER",
        lines=[
            "MobileNetV2-100, ImageNet pretrained.",
            "Shared weights across left and right.",
            "Truncated after stage four; saves 1 M params.",
            "Returns features at 1/2, 1/4, 1/8, 1/16 scale.",
            "0.54 M of the 0.87 M total live here.",
        ])
    _arch_block(s, 5.05, 1.40, 4.40, 3.55,
        title="STAGE 2  ·  TILE HYPOTHESIS INIT",
        lines=[
            "Group correlation cost volume at 1/16 only.",
            "8 groups, 24 disparity candidates.",
            "3D aggregator: two Conv3D + GN + SiLU layers.",
            "Soft argmin over disparity gives initial d.",
            "Per tile state: d, slope x, slope y, 16 channel "
            "feature, confidence in [0, 1].",
        ])
    return s


def build_arch_stack_b(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 07.3 ARCHITECTURE  ·  REFINE + UPSAMPLE",
        title_main="Iterative residuals,",
        title_accent="convex upsample.")
    _arch_block(s, 0.55, 1.40, 4.40, 3.55,
        title="STAGE 3  ·  ITERATIVE REFINEMENT",
        lines=[
            "Each TileRefine iteration: warp right features by "
            "current d, concat with left features, hypothesis state.",
            "3 layer trunk (3x3, hidden 48, GN, SiLU).",
            "Predict residuals: delta d, slope, feature, conf.",
            "Coarse to fine: 1/16 (x2), 1/8 (x3), 1/4 (x3) = 8 iters.",
            "Plane equation upsample between scales.",
        ])
    _arch_block(s, 5.05, 1.40, 4.40, 3.55,
        title="STAGE 4  ·  LEARNED CONVEX UPSAMPLE",
        lines=[
            "RAFT style mask: 9 weights per fine pixel that sum to 1.",
            "Fine disparity = convex average of 9 coarse neighbours.",
            "Mask conditioned on encoder features at the fine scale.",
            "Two blocks chained: 1/4 -> 1/2, 1/2 -> full.",
            "Avoids a third cost volume at fine resolution.",
        ])
    return s


def build_arch_stack_c(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 07.4 ARCHITECTURE  ·  LOSS + BUDGET",
        title_main="Multi scale loss,",
        title_accent="parameter budget.")
    _arch_block(s, 0.55, 1.40, 4.40, 3.55,
        title="SUPERVISION  ·  THREE TERMS",
        lines=[
            "L = sum_k w_k [ L1 + lambda_g L_grad + lambda_h L_hinge ] "
            "+ lambda_s L_smooth.",
            "L1: pixel error on valid GT pixels.",
            "L_grad: L1 on Sobel gradients of disparity.",
            "L_hinge: extra penalty when error > 1 px.",
            "Defaults: lambda_g = 0.5, lambda_h = 0.3, "
            "lambda_s = 0.02.",
        ])
    _arch_block(s, 5.05, 1.40, 4.40, 3.55,
        title="PARAMETER BUDGET  ·  0.87 M TOTAL",
        lines=[
            "Encoder (truncated MobileNetV2)  ·  0.54 M  (62%)",
            "Refine 1/16 (x2 iters)            ·  0.13 M  (15%)",
            "Refine 1/8  (x3 iters)            ·  0.08 M  (9%)",
            "Refine 1/4  (x3 iters)            ·  0.07 M  (8%)",
            "Tile init + cost volume + upsamples ·  0.05 M  (6%)",
            "Checkpoint size on disk: 8.7 MB (fp32).",
        ])
    return s


def build_working_principle(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 08 WORKING PRINCIPLE",
        title_main="Dataset collection",
        title_accent="and training.")
    # Left: training progression GIF (animates in slideshow mode)
    if TRAIN_GIF.exists():
        s.shapes.add_picture(str(TRAIN_GIF),
            Inches(0.55), Inches(1.45), width=Inches(4.80))
    add_text(s, 0.55, 4.60, 4.80, 0.30,
              "Training progression on tracked indoor val pairs.",
              size=9, italic=True, color=SUBINK_HEX, align="center")

    # Right: 4 step blocks
    rx = 5.55
    blocks = [
        ("01  ·  SYNTHETIC PRETRAIN",
         "4,200 pairs from Scene Flow Driving (about 12% of the full "
         "Scene Flow corpus). 30 epochs, Kaggle dual T4."),
        ("02  ·  REAL DATA COLLECTION",
         "1,587 stereo pairs captured with the AR0144 USB camera "
         "indoors at RUET. Quality filter accepts 997."),
        ("03  ·  PSEUDO GROUND TRUTH",
         "FoundationStereo (CVPR 2025, 215 M) generates per pair "
         "disparity targets for the 997 indoor pairs."),
        ("04  ·  REAL DATA FINE TUNE",
         "9,000 steps on RTX 3050.  Mean val EPE drops from 1.54 px "
         "(synthetic only) to 0.515 px (50 indoor val held out)."),
    ]
    by = 1.45
    for t, b in blocks:
        add_text(s, rx, by, 4.00, 0.25, t,
                  size=9, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, rx, by + 0.30, 4.00, 0.65, b,
                  size=10, color=INK_HEX)
        by += 0.95
    return s


def build_results_divider(prs):
    return build_dark_section_slide(prs,
        section_label="§ 09 RESULTS & ANALYSIS",
        title_main="Two domains,",
        title_accent="two checkpoints, one model.",
        big_number="09",
        body_lines=[
            "Section 09.1  ·  Scene Flow Driving (synthetic). The "
            "12% subset shows the architecture trains cleanly; the "
            "remaining 88% is the queued full pretrain.",
            "Section 09.2  ·  Indoor real data (custom). 0.515 px "
            "EPE on 50 held out pairs; the same model evaluated "
            "without rectification still passes.",
            "Section 09.3  ·  Live inference panel video at the "
            "AR0144's native 60 fps stride.",
            "Section 09.4  ·  3D reconstruction from disparity, "
            "demonstrating the geometry is correct enough for "
            "downstream point cloud work.",
        ])


def _add_picture_fit(slide, path, x, y, max_w, max_h):
    """Add picture preserving aspect, fit inside (max_w, max_h)."""
    src = Image.open(path); sw, sh = src.size; src.close()
    aspect = sw / sh
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        w = max_w; h = max_w / aspect
    else:
        h = max_h; w = max_h * aspect
    return slide.shapes.add_picture(str(path),
        Inches(x + (max_w - w) / 2),
        Inches(y + (max_h - h) / 2),
        width=Inches(w), height=Inches(h))


def build_results_sf(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 09.1 RESULTS  ·  SCENE FLOW",
        title_main="Scene Flow Driving",
        title_accent="subset.")

    # Top row: training curves on left, qualitative progression on right
    _add_picture_fit(s, SF_TRAINING_CURVES,
                      x=0.40, y=1.30, max_w=4.40, max_h=2.40)
    _add_picture_fit(s, SF_PROGRESSION,
                      x=5.00, y=1.30, max_w=4.70, max_h=2.40)
    # Captions just under each image
    add_text(s, 0.40, 3.78, 4.40, 0.22,
              "Training curves  ·  Kaggle dual T4, 30 epochs",
              size=8.5, italic=True, color=SUBINK_HEX, align="center")
    add_text(s, 5.00, 3.78, 4.70, 0.22,
              "Qualitative progression  ·  step 500 to 7500",
              size=8.5, italic=True, color=SUBINK_HEX, align="center")

    # Bullets at the bottom (no overlap)
    add_bullet_lines(s, 0.40, 4.10, 9.30, 0.85, [
        "Trained on 12% of the full Scene Flow corpus  ·  "
        "validation EPE 1.54 px on 200 held out pairs.",
        "Loss drops about two orders of magnitude over 30 epochs.",
        "Projected EPE after full Scene Flow pretraining: about "
        "0.71 px (the outlined star on the research gap).",
    ], size=9.5)
    add_text(s, 0.40, 5.00, 9.30, 0.22,
              "Disparity scale: 0 to 120 px on the TURBO colormap, "
              "full resolution.",
              size=8.5, italic=True, color=SUBINK_HEX)
    return s


def build_results_real(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 09.2 RESULTS  ·  INDOOR REAL DATA",
        title_main="Indoor real data",
        title_accent="fine tune.")
    # Shorter image so caption + italic note fit cleanly underneath.
    s.shapes.add_picture(str(FIGS / "realdata_training.png"),
        Inches(1.50), Inches(1.30), width=Inches(7.00))
    add_text(s, 0.55, 4.10, 8.95, 0.30,
              "Teacher: FoundationStereo (215 M params).  Student: "
              "StereoLite (0.87 M).  997 clean pairs, 50 val held out, "
              "9,000 steps in 1 h 35 m on RTX 3050.",
              size=9.5, color=SUBINK_HEX, mono=True)
    add_text(s, 0.55, 4.45, 8.95, 0.30,
              "Camera imperfection tolerant.  Inference runs on raw, "
              "non rectified AR0144 frames and still reaches 0.515 px "
              "mean EPE.",
              size=10.5, italic=True, color=ACCENT_HEX)
    add_text(s, 0.55, 4.85, 8.95, 0.25,
              "Disparity scale here: about 5 to 70 px range across the "
              "indoor scenes.",
              size=8.5, italic=True, color=SUBINK_HEX)
    return s


def build_results_video(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 09.3 RESULTS  ·  LIVE INFERENCE",
        title_main="Live inference",
        title_accent="on indoor stereo video.")
    # Add the trimmed 60 s 1080-wide MP4
    poster = FIGS / "_inference_preview.png"
    if VIDEO_60S.exists():
        s.shapes.add_movie(
            str(VIDEO_60S), Inches(0.55), Inches(1.40),
            Inches(8.95), Inches(2.65),
            poster_frame_image=str(poster),
            mime_type="video/mp4")
    add_text(s, 0.55, 4.20, 8.95, 0.30,
              "60 second clip from a 3 minute walk through, played at "
              "1.5x slowdown so the eye can track the depth gradient. "
              "TURBO colormap, fixed range.",
              size=10, color=SUBINK_HEX)
    add_text(s, 0.55, 4.55, 8.95, 0.30,
              "Disparity scale on the colormap: about 5 to 70 px at "
              "1280x720 (close objects red, far objects blue).",
              size=9, italic=True, color=SUBINK_HEX)
    add_text(s, 0.55, 4.85, 8.95, 0.30,
              "End to end inference on RTX 3050 at 54 ms per frame.  "
              "Expected ~40 ms on Orin Nano with INT8.",
              size=10, italic=True, color=ACCENT_HEX)
    return s


def build_results_3d(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 09.4 RESULTS  ·  3D RECONSTRUCTION",
        title_main="From disparity",
        title_accent="to 3D point clouds.")
    gifs = sorted(PCTOP.glob("pair_*.gif"))[:3]
    panel_w = 2.95
    panel_h = 2.30
    gap = 0.10
    total_w = 3 * panel_w + 2 * gap
    x0 = (10.0 - total_w) / 2.0
    y0 = 1.40
    for i, gif in enumerate(gifs):
        s.shapes.add_picture(str(gif),
            Inches(x0 + i * (panel_w + gap)), Inches(y0),
            width=Inches(panel_w), height=Inches(panel_h))
        epe = gif.stem.split("epe")[-1]
        add_text(s, x0 + i * (panel_w + gap),
                  y0 + panel_h + 0.05, panel_w, 0.25,
                  f"val pair  ·  EPE {epe} px",
                  size=10, mono=True, color=ACCENT_HEX, align="center")

    add_text(s, 0.55, 4.30, 8.95, 0.30,
              "Pinhole projection, AR0144 intrinsics: fx ~= 1005 px, "
              "baseline 52 mm.  Open3D post processing: statistical "
              "outlier removal, 4 mm voxel downsample, Phong shaded "
              "splats.",
              size=9.5, color=SUBINK_HEX)
    add_text(s, 0.55, 4.65, 8.95, 0.30,
              "Disparity scale per cloud: 10 to 50 px range from the "
              "respective frame.",
              size=9, italic=True, color=SUBINK_HEX)
    add_text(s, 0.55, 4.95, 8.95, 0.30,
              "These are evidence that the predicted disparity is "
              "metrically consistent, not just visually plausible.",
              size=10.5, italic=True, color=ACCENT_HEX)
    return s


def build_discussion(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 10 DISCUSSION",
        title_main="Reading the numbers,",
        title_accent="and the potential.")
    cols_x = [0.55, 5.05]
    titles = ["OBSERVATION  ·  CV HEAD PLATEAU",
               "OBSERVATION  ·  REAL DATA TRANSFER"]
    body = [
        "The 1/8 scale CV L1 plateaus near 0.13 px (about 1 px at "
        "full resolution).  This is an effective lower bound on "
        "final EPE until we deepen or multi scale the cost volume.",
        "Indoor EPE drops 3x with only 997 fine tune pairs.  The "
        "architecture transfers cleanly across domains; full Scene "
        "Flow pretraining should narrow the gap further.",
    ]
    for x, t, b in zip(cols_x, titles, body):
        from pptx.enum.shapes import MSO_SHAPE
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.40), Inches(4.40), Inches(1.45))
        card.line.color.rgb = RGBColor.from_string(SUBINK_HEX)
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string("FBF7EE")
        add_text(s, x + 0.18, 1.55, 4.10, 0.25, t,
                  size=8.5, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x + 0.18, 1.85, 4.10, 0.95, b,
                  size=10, color=INK_HEX)

    add_text(s, 0.55, 3.05, 8.95, 0.30,
              "WHAT THE ARCHITECTURE EARNS",
              size=9, bold=True, mono=True, color=ACCENT_HEX)
    add_bullet_lines(s, 0.55, 3.40, 8.95, 1.40, [
        "Tile hypothesis state captures slanted surfaces (roads, "
        "walls, reflective ground) gracefully.",
        "Iterative refinement converges cleanly without overfitting "
        "(narrow train to val gap).",
        "Convex upsample preserves boundaries that bilinear "
        "upsampling would smear.",
    ], size=10.5)

    add_text(s, 0.55, 4.85, 8.95, 0.30,
              "Potential.  StereoLite is competitive at 0.87 M params; "
              "with full Scene Flow pretraining and INT8 deployment it "
              "is a viable foundation for production edge stereo.",
              size=10.5, italic=True, color=ACCENT_HEX)
    return s


def build_challenges(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 11 CHALLENGES & LIMITATIONS",
        title_main="What blocked us,",
        title_accent="what is still open.")
    cards = [
        ("HARDWARE & COMPUTE",
         "Full Scene Flow pretraining requires a high end GPU (A100 "
         "class, 80 GB VRAM) for several days. Our local RTX 3050 and "
         "rented Kaggle T4 sessions caps us at the 12% subset."),
        ("PSEUDO GROUND TRUTH",
         "Indoor real data is supervised by FoundationStereo (215 M) "
         "outputs, not LiDAR.  Absolute accuracy still needs a "
         "calibrated depth reference."),
        ("STRESS TESTING",
         "Inference on raw versus rectified frames shows marginal "
         "difference, but a wider stress suite (multiple cameras, "
         "outdoor lighting, motion blur) needs more hardware."),
    ]
    cols_x = [0.55, 3.70, 6.85]
    for x, (t, b) in zip(cols_x, cards):
        from pptx.enum.shapes import MSO_SHAPE
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(1.40), Inches(2.95), Inches(2.85))
        card.line.color.rgb = RGBColor.from_string(SUBINK_HEX)
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string("FBF7EE")
        add_text(s, x + 0.15, 1.55, 2.65, 0.25, t,
                  size=8.5, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x + 0.15, 1.90, 2.65, 2.30, b,
                  size=10, color=INK_HEX)

    add_text(s, 0.55, 4.45, 8.95, 0.30,
              "Current model already meets the stated objectives.  We "
              "do not want to stop here; we want to compete at global "
              "state of the art with the right resources.",
              size=10.5, italic=True, color=ACCENT_HEX)
    return s


def build_impact(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 12 IMPACT & CONTRIBUTIONS",
        title_main="Original work,",
        title_accent="real deployment paths.")

    add_bullet_lines(s, 0.55, 1.40, 8.95, 1.20, [
        "First published combination of HITNet tile hypotheses + RAFT "
        "iterative refinement at under 1 M parameters. Original work, "
        "not a re implementation of any prior model.",
        "Reproducible reference pipeline:  multi scale supervision, "
        "training curves, public Scene Flow numbers, and a 997 pair "
        "indoor pseudo ground truth dataset.",
    ], size=10.5)

    add_text(s, 0.55, 2.95, 8.95, 0.30,
              "WHERE IT HELPS",
              size=9, bold=True, mono=True, color=ACCENT_HEX)
    cards = [
        ("AUTONOMOUS COMBAT DRONES",
         "Indoor flight obstacle avoidance under tight weight + "
         "wattage budgets."),
        ("FACTORY AUTONOMOUS VEHICLES",
         "AGV / AMR navigation in cluttered warehouses without a "
         "LiDAR mast."),
        ("MOBILE ROBOTICS",
         "Visual SLAM front end with depth from two cheap cameras."),
        ("LOW POWER ADAS",
         "Forward collision warning on embedded ECUs without desktop "
         "GPUs."),
    ]
    # 4 cards in one row, all fitting within 10 in wide slide.
    # Each card 2.10 wide, gap 0.18, total = 4*2.10 + 3*0.18 = 8.94 in
    cw = 2.10
    gap = 0.18
    x0 = (10.0 - (4 * cw + 3 * gap)) / 2.0
    from pptx.enum.shapes import MSO_SHAPE
    for i, (t, b) in enumerate(cards):
        x = x0 + i * (cw + gap)
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(3.30), Inches(cw), Inches(1.10))
        card.line.color.rgb = RGBColor.from_string(SUBINK_HEX)
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string("FBF7EE")
        add_text(s, x + 0.12, 3.40, cw - 0.24, 0.22, t,
                  size=7.8, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x + 0.12, 3.62, cw - 0.24, 0.75, b,
                  size=9, color=INK_HEX)

    add_text(s, 0.55, 4.55, 8.95, 0.30,
              "RESEARCH BASELINE",
              size=9, bold=True, mono=True, color=ACCENT_HEX)
    add_text(s, 0.55, 4.85, 8.95, 0.30,
              "An honest, fully open baseline for the under 1 M "
              "parameter stereo regime that the literature has so "
              "far left empty.",
              size=10.5, color=INK_HEX)
    return s


def build_conclusion(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 13 CONCLUSION",
        title_main="One small model,",
        title_accent="three concrete results.")
    add_bullet_lines(s, 0.55, 1.40, 8.95, 2.00, [
        "0.87 M parameter stereo network that combines a pretrained "
        "MobileNetV2 encoder, a single coarse cost volume, iterative "
        "tile refinement, and learned convex upsampling.",
        "1.54 px EPE on Scene Flow Driving 200 val using only 12% of "
        "the full Scene Flow corpus; 0.515 px on 50 held out indoor "
        "real pairs after fine tuning.",
        "End to end 3D point cloud reconstruction works on raw, "
        "non rectified AR0144 frames, demonstrating camera "
        "imperfection tolerance at edge scale compute.",
    ], size=11)

    add_text(s, 0.55, 3.65, 8.95, 0.30,
              "Both stated objectives are met:  efficient under 1 M, "
              "AI refined, deployable, and tolerant of camera "
              "imperfections.",
              size=11, italic=True, bold=True, color=ACCENT_HEX)

    add_text(s, 0.55, 4.20, 8.95, 0.30,
              "TAKEAWAY",
              size=9, bold=True, mono=True, color=ACCENT_HEX)
    add_text(s, 0.55, 4.50, 8.95, 0.55,
              "Carefully composed light weight modules can push the "
              "accuracy versus size Pareto front into a region the "
              "published literature has not occupied.",
              size=11, color=INK_HEX)
    return s


def build_future_work(prs):
    s = build_cream_chrome_slide(prs,
        section_label="§ 14 FUTURE WORK",
        title_main="Open items,",
        title_accent="next steps.")
    cards = [
        ("01  ·  HARDWARE SUPPORT",
         "Acquire access to A100 class compute for full Scene Flow "
         "pretraining (about 35,000 pairs).  Local RTX 3050 cannot "
         "complete this within the project timeline."),
        ("02  ·  FULL BENCHMARKS",
         "After full pretrain, fine tune and report on KITTI 2015, "
         "Middlebury 2014, and ETH3D.  Target competitive D1-all and "
         "bad-3 metrics."),
        ("03  ·  INTERNATIONAL JOURNAL",
         "Submit a polished version (with the broader benchmarks) to "
         "an international IEEE / Springer journal in computer vision "
         "or robotics."),
        ("04  ·  EDGE DEPLOYMENT",
         "TensorRT INT8 conversion on Jetson Orin Nano.  End to end "
         "ROS 2 node for downstream robotics integration."),
    ]
    cols_x = [0.55, 5.05]
    rows_y = [1.40, 3.10]
    cw, ch = 4.40, 1.55
    from pptx.enum.shapes import MSO_SHAPE
    for i, (t, b) in enumerate(cards):
        x = cols_x[i % 2]
        y = rows_y[i // 2]
        card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(cw), Inches(ch))
        card.line.color.rgb = RGBColor.from_string(SUBINK_HEX)
        card.line.width = Pt(0.5)
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor.from_string("FBF7EE")
        add_text(s, x + 0.18, y + 0.12, cw - 0.36, 0.25, t,
                  size=9, bold=True, mono=True, color=ACCENT_HEX)
        add_text(s, x + 0.18, y + 0.42, cw - 0.36, ch - 0.50, b,
                  size=10, color=INK_HEX)
    return s


# --------------------------------------------------------------------------
# Driver
# --------------------------------------------------------------------------

def renumber_footers(prs):
    n = len(prs.slides)
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            t = sh.text_frame.text.strip()
            if not (t.isdigit() and 1 <= int(t) <= 99): continue
            if sh.left is None or sh.left < Emu(8 * 914400): continue
            if sh.top is None or sh.top < Emu(5.0 * 914400): continue
            paras = list(sh.text_frame.paragraphs)
            if not paras or not paras[0].runs: continue
            paras[0].runs[0].text = f"{i+1:02d}"
            for r in paras[0].runs[1:]: r.text = ""
            break


def main():
    prs = Presentation(str(ORIG))
    print(f"opened {ORIG.name}: {len(prs.slides)} slides")

    edit_title_slide(prs)
    print("  edited title")

    # Build new slides in order; collect references
    new_slides = []

    new_slides.append(build_intro_1_1(prs))           # 3
    new_slides.append(build_intro_1_2(prs))           # 4
    new_slides.append(build_intro_1_3(prs))           # 5
    new_slides.append(build_problem(prs))             # 6
    new_slides.append(build_objectives(prs))          # 7
    new_slides.append(build_literature(prs))          # 8
    new_slides.append(build_research_gap(prs))        # 9
    new_slides.append(build_proposed(prs))            # 10
    new_slides.append(build_methodology(prs))         # 11
    new_slides.append(build_implementation(prs))      # 12
    new_slides.append(build_arch_overview(prs))       # 13
    new_slides.append(build_arch_stack_a(prs))        # 14
    new_slides.append(build_arch_stack_b(prs))        # 15
    new_slides.append(build_arch_stack_c(prs))        # 16
    new_slides.append(build_working_principle(prs))   # 17
    new_slides.append(build_results_divider(prs))     # 18
    new_slides.append(build_results_sf(prs))          # 19
    new_slides.append(build_results_real(prs))        # 20
    new_slides.append(build_results_video(prs))       # 21
    new_slides.append(build_results_3d(prs))          # 22
    new_slides.append(build_discussion(prs))          # 23
    new_slides.append(build_challenges(prs))          # 24
    new_slides.append(build_impact(prs))              # 25
    new_slides.append(build_conclusion(prs))          # 26
    new_slides.append(build_future_work(prs))         # 27

    print(f"  appended {len(new_slides)} new slides")

    # Reorder: title (idx 0), outline (idx 1), then new_slides in order,
    # then Q&A (originally idx 27, now somewhere later).
    qa_slide = prs.slides[T_QA]
    target_idx = 2
    for sl in new_slides:
        move_slide(prs, sl, target_idx)
        target_idx += 1
    move_slide(prs, qa_slide, target_idx)

    # Now rewrite the outline (slide 2) to match the new structure
    build_outline_slide(prs)
    print("  rewrote outline")

    # Delete the original 25 content slides + extras (indices 2 .. 27 - 28
    # before they are pushed down). Easier: identify slides that were NOT
    # in our keep list and delete them.
    keep_ids = {prs.slides[0].slide_id,           # title
                 prs.slides[1].slide_id,           # outline
                 qa_slide.slide_id}
    keep_ids.update(s.slide_id for s in new_slides)
    # Walk and remove all that are not in keep_ids
    xml = prs.slides._sldIdLst
    sld_id_els = list(xml)
    for el in sld_id_els:
        sid = int(el.attrib["id"])
        if sid not in keep_ids:
            xml.remove(el)
    print(f"  pruned to {len(prs.slides)} slides")

    renumber_footers(prs)
    print("  renumbered footers")

    prs.save(str(OUT))
    print(f"\nsaved {OUT}")
    print(f"  {OUT.stat().st_size/1e6:.1f} MB  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
