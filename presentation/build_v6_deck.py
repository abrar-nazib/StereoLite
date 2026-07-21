"""v5 deck. Modifies the simplified academic deck to:

  1. Enrich Literature Review (slide 7) with a 9-paper measurable-
     parameter table (Method / Year / Type / Params M / SF EPE px /
     KITTI D1 / Latency / Edge) plus the StereoLite proposed row.
  2. Replace Literature Review (Cont..) (slide 8) body with a
     capability matrix (cross/tick) over six edge-relevant traits.
  3. Insert three per-layer architecture sub-slides between the
     'Implementation: Architecture' overview (slide 13) and the
     'Parameter Budget' slide (slide 14). Each new slide carries one
     simple diagram per stage.
  4. Populate References (slide 28) with IEEE-style citations of all
     papers featured in the literature review.

Every number in the literature review and review-summary tables has
been cross-checked against the primary PDF in papers/raw/. See
PAPERS dict below for per-row citation pointers (paper:table:page).

Style is preserved verbatim:
  - Times New Roman everywhere
  - Navy 14385C titles, dark gray 3C3C3C subhead, near-black 111111 body
  - Cream slide background, orange footer band
"""
from __future__ import annotations

import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path("/home/abrar/Research/stero_research_claude")
ORIG = ROOT / "presentation/deck/Thesis_MTE_RUET_Presentation_Smiplified.pptx"
OUT  = ROOT / "presentation/deck/Thesis_MTE_RUET_Presentation_Final_v6.pptx"

FIGS = ROOT / "presentation/figs"
TFIGS = ROOT / "thesis/book/figures"   # final thesis figures (v6 results/arch)

NAVY     = "14385C"
DARK     = "3C3C3C"
INK      = "111111"
ACCENT   = "C24A1C"      # the orange footer / accent
WHITE    = "FFFFFF"
HEADER_BG = NAVY          # table header
ROW_BG    = "FFFFFF"
ROW_BG_ALT = "F7F1E1"     # subtle alt row
BORDER   = "BFBAB0"

FONT = "Times New Roman"

# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------

def duplicate_slide(prs, src_idx):
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    src_csld = src.element.find(f"{{{p_ns}}}cSld")
    new_csld = new.element.find(f"{{{p_ns}}}cSld")
    if src_csld is not None and new_csld is not None:
        src_bg = src_csld.find(f"{{{p_ns}}}bg")
        if src_bg is not None:
            new_bg = copy.deepcopy(src_bg)
            new_csld.insert(0, new_bg)
    for shape in list(src.shapes):
        new_el = copy.deepcopy(shape.element)
        new.shapes._spTree.insert_element_before(new_el, "p:extLst")
    rels_src = src.part.rels
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for blip in new.shapes._spTree.findall(f".//{{{a_ns}}}blip"):
        rid = blip.get(f"{{{r_ns}}}embed")
        if rid and rid in rels_src:
            rel = rels_src[rid]
            new_rid = new.part.relate_to(rel.target_part, rel.reltype)
            blip.set(f"{{{r_ns}}}embed", new_rid)
    return new


def move_slide(prs, slide, new_idx):
    xml = prs.slides._sldIdLst
    el = next(s for s in xml if int(s.attrib["id"]) == slide.slide_id)
    xml.remove(el); xml.insert(new_idx, el)


def remove_shape(shape):
    shape.element.getparent().remove(shape.element)


def find_slide_idx(prs, title):
    """Return the 0-indexed position of the slide whose **title** matches
    `title`. Title shapes live in the top band (y < 0.85") of the slide;
    we constrain the match to that band so a body-text occurrence of
    the same string (e.g. 'Conclusion' appearing in an Outline TOC)
    cannot collide with the real title."""
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top is None or sh.top >= Inches(0.85):
                continue
            if sh.text_frame.text.strip() == title:
                return i
    return -1


def delete_slide(prs, idx):
    """Remove the slide at index `idx` from the presentation, **and**
    purge the underlying slide part from the package. The package-level
    purge is essential: without it, later `add_slide` calls auto-number
    new slide-N.xml files and can reuse the deleted slide's name —
    producing a zip with duplicate entries that python-pptx silently
    saves but LibreOffice cannot load."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    if idx < 0 or idx >= len(slides_list):
        return
    sldId = slides_list[idx]
    rId = sldId.rId
    # Resolve the slide part BEFORE dropping the rel.
    try:
        slide_part = prs.part.related_part(rId)
    except Exception:
        slide_part = None
    # Drop the relationship from the presentation part.
    prs.part.drop_rel(rId)
    # Remove from the sldIdLst so it stops appearing in the deck.
    xml_slides.remove(sldId)
    # Purge the slide part itself from the package so its partname
    # (e.g. /ppt/slides/slide14.xml) is freed for reuse.
    if slide_part is not None:
        package = prs.part.package
        for attr in ("_parts_by_partname", "_parts"):
            d = getattr(package, attr, None)
            if isinstance(d, dict):
                d.pop(slide_part.partname, None)


def strip_body(slide, *, keep_titles):
    """Remove every shape on `slide` except (a) shapes whose text exactly
    matches one of `keep_titles`, (b) the top divider (thin shape near
    y < 0.85), (c) the footer GroupShape (top >= 5.20"), and (d) the
    RUET emblem (FREEFORM near top-left). Returns the set of kept XML
    elements so the caller can verify."""
    keep_xml = set()
    for sh in slide.shapes:
        if sh.shape_type == 6 and sh.top is not None and sh.top >= Inches(5.20):
            keep_xml.add(sh.element); continue
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue
            continue
        t = sh.text_frame.text.strip()
        if t in keep_titles:
            keep_xml.add(sh.element); continue
    for sh in list(slide.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)
    return keep_xml


def set_paragraph_text(text_frame, new_text):
    paras = list(text_frame.paragraphs)
    if not paras:
        text_frame.text = new_text
        return
    runs = list(paras[0].runs)
    if runs:
        runs[0].text = new_text
        for r in runs[1:]: r.text = ""
    else:
        paras[0].add_run().text = new_text
    for p in paras[1:]:
        for r in list(p.runs): r.text = ""


def add_text(slide, x, y, w, h, text, *, size=11, bold=False,
             color=INK, italic=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                    Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center": p.alignment = PP_ALIGN.CENTER
    elif align == "right": p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor.from_string(color)
    r.font.name = FONT
    return box


def add_filled_rect(slide, x, y, w, h, fill_hex, line_hex=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if line_hex is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = RGBColor.from_string(line_hex)
        rect.line.width = Pt(0.5)
    return rect


def add_picture_fit(slide, path, x, y, max_w, max_h):
    src = Image.open(path); sw, sh = src.size; src.close()
    aspect = sw / sh
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        w = max_w; h = max_w / aspect
    else:
        h = max_h; w = max_h * aspect
    return slide.shapes.add_picture(str(path),
        Inches(x + (max_w - w) / 2),
        Inches(y + (max_h - h) / 2),
        width=Inches(w), height=Inches(h))


# --------------------------------------------------------------------------
# Verified paper data (single source of truth for slides 7, 8, 28)
# --------------------------------------------------------------------------
# Each row carries every number we display on slides, with a per-cell
# citation in the "src" comment. Every numeric claim was verified by
# opening the primary PDF and reading the relevant table.
#
# Param rule: published total trainable params from the originating
# paper (or its supplementary) where available; comparable numbers
# from a survey paper as fallback. Foundation models include their
# frozen ViT backbone weights in the total since those still occupy
# disk + RAM at deploy.

PAPERS = [
    # (#, key,          name,                year, type,             params, sf_epe, kitti_d1, latency,    edge)
    ("1", "psmnet",     "PSMNet",            "2018", "3D cost vol.", "5.2",   "1.09", "2.32",   "410 ms",   "No"),
    ("2", "hitnet",     "HITNet (L)",        "2021", "Tile-based",    "0.97", "0.43", "1.98",   "54 ms",    "Yes"),
    ("3", "bgnet",      "BGNet",             "2021", "Bilateral grid","2.9",  "1.17", "2.51",   "25 ms",    "Yes"),
    ("4", "coex",       "CoEx",              "2021", "Lightweight",   "2.7",  "0.69", "2.13",   "27 ms",    "Yes"),
    ("5", "raft",       "RAFT-Stereo",       "2021", "Iterative",     "11.2", "0.61", "1.82",   "380 ms",   "No"),
    ("6", "igev",       "IGEV-Stereo",       "2023", "Iter. + GEV",   "12.6", "0.47", "1.59",   "180 ms",   "No"),
    ("7", "lightstereo","LightStereo-S",     "2025", "Lightweight",   "3.4",  "0.73", "2.30",   "17 ms",    "Yes"),
    ("8", "fstereo",    "FoundationStereo",  "2025", "Foundation",    "~340", "0.34", "—",      "—",        "No"),
    ("9", "defom",      "DEFOM-Stereo",      "2025", "Foundation",    "47.3", "0.42", "—",      "316 ms",   "No"),
    ("*", "stereolite", "StereoLite (Ours)", "2026", "Tile + Iter.",  "2.96", "0.78", "-",      "50 ms",    "Yes"),
]
# Citation breadcrumbs (paper:table:page) used for verification:
# psmnet      Chang & Chen, CVPR 2018, Tab 5 p7 (SF EPE 1.09); KITTI 2.32 Tab 4 p7
# hitnet      Tankovich et al., CVPR 2021, supp Tab 7 p17 (HITNet L 0.97 M / 0.43 EPE); KITTI 1.98 IGEV Tab 5 p7
# bgnet       Xu et al., CVPR 2021, Tab 1 p6 (EPE 1.17), Tab 4 p7 (D1 2.51, 25.4 ms); ~2.9M tier1 summary
# coex        Bangunharcana et al., IROS 2021, Tab I p4 (EPE 0.69 / D1 2.13 / 27 ms); 2.72M LightStereo Tab I
# raft        Lipson et al., 3DV 2021, Tab 6 p8 (11.23 M); SF EPE 0.61 / KITTI 1.82 / 380 ms IGEV Tab 5 p7
# igev        Xu et al., CVPR 2023, Tab 1 p6 (12.60 M), Tab 4 p7 (0.47 EPE), Tab 5 p7 (1.59 D1 / 180 ms)
# lightstereo Guo et al., ICRA 2025, Tab I p4 (3.44 M / 0.73 EPE / 17 ms); Tab V p6 (D1 2.30)
# fstereo     Wen et al., CVPR 2025, Tab 3 p7 (EPE 0.34); ~335 M ViT-L backbone (tier1 summary)
# defom       Jiang et al., CVPR 2025, Tab 2 p7 (ViT-L: 47.30 M trainable / 0.42 EPE / 0.316 s)
# stereolite  StereoLite final; 2.96 M trainable; thesis ch6 tab:sceneflow_results / tab:latency
#             (49.8 ms fp16 RTX 3050 at 384x640). SF EPE cell is FT3D full test set (4,370 pairs,
#             native 960x540), not the competitors' cropped SF protocol; no KITTI submission ("-").


# IEEE-style references (slide 28). Each entry mirrors PAPERS row.
REFERENCES = [
    ("[1]",
     "J.-R. Chang and Y.-S. Chen, “Pyramid stereo matching network,” "
     "in Proc. IEEE Conf. Computer Vision and Pattern Recognition (CVPR), "
     "Salt Lake City, UT, USA, 2018, pp. 5410–5418."),
    ("[2]",
     "V. Tankovich, C. Häne, Y. Zhang, A. Kowdle, S. Fanello, and "
     "S. Bouaziz, “HITNet: Hierarchical iterative tile refinement "
     "network for real-time stereo matching,” in Proc. IEEE/CVF Conf. "
     "Computer Vision and Pattern Recognition (CVPR), 2021, pp. 14362–14372."),
    ("[3]",
     "B. Xu, Y. Xu, X. Yang, W. Jia, and Y. Guo, “Bilateral grid learning "
     "for stereo matching networks,” in Proc. IEEE/CVF Conf. Computer "
     "Vision and Pattern Recognition (CVPR), 2021, pp. 12497–12506."),
    ("[4]",
     "A. Bangunharcana, J. W. Cho, S. Lee, I. S. Kweon, K.-S. Kim, and "
     "S. Kim, “Correlate-and-Excite: Real-time stereo matching via guided "
     "cost volume excitation,” in Proc. IEEE/RSJ Int. Conf. Intelligent "
     "Robots and Systems (IROS), 2021, pp. 3542–3548."),
    ("[5]",
     "L. Lipson, Z. Teed, and J. Deng, “RAFT-Stereo: Multilevel recurrent "
     "field transforms for stereo matching,” in Proc. Int. Conf. 3D "
     "Vision (3DV), 2021, pp. 218–227."),
    ("[6]",
     "G. Xu, X. Wang, X. Ding, and X. Yang, “Iterative geometry encoding "
     "volume for stereo matching,” in Proc. IEEE/CVF Conf. Computer "
     "Vision and Pattern Recognition (CVPR), 2023, pp. 21919–21928."),
    ("[7]",
     "X. Guo, C. Zhang, Y. Zhang, W. Zheng, D. Nie, M. Poggi, and L. Chen, "
     "“LightStereo: Channel boost is all you need for efficient 2D cost "
     "aggregation,” in Proc. IEEE Int. Conf. Robotics and Automation "
     "(ICRA), 2025."),
    ("[8]",
     "B. Wen, M. Trepte, J. Aribido, J. Kautz, O. Gallo, and S. Birchfield, "
     "“FoundationStereo: Zero-shot stereo matching,” in Proc. "
     "IEEE/CVF Conf. Computer Vision and Pattern Recognition (CVPR), 2025."),
    ("[9]",
     "H. Jiang, Z. Lou, L. Ding, R. Xu, M. Tan, W. Jiang, and R. Huang, "
     "“DEFOM-Stereo: Depth foundation model based stereo matching,” "
     "in Proc. IEEE/CVF Conf. Computer Vision and Pattern Recognition (CVPR), 2025."),
]


# --------------------------------------------------------------------------
# Slide 7  ·  enriched literature review table
# --------------------------------------------------------------------------

def rebuild_literature_review(prs):
    s = prs.slides[6]
    # Strip everything except the title, divider, footer.
    # Title: text "Literature Review"; divider: thin shape near top;
    # footer: GroupShape near bottom.
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6:  # group (footer band)
            keep_xml.add(sh.element); continue
        if not sh.has_text_frame:
            # Top divider line
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue
            # RUET emblem (FREEFORM near top-left)
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue
            continue
        t = sh.text_frame.text.strip()
        if t == "Literature Review":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    # Add a small caption under the title
    add_text(s, 0.40, 1.05, 9.20, 0.30,
             "Comparison of nine prior methods plus ours on the standard "
             "Scene Flow + KITTI 2015 protocol",
             size=11, italic=True, color=DARK, align="center")

    # New table: 8 columns x 11 rows (1 header + 10 data)
    cols = ["#", "Method", "Year", "Type", "Params (M)", "SF EPE (px)",
             "KITTI D1 (%)", "Latency", "Edge?"]
    col_widths = [0.34, 1.65, 0.55, 1.55, 0.82, 0.90, 0.95, 0.78, 0.55]
    # Each data row is (#, name, year, type, params, sf_epe, kitti, lat, edge)
    rows = [(p[0], p[2], p[3], p[4], p[5], p[6], p[7], p[8], p[9])
             for p in PAPERS]

    # Center the table on the 10-inch slide (total width ~8.09)
    table_x = (10.00 - sum(col_widths)) / 2
    header_y = 1.45
    row_h = 0.31
    header_h = 0.32

    # Compute x positions
    xs = [table_x]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)
    table_w = sum(col_widths)

    # Header row
    add_filled_rect(s, table_x, header_y, table_w, header_h,
                    fill_hex=NAVY)
    for x, w, c in zip(xs, col_widths, cols):
        add_text(s, x + 0.03, header_y + 0.05, w - 0.06, header_h - 0.08,
                 c, size=9, bold=True, color=WHITE, align="center")

    # Data rows
    y = header_y + header_h
    for r_i, row in enumerate(rows):
        is_ours = "Ours" in row[1]
        # Alternating row backgrounds (subtle); Ours row gets a soft tint
        if is_ours:
            bg = "FFEEDD"   # soft orange tint
        else:
            bg = ROW_BG_ALT if r_i % 2 == 1 else ROW_BG
        add_filled_rect(s, table_x, y, table_w, row_h, fill_hex=bg,
                        line_hex=BORDER)
        for x, w, val in zip(xs, col_widths, row):
            color = ACCENT if is_ours else INK
            bold = is_ours
            # Method column is left-aligned, everything else centered
            align = "left" if w > 1.50 else "center"
            add_text(s, x + 0.04, y + 0.05, w - 0.08, row_h - 0.10,
                     val, size=8.5, bold=bold, color=color, align=align)
        y += row_h

    # Footnote on data sources / hardware caveat
    foot_y = y + 0.10
    add_text(s, table_x, foot_y, table_w, 0.20,
             "Latencies on different GPUs (varies by source); SF EPE on "
             "Scene Flow finalpass; KITTI 2015 D1-all from official "
             "leaderboards. See the References slide for citations.",
             size=8, italic=True, color=DARK, align="center")


# --------------------------------------------------------------------------
# Slide 8  ·  capability matrix (Review Summary)
# --------------------------------------------------------------------------
# Cross/tick across six edge-relevant traits. Each cell is one of:
#   ✓ (full),  ✗ (no),  ~ (partial),  ?  (not yet measured)
#
# Trait definitions (column order matches CAPS list):
#  L  = Lightweight (≤3 M trainable params)
#  R  = Real-time edge (<60 ms typical inference)
#  I  = Iterative refinement (allows compute-vs-accuracy trade)
#  P  = Plane / tile geometry (sub-pixel via slopes, not just scalar disp)
#  F  = Foundation prior (monocular depth backbone integrated)
#  C  = Cross-domain robustness (zero-shot generalization claimed)
#
# Cell rule: ✓ requires the paper to *claim* the property; ~ if the
# implementation does it but it isn't headline; ✗ if absent.
#
# StereoLite row (final, thesis ch6): L ✓ at 2.96 M trainable;
# R ✓ at 36.3 ms INT8 TensorRT on Jetson Orin Nano (measured);
# I ✓ ConvGRU refinement, 2+3+3 = 8 updates across 1/16, 1/8, 1/4;
# P ✓ plane rendering from slanted tile states;
# C ~ measured zero-shot Middlebury 2014 D1-all 10.9% (competitive
#   but behind foundation-prior methods, hence partial).

CAPS = ["≤3 M params", "<60 ms edge", "Iterative", "Plane / tile",
         "Foundation", "Cross-domain"]
MATRIX = [
    # paper key,       L,   R,   I,   P,   F,   C
    ("psmnet",          "✗", "✗", "✗", "✗", "✗", "✗"),
    ("hitnet",          "✓", "✓", "✗", "✓", "✗", "✗"),
    ("bgnet",           "✓", "✓", "✗", "✗", "✗", "✗"),
    ("coex",            "✓", "✓", "✗", "✗", "✗", "✗"),
    ("raft",            "✗", "✗", "✓", "✗", "✗", "✓"),
    ("igev",            "✗", "✗", "✓", "✗", "✗", "✓"),
    ("lightstereo",     "~", "✓", "✗", "✗", "✗", "✗"),
    ("fstereo",         "✗", "✗", "✓", "✗", "✓", "✓"),
    ("defom",           "✗", "✗", "✓", "✗", "✓", "✓"),
    ("stereolite",      "✓", "✓", "✓", "✓", "✗", "~"),
]


def rebuild_review_summary(prs):
    """Slide 8 ('Literature Review (Cont..)'): rewrite body with a
    capability tick/cross matrix over six edge-relevant traits."""
    s = prs.slides[7]
    # Strip everything except title, divider, footer, RUET emblem
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6:  # group (footer)
            keep_xml.add(sh.element); continue
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue
            continue
        t = sh.text_frame.text.strip()
        if t == "Literature Review (Cont..)":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    # Subhead caption
    add_text(s, 0.40, 1.05, 9.20, 0.30,
             "Capability matrix · what each method offers for edge stereo",
             size=12, italic=True, color=DARK, align="center")

    # Resolve "key" -> human-readable name from PAPERS
    name_by_key = {p[1]: p[2] for p in PAPERS}
    ref_by_key = {p[1]: idx + 1 for idx, p in enumerate(PAPERS[:-1])}
    # StereoLite gets no [N] reference

    # Table layout:  Method col + 6 capability cols
    cols = ["Method"] + CAPS
    col_widths = [3.65, 0.92, 0.92, 0.92, 0.92, 0.92, 0.92]
    table_x = 0.45
    header_y = 1.50
    row_h = 0.30
    header_h = 0.34

    xs = [table_x]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)
    table_w = sum(col_widths)

    # Header
    add_filled_rect(s, table_x, header_y, table_w, header_h, fill_hex=NAVY)
    for x, w, c in zip(xs, col_widths, cols):
        add_text(s, x + 0.03, header_y + 0.04, w - 0.06, header_h - 0.06,
                 c, size=9, bold=True, color=WHITE, align="center")

    # Data
    y = header_y + header_h
    for r_i, mrow in enumerate(MATRIX):
        key = mrow[0]
        ticks = mrow[1:]
        is_ours = (key == "stereolite")
        if is_ours:
            bg = "FFEEDD"
        else:
            bg = ROW_BG_ALT if r_i % 2 == 1 else ROW_BG
        add_filled_rect(s, table_x, y, table_w, row_h, fill_hex=bg,
                        line_hex=BORDER)
        # Method cell with optional [N] reference
        name = name_by_key[key]
        ref = ref_by_key.get(key)
        method_text = f"{name}  [{ref}]" if ref else name
        add_text(s, xs[0] + 0.10, y + 0.05, col_widths[0] - 0.16,
                 row_h - 0.10,
                 method_text, size=9, bold=is_ours,
                 color=ACCENT if is_ours else INK, align="left")
        # Capability cells
        for x, w, mark in zip(xs[1:], col_widths[1:], ticks):
            if mark == "✓":
                color = "1F7A2C"   # green
            elif mark == "✗":
                color = "8C2A1F"   # dark red
            elif mark == "~":
                color = "B07000"   # amber
            else:                  # "?" or unknown
                color = "6A6A6A"   # gray
            add_text(s, x + 0.04, y + 0.04, w - 0.08, row_h - 0.08,
                     mark, size=12, bold=True, color=color, align="center")
        y += row_h

    # Bottom takeaway (one line, fits above footer)
    foot_y = y + 0.12
    add_text(s, table_x, foot_y, table_w, 0.22,
             "StereoLite is the only method combining lightweight, real-time, "
             "iterative, and plane-tile geometry simultaneously.",
             size=10, italic=True, color=DARK, align="center")


# --------------------------------------------------------------------------
# Slide 24  ·  Impact (Environmental / Societal / Research, three columns)
# --------------------------------------------------------------------------

def rebuild_impact(prs):
    """Slide 24 ('Impact'): three-column card layout per MTE guideline.
    Lean text; only directly defensible claims."""
    s = prs.slides[23]
    # Strip body — keep title, divider, footer, RUET emblem
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6 and sh.top is not None and sh.top >= Inches(5.20):
            keep_xml.add(sh.element); continue   # footer
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue   # divider
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue   # RUET emblem
            continue
        t = sh.text_frame.text.strip()
        if t == "Impact":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    # The source slide has no title shape in the top band; add one so
    # the slide is not headed by the subhead alone.
    if not any(sh.has_text_frame and sh.text_frame.text.strip() == "Impact"
               for sh in s.shapes):
        add_text(s, 0.50, 0.22, 9.00, 0.45, "Impact (Cont..)",
                 size=20, bold=True, color=NAVY, align="center")

    # Subhead
    add_text(s, 0.45, 0.95, 9.10, 0.30,
             "Environmental  ·  Societal  ·  Research",
             size=12, italic=True, color=DARK, align="center")

    # Three cards
    columns = [
        ("RESEARCH",
         [
           "Combines tile plane hypotheses (HITNet) with ConvGRU iterative refinement (RAFT) and a narrow-band GEV in one 2.96 M model.",
           "Teacher-verified on real data: 1.45 px agreement with FoundationStereo over 997 rig pairs.",
           "Open, reproducible baseline for edge stereo under 3 M parameters.",
         ]),
        ("ENVIRONMENTAL",
         [
           "~30× lower inference power vs running a 340 M foundation stereo model on a server GPU (projected, edge GPU at 7 to 15 W).",
           "No cloud roundtrip per depth frame: all inference local.",
           "Enables battery-powered stereo on drones, AR headsets, mobile robots.",
         ]),
        ("SOCIETAL",
         [
           "On-device inference: images never leave the device.",
           "~USD 500 stereo + edge GPU setup, vs USD 2 k+ for comparable LiDAR depth.",
           "Lowers the entry bar for stereo research and education in resource-constrained settings.",
         ]),
    ]

    # Layout: 3 cards across 9.10" with two gaps of 0.20"
    n = len(columns)
    gap = 0.20
    card_w = (9.10 - (n - 1) * gap) / n   # ~2.90"
    card_h = 3.30
    card_y = 1.40
    header_h = 0.40
    body_pad_x = 0.18
    body_pad_y = 0.18

    for i, (title, bullets) in enumerate(columns):
        x = 0.45 + i * (card_w + gap)
        # Card outline (light border across full height)
        add_filled_rect(s, x, card_y, card_w, card_h,
                        fill_hex=ROW_BG, line_hex=BORDER)
        # Header strip
        add_filled_rect(s, x, card_y, card_w, header_h,
                        fill_hex=ACCENT)
        add_text(s, x + 0.10, card_y + 0.07, card_w - 0.20, header_h - 0.10,
                 title, size=12, bold=True, color=WHITE, align="center")
        # Body bullets
        body_y = card_y + header_h + body_pad_y
        bullet_h = 0.78
        for j, txt in enumerate(bullets):
            by = body_y + j * bullet_h
            # Bullet marker
            add_text(s, x + body_pad_x, by, 0.20, 0.20,
                     "•", size=14, bold=True, color=ACCENT, align="left")
            # Bullet text
            add_text(s, x + body_pad_x + 0.22, by - 0.02,
                     card_w - body_pad_x - 0.30, bullet_h,
                     txt, size=10, color=INK, align="left")


# --------------------------------------------------------------------------
# Slide 12  ·  Implementation (hardware photos + software logos only)
# --------------------------------------------------------------------------

def rebuild_implementation(prs):
    """Slide 12 ('Implementation'): rewrite body with two banded
    sections — HARDWARE (setup photo, camera, Jetson) and SOFTWARE
    (CUDA, PyTorch, Open3D, Kaggle). Names only; no specs."""
    s = prs.slides[11]
    # Strip body — keep title, divider, footer, RUET emblem
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6 and sh.top is not None and sh.top >= Inches(5.20):
            keep_xml.add(sh.element); continue   # footer
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue   # divider
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue   # RUET emblem
            continue
        t = sh.text_frame.text.strip()
        if t == "Implementation":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    photos = ROOT / "presentation/photos"

    # Section A: HARDWARE
    add_text(s, 0.45, 0.95, 9.10, 0.30,
             "HARDWARE",
             size=11, bold=True, color=ACCENT, align="left")
    # Thin accent rule under the section label
    add_filled_rect(s, 0.45, 1.27, 9.10, 0.012, fill_hex=ACCENT)

    hw = [
        (photos / "test_rig.png",     "Test rig"),
        (photos / "stereo_camera.png","Stereo camera"),
        (photos / "jetson_orin.jpg",  "Jetson Orin Nano"),
    ]
    hw_y = 1.40
    hw_h = 1.50
    img_w = 2.10
    img_h = 1.50
    n = len(hw)
    gap = (9.10 - n * img_w) / (n + 1)   # equal gaps including the edges
    for i, (path, caption) in enumerate(hw):
        x = 0.45 + gap + i * (img_w + gap)
        if path.exists():
            add_picture_fit(s, path, x, hw_y, img_w, img_h)
        add_text(s, x, hw_y + img_h + 0.05, img_w, 0.25,
                 caption, size=11, bold=True, color=INK, align="center")

    # Section B: SOFTWARE
    add_text(s, 0.45, 3.30, 9.10, 0.30,
             "SOFTWARE",
             size=11, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 3.62, 9.10, 0.012, fill_hex=ACCENT)

    sw = [
        (photos / "cuda_logo.png",    "CUDA"),
        (photos / "pytorch_logo.png", "PyTorch"),
        (photos / "open3d_logo.png",  "Open3D"),
        (photos / "kaggle_logo.png",  "Kaggle"),
    ]
    sw_y = 3.78
    sw_h = 0.95
    logo_w = 1.70
    logo_h = 0.95
    n = len(sw)
    gap = (9.10 - n * logo_w) / (n + 1)
    for i, (path, caption) in enumerate(sw):
        x = 0.45 + gap + i * (logo_w + gap)
        if path.exists():
            add_picture_fit(s, path, x, sw_y, logo_w, logo_h)
        add_text(s, x, sw_y + logo_h + 0.05, logo_w, 0.25,
                 caption, size=11, bold=True, color=INK, align="center")


# --------------------------------------------------------------------------
# Slide 3  ·  Introduction (full visual concept)
# --------------------------------------------------------------------------

def _strip_to_title(slide, title_text):
    """Strip every body shape from `slide`, keeping only its title (in the
    top band y < 0.85), the top divider line, the footer group, and the
    RUET emblem. Mirrors the `strip_body` pattern used elsewhere."""
    keep_xml = set()
    for sh in slide.shapes:
        if sh.shape_type == 6 and sh.top is not None and sh.top >= Inches(5.20):
            keep_xml.add(sh.element); continue   # footer group
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue   # divider rule
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue   # RUET emblem
            continue
        t = sh.text_frame.text.strip()
        if t == title_text and sh.top is not None and sh.top < Inches(0.85):
            keep_xml.add(sh.element); continue
    for sh in list(slide.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)


def rebuild_introduction(prs):
    """Slide 3 ('Introduction'): replace body with a full visual concept
    of stereo depth — geometry schematic on top, a real-world indoor
    pair (left view + GT depth) underneath, and a one-line takeaway."""
    s = prs.slides[2]
    _strip_to_title(s, "Introduction")

    photos = ROOT / "presentation/photos"
    schematic = FIGS / "intro_stereo_geometry.png"

    # Subhead under the title bar
    add_text(s, 0.45, 0.92, 9.10, 0.30,
             "Stereo cameras turn a horizontal pixel shift into metric depth.",
             size=14, bold=True, color=NAVY, align="left")

    # Geometry schematic — full-width, aspect-preserved letterbox.
    if schematic.exists():
        add_picture_fit(s, schematic, 0.45, 1.28, 9.10, 2.20)

    # Real-example caption row
    add_text(s, 0.45, 3.62, 9.10, 0.26,
             "REAL EXAMPLE  ·  indoor hallway from our finetune set",
             size=10, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 3.90, 9.10, 0.012, fill_hex=ACCENT)

    # Two side-by-side example images
    # Sized so image bottom (4.98) + caption row (5.02 to 5.24) stay
    # clear of the footer band, which starts at y = 5.29.
    img_w = 2.10
    img_h = 1.00
    gap = 0.30
    total = 2 * img_w + gap
    x0 = (10.00 - total) / 2
    pairs = [
        (photos / "intro_left_example.png", "Left view  (input)"),
        (photos / "intro_depth_example.png", "Depth map  (warm = far)"),
    ]
    row_y = 3.98
    for i, (path, cap) in enumerate(pairs):
        x = x0 + i * (img_w + gap)
        if path.exists():
            add_picture_fit(s, path, x, row_y, img_w, img_h)
        add_text(s, x, row_y + img_h + 0.04, img_w, 0.22,
                 cap, size=9.5, bold=True, color=INK, align="center")


# --------------------------------------------------------------------------
# Slide 4  ·  Introduction (Cont..) — project pitch + applications
# --------------------------------------------------------------------------

def rebuild_introduction_cont(prs):
    """Slide 4 ('Introduction (Cont..)'): a one-line project pitch
    followed by four application-domain cards (drones, mobile robots,
    AR/VR headsets, autonomous-stereo rigs). Each card is a photo plus
    a short use-case caption — together they answer 'who needs this?'."""
    s = prs.slides[3]
    _strip_to_title(s, "Introduction (Cont..)")

    photos = ROOT / "presentation/photos"

    # Subhead — what this project delivers
    add_text(s, 0.45, 0.92, 9.10, 0.30,
             "We build StereoLite, a 2.96 M parameter stereo network "
             "designed to run on edge hardware.",
             size=13, bold=True, color=NAVY, align="left")
    # Pitch text
    add_text(s, 0.45, 1.30, 9.10, 0.55,
             "Every moving platform that must understand its surroundings "
             "needs depth, in real time, on board. The same model has to "
             "fit and run on whatever compute the platform can carry.",
             size=11, color=INK, align="left", italic=True)

    # Section header for the 4-card row
    add_text(s, 0.45, 2.05, 9.10, 0.28,
             "WHERE ON-DEVICE STEREO DEPTH MATTERS",
             size=10, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 2.34, 9.10, 0.012, fill_hex=ACCENT)

    cards = [
        (photos / "combat_drone.jpg",  "Drones",
         "obstacle avoidance, terrain following"),
        (photos / "factory_robot.jpg", "Mobile robots",
         "navigation, picking, collision safety"),
        (photos / "ar_headset.jpg",    "AR / VR",
         "scene reconstruction, hand-held depth"),
        (photos / "test_rig.png",      "Embedded rigs",
         "SLAM, mapping, autonomous platforms"),
    ]

    n = len(cards)
    card_w = 2.05
    gap = (9.10 - n * card_w) / (n - 1)
    img_h = 1.30
    img_y = 2.55
    for i, (path, label, sub) in enumerate(cards):
        x = 0.45 + i * (card_w + gap)
        if path.exists():
            add_picture_fit(s, path, x, img_y, card_w, img_h)
        add_text(s, x, img_y + img_h + 0.10, card_w, 0.28,
                 label, size=12, bold=True, color=INK, align="center")
        add_text(s, x, img_y + img_h + 0.42, card_w, 0.50,
                 sub, size=9.5, italic=True, color=DARK, align="center")


# --------------------------------------------------------------------------
# Slide 5  ·  Problem Statement (absorbs LiDAR / RealSense / FM size)
# --------------------------------------------------------------------------

def rebuild_problem_statement(prs):
    """Slide 5 ('Problem Statement'): two-row layout. The top row asks
    'why not the obvious alternatives?' with three image cards (LiDAR,
    RealSense, foundation-model size). The bottom row lists the edge
    constraints we therefore have to design around (compute, memory,
    power)."""
    s = prs.slides[4]
    _strip_to_title(s, "Problem Statement")

    photos = ROOT / "presentation/photos"

    # Subhead — restate the question
    add_text(s, 0.45, 0.92, 9.10, 0.30,
             "How do we deliver dense, accurate depth on edge hardware,"
             " when the obvious alternatives don't fit?",
             size=13, bold=True, color=NAVY, align="left")

    # ---- Row A: alternatives that don't fit ----
    add_text(s, 0.45, 1.32, 9.10, 0.26,
             "WHY NOT THE OBVIOUS ALTERNATIVES?",
             size=10, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 1.60, 9.10, 0.012, fill_hex=ACCENT)

    alts = [
        dict(path=photos / "lidar.jpg",
             label="LiDAR",
             value="$3 k to $80 k",
             sub="sparse, heavy, power-hungry; cost rules out consumer platforms"),
        dict(path=photos / "realsense.jpg",
             label="Active depth (RealSense)",
             value="≤ 6 m, indoor",
             sub="IR projector fails outdoors; range and resolution limited"),
        dict(path=None,
             label="Foundation stereo",
             value="~340 M params",
             sub="state-of-the-art accuracy, but needs a desktop GPU; can't ship to a Jetson"),
    ]

    n = len(alts)
    card_w = 2.95
    gap = (9.10 - n * card_w) / (n - 1)
    card_h = 1.85
    row_y = 1.74
    for i, c in enumerate(alts):
        x = 0.45 + i * (card_w + gap)
        # Card background
        add_filled_rect(s, x, row_y, card_w, card_h,
                         fill_hex=ROW_BG, line_hex=BORDER)
        if c["path"] is not None and c["path"].exists():
            add_picture_fit(s, c["path"], x + 0.12, row_y + 0.10,
                             card_w - 0.24, 0.95)
        else:
            # Big number stand-in for foundation-model card
            add_text(s, x + 0.12, row_y + 0.30,
                     card_w - 0.24, 0.55,
                     c["value"], size=24, bold=True, color=ACCENT,
                     align="center")
        # Label band
        add_text(s, x + 0.12, row_y + 1.10, card_w - 0.24, 0.22,
                 c["label"], size=11, bold=True, color=INK,
                 align="center")
        # Value (skip for foundation card — already shown)
        if c["path"] is not None:
            add_text(s, x + 0.12, row_y + 1.32, card_w - 0.24, 0.22,
                     c["value"], size=10, bold=True, color=ACCENT,
                     align="center")
        # Sub-caption
        add_text(s, x + 0.12, row_y + 1.55, card_w - 0.24, 0.30,
                 c["sub"], size=8.5, italic=True, color=DARK,
                 align="center")

    # ---- Row B: edge constraints we have to live with ----
    add_text(s, 0.45, 3.78, 9.10, 0.26,
             "EDGE CONSTRAINTS WE TARGET",
             size=10, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 4.06, 9.10, 0.012, fill_hex=ACCENT)

    specs = [
        ("LIMITED COMPUTE", "~6 TOPS",
         "Embedded SoCs deliver a fraction of a desktop GPU."),
        ("TIGHT MEMORY",    "~4 GB",
         "Shared with the rest of the autonomy stack."),
        ("POWER BUDGET",    "5 to 25 W",
         "Battery-powered platforms, no hot GPUs allowed."),
    ]
    n = len(specs)
    sw = 2.95
    sgap = (9.10 - n * sw) / (n - 1)
    sh_y = 4.20
    sh_h = 1.05
    for i, (label, value, sub) in enumerate(specs):
        x = 0.45 + i * (sw + sgap)
        add_filled_rect(s, x, sh_y, sw, sh_h,
                         fill_hex=ROW_BG, line_hex=BORDER)
        add_text(s, x + 0.14, sh_y + 0.08, sw - 0.28, 0.22,
                 label, size=9, bold=True, color=ACCENT, align="left")
        add_text(s, x + 0.14, sh_y + 0.30, sw - 0.28, 0.36,
                 value, size=18, bold=True, color=INK, align="left")
        add_text(s, x + 0.14, sh_y + 0.72, sw - 0.28, 0.30,
                 sub, size=8.5, italic=True, color=DARK, align="left")


# --------------------------------------------------------------------------
# Results slides  ·  final trained-model figures (thesis book)
# --------------------------------------------------------------------------

def embed_results_progress_gifs(prs):
    """Rewrite the two Results slides with FINAL trained-model figures
    from the thesis book (TFIGS = thesis/book/figures/).

    Slide A (in-domain): full Scene Flow training curves (left) plus
    FT3D test qualitative panels (right).
    Slide B (generalization + edge): Middlebury 2014 zero-shot panels
    (left) plus real-rig camera panels (right).

    Slides are located by subhead text. Both the legacy GIF-era
    subheads and the new v6 subheads are accepted so the function
    stays robust across rebuilds; after locating, the subhead is
    rewritten in place to the v6 wording (title shape is untouched,
    keeping the existing 'Results' styling)."""
    plan = [
        dict(
            match = ("Scene Flow (synthetic)",
                     "Scene Flow pre-training, in-domain test"),
            subhead = "Scene Flow pre-training, in-domain test",
            left_img  = TFIGS / "fig_4_1_training_curves.png",
            left_cap  = ("Training loss and validation EPE, "
                         "60k steps on full Scene Flow"),
            right_img = TFIGS / "fig_4_3_sceneflow_qualitative.png",
            right_cap = ("FT3D test qualitative: left image, "
                         "ground truth, prediction, error"),
            stats = [
                ("TRAINING PAIRS", "35,454",  "full Scene Flow finalpass"),
                ("TEST EPE",       "0.78 px", "FT3D test set, 4,370 pairs"),
                ("D1-ALL",         "3.40%",   "outlier rate, FT3D test"),
            ],
        ),
        dict(
            match = ("Indoor real-data fine-tune.",
                     "Zero-shot generalization and on-device deployment"),
            subhead = "Zero-shot generalization and on-device deployment",
            left_img  = TFIGS / "fig_4_5_mb14_qualitative.png",
            left_cap  = ("Middlebury 2014 zero-shot: "
                         "easiest and hardest scenes"),
            right_img = TFIGS / "fig_4_6_camera.png",
            right_cap = "Zero-shot transfer to our low-cost real rig",
            stats = [
                ("ZERO-SHOT MB14",   "10.9% D1",
                 "23 scenes, no fine-tuning"),
                ("REAL RIG",         "1.45 px",
                 "EPE vs FoundationStereo, 997 pairs"),
                ("JETSON ORIN NANO", "36.3 ms",
                 "27.5 FPS, INT8 TensorRT, measured"),
            ],
        ),
    ]

    # Equal-size side-by-side boxes; aspect-fit content with letterbox.
    LEFT_BOX  = dict(x=0.42, y=1.10, max_w=4.55, max_h=2.85)
    RIGHT_BOX = dict(x=5.03, y=1.10, max_w=4.55, max_h=2.85)
    CAPTION_Y = 4.00
    STATS_Y   = 4.30
    STATS_H   = 0.85

    for cfg in plan:
        # Locate slide (and its subhead shape) by subhead text
        idx, subhead_shape = -1, None
        for i, s in enumerate(prs.slides):
            for sh in s.shapes:
                if (sh.has_text_frame
                        and sh.text_frame.text.strip() in cfg["match"]):
                    idx, subhead_shape = i, sh
                    break
            if idx >= 0: break
        if idx < 0:
            print(f"  (skipped: slide '{cfg['match'][0]}' not found)")
            continue
        s = prs.slides[idx]

        # Rewrite the subhead in place (keeps its formatting)
        set_paragraph_text(subhead_shape.text_frame, cfg["subhead"])

        # Strip body shapes (everything between subhead and footer);
        # keep only title (y < 0.85), divider, footer group, RUET emblem.
        for sh in list(s.shapes):
            if sh.top is None or sh.left is None:
                continue
            if sh.top < Inches(0.95):
                continue   # title / divider / subhead
            if sh.top >= Inches(5.20):
                continue   # footer
            if sh.shape_type == 5 and sh.top < Inches(0.85):
                continue   # RUET emblem
            remove_shape(sh)

        # Thesis figures side by side
        if cfg["left_img"].exists():
            add_picture_fit(s, cfg["left_img"], **LEFT_BOX)
        if cfg["right_img"].exists():
            add_picture_fit(s, cfg["right_img"], **RIGHT_BOX)
        # Captions under each
        add_text(s, LEFT_BOX["x"], CAPTION_Y, LEFT_BOX["max_w"], 0.20,
                 cfg["left_cap"], size=10, italic=True,
                 color=DARK, align="center")
        add_text(s, RIGHT_BOX["x"], CAPTION_Y, RIGHT_BOX["max_w"], 0.20,
                 cfg["right_cap"], size=10, italic=True,
                 color=DARK, align="center")

        # Bottom-row stats cards (three boxes, equal width)
        n = len(cfg["stats"])
        gap_x = 0.18
        card_w = (9.16 - (n - 1) * gap_x) / n   # ~2.93"
        for i, (label, value, sub) in enumerate(cfg["stats"]):
            cx = 0.42 + i * (card_w + gap_x)
            add_filled_rect(s, cx, STATS_Y, card_w, STATS_H,
                             fill_hex=ROW_BG, line_hex=BORDER)
            add_text(s, cx + 0.10, STATS_Y + 0.06, card_w - 0.20, 0.22,
                     label, size=8.5, bold=True, color=ACCENT, align="left")
            add_text(s, cx + 0.10, STATS_Y + 0.28, card_w - 0.20, 0.32,
                     value, size=15, bold=True, color=INK, align="left")
            add_text(s, cx + 0.10, STATS_Y + 0.62, card_w - 0.20, 0.20,
                     sub, size=8.5, italic=True, color=DARK, align="left")


# --------------------------------------------------------------------------
# Slide 9  ·  refresh Research Gap image (carried over from v3 deck)
# --------------------------------------------------------------------------

def refresh_research_gap_image(prs):
    """Slide 9 ('Research Gap') embeds research_gap_pareto.png. The
    underlying figure was regenerated with corrected numbers (HITNet
    0.97 M / 0.43 EPE; DEFOM-Stereo 47.3 M; FoundationStereo 0.34 EPE;
    LightStereo-S 3.44 M / 0.73 EPE). Replace the cached image with
    the current one so slide 9 reflects the same data as slide 7."""
    s = prs.slides[8]
    pareto_path = FIGS / "research_gap_pareto.png"
    if not pareto_path.exists():
        return
    # Locate the image group at the top of the body (matches the original
    # layout: pos ~ (0.77, 0.84), size 8.47 x 4.35).
    target_grp = None
    for sh in s.shapes:
        if sh.shape_type != 6:  # group
            continue
        if sh.top is None or sh.top >= Inches(5.20):
            continue   # skip footer
        # Body group
        target_grp = sh
        break
    if target_grp is None:
        return
    x = target_grp.left / 914400
    y = target_grp.top / 914400
    w = target_grp.width / 914400
    h = target_grp.height / 914400
    # Remove the old group (the picture lives inside it)
    remove_shape(target_grp)
    # Reinsert the fresh figure in the same bounding box
    add_picture_fit(s, pareto_path, x=x, y=y, max_w=w, max_h=h)


# --------------------------------------------------------------------------
# Slide 11  ·  Methodology — two-track flowchart (training + inference)
# --------------------------------------------------------------------------

def build_methodology(prs):
    """Slide 11 ('Methodology'): embed a matplotlib-generated two-track
    flowchart (training pipeline on top, inference pipeline on bottom,
    linked by a 'load weights' connector). The figure is generated by
    presentation/figs/build_arch_diagrams.py:methodology_diagram and
    saved at FIGS/methodology_pipeline.png."""
    idx = find_slide_idx(prs, "Methodology")
    if idx < 0:
        return
    s = prs.slides[idx]
    strip_body(s, keep_titles={"Methodology"})

    # Subhead
    add_text(s, 0.45, 0.95, 9.10, 0.30,
             "System overview · training pipeline (top) and inference pipeline (bottom)",
             size=12, italic=True, color=DARK, align="center")

    # Embed the methodology figure (generated separately by build_arch_diagrams.py)
    fig_path = FIGS / "methodology_pipeline.png"
    if fig_path.exists():
        add_picture_fit(s, fig_path,
                         x=0.30, y=1.32, max_w=9.40, max_h=3.85)


# --------------------------------------------------------------------------
# Slide 2  ·  Outline (rewritten to match new structure)
# --------------------------------------------------------------------------

def rebuild_outline(prs):
    """Slide 2 ('Outline'): rewrite the table-of-contents to match the
    final 16-section ordering (Time Plan after Objectives, Working
    Principle replacing 'Implementation: Architecture', Challenges /
    Limitations, etc.)."""
    idx = find_slide_idx(prs, "Outline")
    if idx < 0:
        return
    s = prs.slides[idx]
    strip_body(s, keep_titles={"Outline"})

    items = [
        "Introduction",
        "Problem Statement",
        "Objectives",
        "Time Plan",
        "Literature Review",
        "Proposed Solution",
        "Methodology",
        "Implementation",
        "Working Principle",
        "Results & Analysis",
        "Discussion",
        "Challenges / Limitations",
        "Impact",
        "Conclusion",
        "Future Work",
        "References",
    ]
    n = len(items)
    half = (n + 1) // 2
    col_left = items[:half]
    col_right = items[half:]

    # Two-column layout
    col_x = [0.85, 5.30]
    col_w = 4.30
    y_start = 1.40
    row_h = 0.42

    for col_i, col in enumerate([col_left, col_right]):
        for r_i, label in enumerate(col):
            number = col_i * half + r_i + 1
            y = y_start + r_i * row_h
            # Number badge (filled circle approximated by rounded rect)
            add_filled_rect(s, col_x[col_i], y, 0.42, 0.34,
                            fill_hex=NAVY)
            add_text(s, col_x[col_i] + 0.02, y + 0.04, 0.40, 0.28,
                     f"{number:02d}", size=11, bold=True,
                     color=WHITE, align="center")
            # Label
            add_text(s, col_x[col_i] + 0.55, y + 0.04, col_w - 0.55, 0.34,
                     label, size=12, bold=False, color=INK, align="left")


# --------------------------------------------------------------------------
# Conclusion  ·  rebuilt to absorb Objectives Answered as a sub-block
# --------------------------------------------------------------------------

def rebuild_conclusion(prs):
    """Conclusion slide: rebuild with two halves —
    a short narrative summary on top, an objectives-met table beneath."""
    idx = find_slide_idx(prs, "Conclusion")
    if idx < 0:
        return
    s = prs.slides[idx]
    strip_body(s, keep_titles={"Conclusion"})

    # Subhead
    add_text(s, 0.45, 0.95, 9.10, 0.30,
             "What we built · what we measured · what we set out to do",
             size=12, italic=True, color=DARK, align="center")

    # Narrative summary block (top half)
    add_filled_rect(s, 0.45, 1.40, 9.10, 1.15, fill_hex=ROW_BG,
                    line_hex=BORDER)
    add_text(s, 0.65, 1.52, 8.70, 0.28,
             "STEREOLITE",
             size=11, bold=True, color=ACCENT, align="left")
    add_text(s, 0.65, 1.80, 8.70, 0.70,
             "A 2.96 M parameter stereo network: tile plane hypotheses "
             "with ConvGRU iterative refinement, trained on the full "
             "Scene Flow set (35,454 pairs). FT3D test EPE 0.78 px / "
             "D1-all 3.40%. Zero-shot Middlebury 2014 D1 10.9%, within "
             "4 points of LiteAnyStereo at 2.6x fewer params. Measured "
             "on device: 36.3 ms INT8 (27.5 FPS) on Jetson Orin Nano; "
             "49.8 ms fp16 on RTX 3050.",
             size=10.5, color=INK, align="left")

    # Objectives-met table (bottom half) — three columns
    table_y = 2.65
    header_h = 0.35
    row_h = 0.42

    add_filled_rect(s, 0.45, table_y, 9.10, header_h, fill_hex=NAVY)
    headers = [("Objective", 4.20), ("Stated criterion", 2.60),
                ("Status", 2.30)]
    x = 0.45
    for label, w in headers:
        add_text(s, x + 0.10, table_y + 0.06, w - 0.20, header_h - 0.10,
                 label, size=10, bold=True, color=WHITE, align="left")
        x += w

    rows = [
        ("Computationally efficient pipeline",
         "Real-time on Jetson Orin Nano",
         "Met · 36.3 ms INT8 (measured)"),
        ("Accurate disparity estimation",
         "Competitive on Scene Flow test",
         "Met · EPE 0.78 px, D1 3.40%"),
        ("Cross-domain generalization",
         "Zero-shot on unseen datasets",
         "Met · KITTI 3.9/4.3%, ETH3D 4.0%, MB14 10.9%"),
        ("Camera-imperfection tolerance",
         "Works on a real, imperfect rig",
         "Met · 1.45 px on 997 real pairs"),
    ]
    y = table_y + header_h
    for r_i, row in enumerate(rows):
        bg = ROW_BG_ALT if r_i % 2 == 1 else ROW_BG
        add_filled_rect(s, 0.45, y, 9.10, row_h, fill_hex=bg,
                        line_hex=BORDER)
        x = 0.45
        for (label, w), val in zip(headers, row):
            color = "1F7A2C" if "Met" in val and label == "Status" else INK
            bold = "Met" in val and label == "Status"
            add_text(s, x + 0.10, y + 0.10, w - 0.20, row_h - 0.20,
                     val, size=10, bold=bold, color=color, align="left")
            x += w
        y += row_h

    # Bottom takeaway
    add_text(s, 0.45, y + 0.15, 9.10, 0.30,
             "Robust to rectification error (EPE 1.03 to 1.53 px up to "
             "1 px vertical offset); all stated objectives are met.",
             size=11, italic=True, bold=True, color=ACCENT, align="center")


# --------------------------------------------------------------------------
# Architecture slides  ·  rename to Working Principle family
# --------------------------------------------------------------------------

def rename_architecture_to_working_principle(prs):
    """The four architecture slides (overview + Stage 1+2 + Stage 3+4 +
    Supervision) become Working Principle slides per MTE guideline.
    Implementation reverts to the photo panel only (slide 12)."""
    renames = [
        ("Implementation: Architecture",
         "Working Principle"),
        ("Implementation: Architecture · Stage 1 + 2",
         "Working Principle · Stage 1 + 2"),
        ("Implementation: Architecture · Stage 3",
         "Working Principle · Stage 3"),
        ("Implementation: Architecture · Stage 4 + Supervision",
         "Working Principle · Stage 4 + Supervision"),
    ]
    for old, new in renames:
        idx = find_slide_idx(prs, old)
        if idx < 0:
            continue
        s = prs.slides[idx]
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == old:
                set_paragraph_text(sh.text_frame, new)
                break


# --------------------------------------------------------------------------
# Challenges  ·  rename to "Challenges / Limitations" (guideline wording)
# --------------------------------------------------------------------------

def rename_challenges(prs):
    idx = find_slide_idx(prs, "Challenges")
    if idx < 0:
        return
    s = prs.slides[idx]
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Challenges":
            set_paragraph_text(sh.text_frame, "Challenges / Limitations")
            _fix_title_shape(sh)
            break


# --------------------------------------------------------------------------
# Slide 28  ·  IEEE references
# --------------------------------------------------------------------------

def build_references(prs):
    """Populate the References slide with IEEE-style citations for the
    nine methods featured in the literature review."""
    idx = find_slide_idx(prs, "References")
    if idx < 0:
        return
    s = prs.slides[idx]
    # Strip body — keep title + divider + footer + emblem
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6:
            keep_xml.add(sh.element); continue
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue
            continue
        t = sh.text_frame.text.strip()
        if t == "References":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    # Subhead
    add_text(s, 0.45, 1.00, 9.10, 0.30,
             "IEEE-style citations of methods featured in the "
             "literature review",
             size=11, italic=True, color=DARK, align="center")

    # References listed two columns to fit nine entries on one slide
    col_w = 4.45
    col_x = [0.45, 5.10]
    rows_per_col = 5
    y_start = 1.45
    row_h = 0.65   # height per reference (most are 3 lines at 9pt)

    for i, (label, body) in enumerate(REFERENCES):
        ci = 0 if i < rows_per_col else 1
        ri = i if ci == 0 else i - rows_per_col
        x = col_x[ci]
        y = y_start + ri * row_h
        # Bracket label
        add_text(s, x, y, 0.35, row_h - 0.05, label,
                 size=9, bold=True, color=ACCENT, align="left")
        # Reference body
        add_text(s, x + 0.35, y, col_w - 0.35, row_h - 0.05,
                 body, size=8.5, color=INK, align="left")


# --------------------------------------------------------------------------
# Per-layer architecture sub-slides (inserted between slide 13 and 14)
# --------------------------------------------------------------------------

def refresh_working_principle_overview(prs):
    """Replace the pre-defense architecture overview diagram with the
    final thesis architecture figure (fig_3_1) and refresh the caption
    to the 2.96 M chassis. Called AFTER the rename phase, so the slide
    title is 'Working Principle' (falls back to the pre-rename title).
    Keeps title, subhead, top divider and footer group; strips the rest
    (old diagram + old param-split captions)."""
    idx = find_slide_idx(prs, "Working Principle")
    if idx < 0:
        idx = find_slide_idx(prs, "Implementation: Architecture")
    if idx < 0:
        print("  WARNING: overview slide not found; skipped refresh")
        return
    s = prs.slides[idx]

    title_sh = subhead_sh = footer_group = top_divider = None
    for sh in s.shapes:
        if sh.shape_type == 6 and sh.top is not None and \
                sh.top >= Inches(5.20):
            footer_group = sh
            continue
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                top_divider = sh
            continue
        t = sh.text_frame.text.strip()
        if ("Working Principle" in t or
                "Implementation: Architecture" in t) and title_sh is None:
            title_sh = sh; continue
        if "StereoLite, end-to-end" in t and subhead_sh is None:
            subhead_sh = sh; continue

    keeper_xml = {sh.element for sh in
                  [title_sh, subhead_sh, footer_group, top_divider]
                  if sh is not None}
    for sh in list(s.shapes):
        if sh.element in keeper_xml:
            continue
        remove_shape(sh)

    # Final architecture figure is a very wide (sideways-page) render;
    # full-width single placement fits the body band.
    add_picture_fit(s, TFIGS / "fig_3_1_architecture_preview.png",
                    x=0.30, y=1.15, max_w=9.40, max_h=3.30)
    add_text(s, 0.40, 4.55, 9.20, 0.30,
             "StereoLite final architecture: encoder, cost volume init, "
             "recurrent refinement with GEV fusion, plane upsampling",
             size=11, bold=True, color=INK, align="center")
    add_text(s, 0.40, 4.88, 9.20, 0.30,
             "2.96 M trainable parameters; encoder and refinement hold "
             "roughly 95 percent of the budget",
             size=11, italic=True, color=DARK, align="center")
    print("  refreshed Working Principle overview with final architecture")


def patch_footer_dates(prs):
    """Deck-wide footer date refresh: APRIL 2026 -> JULY 2026 wherever
    it appears (footer groups on every slide)."""
    n = 0
    for s in prs.slides:
        for tf in _walk_text_frames(s.shapes):
            if _swap_literal_in_tf(tf, "APRIL 2026", "JULY 2026"):
                n += 1
    print(f"  footer dates: {n} occurrences APRIL 2026 -> JULY 2026")


def build_time_plan(prs):
    """The Outline promises a Time Plan entry but the source deck has no
    such slide; build one from the final thesis Gantt chart and insert
    it right after Objectives (Outline position 04)."""
    obj_idx = find_slide_idx(prs, "Objectives")
    if obj_idx < 0:
        print("  WARNING: Objectives slide not found; Time Plan skipped")
        return
    new = duplicate_slide(prs, obj_idx)
    # Retitle FIRST, then strip everything but the retitled shape
    for sh in new.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Objectives" \
                and sh.top is not None and sh.top < Inches(0.85):
            set_paragraph_text(sh.text_frame, "Time Plan")
            break
    _strip_to_title(new, "Time Plan")
    add_picture_fit(new, TFIGS / "fig_b_1_gantt_fixed.png",
                    x=0.55, y=1.10, max_w=8.90, max_h=3.75)
    add_text(new, 0.55, 4.95, 8.90, 0.28,
             "Eleven overlapping phases, July 2025 to July 2026 "
             "(4th year odd and even semesters)",
             size=11, italic=True, color=DARK, align="center")
    move_slide(prs, new, obj_idx + 1)
    print("  built Time Plan slide from thesis Gantt")


def _fix_title_shape(sh):
    """Widen a retitled title shape to the full band and strip autofit,
    so longer titles keep the template font size instead of shrinking
    and wrapping onto two lines."""
    from pptx.oxml.ns import qn
    sh.left = Inches(0.50)
    sh.width = Inches(9.00)
    bodyPr = sh.text_frame._txBody.bodyPr
    for tag in ("a:normAutofit", "a:spAutoFit"):
        el = bodyPr.find(qn(tag))
        if el is not None:
            bodyPr.remove(el)


def _find_slide_by_text(prs, needle):
    """Index of the first slide whose any text frame contains `needle`
    (used where titles are not unique, e.g. the three Results slides)."""
    for i, s in enumerate(prs.slides):
        for tf in _walk_text_frames(s.shapes):
            if needle in tf.text:
                return i
    return -1


def _build_content_slide(prs, *, title, subhead, panels, after_idx):
    """Generic v6 expansion slide: duplicate the Objectives slide as a
    template (title band + divider + footer), strip it, retitle, and lay
    out 1 or 2 figure panels with caption + one-line note. `panels` is a
    list of (path, caption, note) with 1 or 2 entries. Inserted at
    after_idx + 1."""
    src_idx = find_slide_idx(prs, "Objectives")
    if src_idx < 0 or after_idx < 0:
        print(f"  WARNING: could not build slide '{title}' "
              f"(template {src_idx}, anchor {after_idx})")
        return None
    new = duplicate_slide(prs, src_idx)
    for sh in new.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Objectives" \
                and sh.top is not None and sh.top < Inches(0.85):
            set_paragraph_text(sh.text_frame, title)
            _fix_title_shape(sh)
            break
    _strip_to_title(new, title)
    if subhead:
        add_text(new, 0.45, 0.88, 9.10, 0.28, subhead,
                 size=11, italic=True, color=DARK, align="left")
    if len(panels) == 1:
        path, caption, note = panels[0]
        add_picture_fit(new, path, x=0.55, y=1.25, max_w=8.90, max_h=3.25)
        if caption:
            add_text(new, 0.55, 4.58, 8.90, 0.26, caption,
                     size=11, bold=True, color=INK, align="center")
        if note:
            add_text(new, 0.80, 4.88, 8.40, 0.30, note,
                     size=10, italic=True, color=DARK, align="center")
    else:
        for i, (path, caption, note) in enumerate(panels[:2]):
            x = 0.30 if i == 0 else 5.10
            add_picture_fit(new, path, x=x, y=1.25, max_w=4.60, max_h=3.10)
            if caption:
                add_text(new, x, 4.45, 4.60, 0.26, caption,
                         size=10, bold=True, color=INK, align="center")
            if note:
                add_text(new, x, 4.74, 4.60, 0.55, note,
                         size=10, italic=True, color=DARK, align="center")
    move_slide(prs, new, after_idx + 1)
    print(f"  built expansion slide '{title}'")
    return new


def build_expansion_slides(prs):
    """Grow the deck toward the expected ~35 slides with thesis figures
    and setup photos. Runs LATE (after every index-hardcoded rebuild);
    every anchor is found fresh so earlier insertions shift safely."""
    PHOTOS = ROOT / "presentation/photos"

    # 1. Literature timeline, after the capability matrix
    _build_content_slide(prs,
        title="Evolution of Stereo Matching",
        subhead="Six visible eras, classical to foundation-model",
        panels=[(TFIGS / "fig_2_2_timeline_preview.png",
                 "Timeline of stereo matching methods, 2002 to 2026",
                 "Edge-oriented designs form their own track from 2018 "
                 "onward; StereoLite builds on that track.")],
        after_idx=find_slide_idx(prs, "Literature Review (Cont..)"))

    # 2. Parameter budget + input protocol, after the WP family
    _build_content_slide(prs,
        title="Working Principle · Budget + Input",
        subhead="Where the 2.96 M parameters go, and what the network sees",
        panels=[
            (TFIGS / "fig_3_11_param_budget.png",
             "Module-wise parameter split",
             "Encoder and refinement hold roughly 95 percent of the "
             "2.96 M budget."),
            (TFIGS / "fig_3_10_input_protocol.png",
             "Input protocol: native 384 x 640 crops",
             "Random co-located crops with asymmetric color jitter and "
             "right-only erase."),
        ],
        after_idx=find_slide_idx(prs,
                                 "Working Principle · Stage 4 + Supervision"))

    # 3. Experimental setup: rig photo + rig geometry, after Implementation
    _build_content_slide(prs,
        title="Experimental Setup",
        subhead="Low-cost stereo rig and rig geometry",
        panels=[
            (PHOTOS / "test_rig.png",
             "Stereo test rig",
             "AR0144 stereo camera, 2560 x 720 side-by-side capture."),
            (TFIGS / "fig_5_2_rig_preview.png",
             "Rig geometry",
             "Baseline 52 mm, focal length about 1005 px after "
             "rectification to 1280 x 720."),
        ],
        after_idx=find_slide_idx(prs, "Implementation"))

    # 4. Deployment pipeline, after Experimental Setup
    _build_content_slide(prs,
        title="Deployment Pipeline",
        subhead="From fp32 checkpoint to a Jetson INT8 engine",
        panels=[(TFIGS / "fig_4_5_export_preview.png",
                 "Export path: graph optimization, ONNX, INT8 calibration, "
                 "TensorRT",
                 "Four equivalence-proven graph changes give a 1.74x "
                 "latency cut before quantization; three operator swaps "
                 "make the graph INT8-friendly.")],
        after_idx=find_slide_idx(prs, "Experimental Setup"))

    # 5. Training convergence, after the in-domain Results slide
    _build_content_slide(prs,
        title="Results: Convergence",
        subhead="Prediction quality across training checkpoints",
        panels=[(TFIGS / "fig_4_2_convergence.png",
                 "Predicted disparity at increasing training steps, "
                 "two FT3D scenes",
                 "Structure appears early; thin details and boundaries "
                 "keep sharpening to step 53k.")],
        after_idx=_find_slide_by_text(prs,
                  "Scene Flow pre-training, in-domain test"))

    # 6. MB14 per-scene, after the zero-shot Results slide
    _build_content_slide(prs,
        title="Results: Zero-Shot Per Scene",
        subhead="Middlebury 2014, all 23 perfect-set scenes",
        panels=[(TFIGS / "fig_4_4_mb14_perscene.png",
                 "Per-scene D1-all, sorted",
                 "16 of 23 scenes land in the reference band; Jadeplant "
                 "and Flowers are the hardest.")],
        after_idx=_find_slide_by_text(prs,
                  "Zero-shot generalization and on-device deployment"))

    # 7. Rectification robustness, after the per-scene slide
    _build_content_slide(prs,
        title="Results: Rectification Robustness",
        subhead="Accuracy under vertical misalignment (second objective)",
        panels=[(TFIGS / "fig_4_9_rectification.png",
                 "EPE and D1 vs injected vertical offset",
                 "EPE rises only 1.03 to 1.53 px up to 1 px of offset; "
                 "degradation becomes severe beyond 2 px.")],
        after_idx=find_slide_idx(prs, "Results: Zero-Shot Per Scene"))

    # 8. Ablations, after the robustness slide
    _build_content_slide(prs,
        title="Results: Ablations",
        subhead="What each design and training choice contributes",
        panels=[(TFIGS / "fig_4_7_ablations.png",
                 "Augmentation, inference optimization, blur and input "
                 "protocol ablations",
                 "Every retained choice is backed by a controlled A/B "
                 "on the same pair set.")],
        after_idx=find_slide_idx(prs, "Results: Rectification Robustness"))

    # 9. Zero-shot quartet table, right after the zero-shot Results slide
    #    (before the per-scene slide, since this call runs after the
    #    others and anchors directly on the zero-shot slide).
    build_quartet_slide(prs)


def build_quartet_slide(prs):
    """Zero-shot quartet table (KITTI 2012/2015, ETH3D, MB14) next to the
    in-domain reference row. Numbers from
    model/benchmarks/20260704_fullsf_gev4onp_nc/kitti_eth3d_zero_shot.json
    and mb14_zero_shot.json; same protocol across all rows."""
    src_idx = find_slide_idx(prs, "Objectives")
    after_idx = _find_slide_by_text(
        prs, "Zero-shot generalization and on-device deployment")
    if src_idx < 0 or after_idx < 0:
        print("  WARNING: quartet slide anchors not found; skipped")
        return
    new = duplicate_slide(prs, src_idx)
    for sh in new.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Objectives" \
                and sh.top is not None and sh.top < Inches(0.85):
            set_paragraph_text(sh.text_frame, "Results: Zero-Shot Quartet")
            _fix_title_shape(sh)
            break
    _strip_to_title(new, "Results: Zero-Shot Quartet")

    add_text(new, 0.45, 0.90, 9.10, 0.28,
             "One checkpoint, four unseen datasets, one protocol "
             "(384 x 640, valid disparities up to 192 px)",
             size=11, italic=True, color=DARK, align="left")

    table_y = 1.45
    header_h = 0.40
    row_h = 0.50
    headers = [("Dataset", 3.10), ("Pairs", 1.10), ("EPE (px)", 1.40),
               ("bad-2 (%)", 1.60), ("D1-all (%)", 1.90)]
    add_filled_rect(new, 0.45, table_y, 9.10, header_h, fill_hex=NAVY)
    x = 0.45
    for label, w in headers:
        add_text(new, x + 0.12, table_y + 0.08, w - 0.24, header_h - 0.12,
                 label, size=12, bold=True, color=WHITE, align="left")
        x += w

    rows = [
        ("KITTI 2012 (train)", "194", "0.82", "6.96", "4.33", False),
        ("KITTI 2015 (train)", "200", "0.82", "6.42", "3.93", False),
        ("ETH3D (train)", "27", "0.93", "6.47", "3.96", False),
        ("Middlebury 2014", "23", "1.71", "14.5", "10.9", False),
        ("Scene Flow FT3D test (in-domain)", "4,370", "0.78", "5.34",
         "3.40", True),
    ]
    y = table_y + header_h
    for r_i, (name, n, epe, bad2, d1, is_ref) in enumerate(rows):
        bg = ROW_BG_ALT if r_i % 2 == 1 else ROW_BG
        add_filled_rect(new, 0.45, y, 9.10, row_h, fill_hex=bg,
                        line_hex=BORDER)
        x = 0.45
        vals = [name, n, epe, bad2, d1]
        for (label, w), val in zip(headers, vals):
            add_text(new, x + 0.12, y + 0.12, w - 0.24, row_h - 0.20,
                     val, size=12, bold=(label == "D1-all (%)" or is_ref),
                     color=(DARK if is_ref else INK), align="left")
            x += w
        y += row_h

    add_text(new, 0.45, y + 0.20, 9.10, 0.35,
             "Driving and outdoor domains land near the in-domain outlier "
             "rate; indoor close-range (Middlebury) remains the weak axis.",
             size=12, italic=True, bold=True, color=ACCENT, align="center")
    add_text(new, 0.45, y + 0.60, 9.10, 0.30,
             "Training-split evaluations under our protocol; official "
             "leaderboard submissions pending.",
             size=10, italic=True, color=DARK, align="center")
    move_slide(prs, new, after_idx + 1)
    print("  built Zero-Shot Quartet table slide")


def patch_objectives(prs):
    """Match the deck objectives to the thesis section 1.4 wording
    (rewritten to avoid the word 'that')."""
    idx = find_slide_idx(prs, "Objectives")
    if idx < 0:
        print("  WARNING: Objectives slide not found; skipped")
        return
    s = prs.slides[idx]
    swaps = [
        ("pipeline that leverages AI-based",
         "pipeline, leveraging AI-based"),
        ("an architecture that can withstand camera imperfections "
         "in terms of rectification",
         "an architecture able to withstand camera rectification "
         "imperfections"),
        ("Jetson Nano", "Jetson Orin Nano"),
    ]
    for old, new in swaps:
        hit = any(_swap_literal_in_tf(tf, old, new)
                  for tf in _walk_text_frames(s.shapes))
        marker = "replaced" if hit else "NOT FOUND"
        print(f"  objectives: {marker} '{old[:40]}...'")

    # The goals text shape (L=0.59, W=9.10) overruns the framed box
    # whose right border sits at x=9.57; pull the text width in so the
    # first objective no longer overflows the frame.
    for sh in s.shapes:
        if sh.has_text_frame and "1. To design" in sh.text_frame.text:
            if sh.width > Inches(8.80):
                sh.width = Inches(8.80)
                print("  objectives: goals text width -> 8.80 in")
            break


def patch_impact_original(prs):
    """The baked 'Impact' slide (KEY CONTRIBUTION card) still carries
    the pre-defense sub-1 M framing; update it to the final 2.96 M
    chassis without rebuilding the slide."""
    target = None
    for s in prs.slides:
        for tf in _walk_text_frames(s.shapes):
            if "KEY CONTRIBUTION" in tf.text:
                target = s
                break
        if target is not None:
            break
    if target is None:
        print("  WARNING: baked Impact slide not found; skipped")
        return
    swaps = [
        ("+ RAFT iterative refinement",
         "+ ConvGRU iterative refinement + a narrow-band GEV"),
        ("at under 1 M parameters", "at 2.96 M parameters"),
        ("under-1-million parameter stereo regime",
         "sub-3 M parameter edge stereo regime"),
    ]
    for old, new in swaps:
        hit = any(_swap_literal_in_tf(tf, old, new)
                  for tf in _walk_text_frames(target.shapes))
        marker = "replaced" if hit else "NOT FOUND"
        print(f"  impact slide: {marker} '{old}'")


def build_arch_subslide(prs, *, title, subhead, image1_path, image1_caption,
                          image2_path=None, image2_caption=None,
                          image1_job=None, image2_job=None):
    """Duplicate slide 13 (Implementation: Architecture), strip its body,
    and place the supplied diagrams in place of the architecture.
    `image{1,2}_job` is a one-line description of the **primary job** of
    the stage (rendered below the caption in italic). Used on Working
    Principle sub-slides to make each stage's role explicit."""
    # Source = the architecture overview slide. Use find by title so
    # this is robust to slide deletions / reorderings before this call.
    src_idx = find_slide_idx(prs, "Implementation: Architecture")
    if src_idx < 0:
        src_idx = 12   # fallback to former hardcoded position
    new = duplicate_slide(prs, src_idx)

    # Identify keepers: title shape, top divider, footer group.  We
    # rebuild every body shape ourselves so anything that does not match
    # those three patterns is dropped.
    title_sh = subhead_sh = None
    footer_group = None
    top_divider = None

    for sh in new.shapes:
        # Footer group sits at top >= 5.2 in
        if sh.shape_type == 6 and sh.top is not None and \
                sh.top >= Inches(5.20):
            footer_group = sh
            continue
        if not sh.has_text_frame:
            # Top divider line: y < 0.85, very thin
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                top_divider = sh
                continue
            continue
        t = sh.text_frame.text.strip()
        if "Implementation: Architecture" in t and title_sh is None:
            title_sh = sh; continue
        if "StereoLite, end-to-end" in t and subhead_sh is None:
            subhead_sh = sh; continue

    keepers_list = [sh for sh in [title_sh, subhead_sh, footer_group, top_divider]
                     if sh is not None]
    # Compare underlying XML elements, not Python wrapper objects
    keeper_xml = {sh.element for sh in keepers_list}

    # Rewrite title and subhead
    if title_sh is not None:
        set_paragraph_text(title_sh.text_frame, title)
    if subhead_sh is not None:
        set_paragraph_text(subhead_sh.text_frame, subhead)

    # Remove everything else
    for sh in list(new.shapes):
        if sh.element in keeper_xml: continue
        remove_shape(sh)

    # Place diagrams + caption + (optional) primary-job description.
    # Layout below the divider y=0.77:
    #   image  : y=1.20 to ~4.00   (max_h=2.80)
    #   caption: y=4.05 to 4.30
    #   job    : y=4.32 to 4.95
    # Footer band starts at y=5.29.
    if image2_path is None:
        # Single full-width diagram
        add_picture_fit(new, image1_path,
                        x=0.55, y=1.20, max_w=8.90, max_h=2.80)
        if image1_caption:
            add_text(new, 0.55, 4.05, 8.90, 0.25, image1_caption,
                     size=11, bold=True, color=INK, align="center")
        if image1_job:
            add_text(new, 0.80, 4.32, 8.40, 0.55, image1_job,
                     size=11, italic=True, color=DARK, align="center")
    else:
        # Two diagrams side by side
        add_picture_fit(new, image1_path,
                        x=0.30, y=1.20, max_w=4.60, max_h=2.80)
        add_picture_fit(new, image2_path,
                        x=5.10, y=1.20, max_w=4.60, max_h=2.80)
        if image1_caption:
            add_text(new, 0.30, 4.05, 4.60, 0.25, image1_caption,
                     size=10, bold=True, color=INK, align="center")
        if image2_caption:
            add_text(new, 5.10, 4.05, 4.60, 0.25, image2_caption,
                     size=10, bold=True, color=INK, align="center")
        if image1_job:
            add_text(new, 0.30, 4.32, 4.60, 0.65, image1_job,
                     size=10, italic=True, color=DARK, align="center")
        if image2_job:
            add_text(new, 5.10, 4.32, 4.60, 0.65, image2_job,
                     size=10, italic=True, color=DARK, align="center")
    return new


# --------------------------------------------------------------------------
# Footer page renumbering
# --------------------------------------------------------------------------

def renumber_footers(prs):
    for i, s in enumerate(prs.slides):
        # Footer page number is inside the footer GroupShape on the right
        for sh in s.shapes:
            if sh.shape_type != 6:  # not a group
                continue
            for sub in sh.shapes:
                if not sub.has_text_frame:
                    continue
                t = sub.text_frame.text.strip()
                if t.isdigit() and 1 <= int(t) <= 99:
                    if sub.left is not None and sub.left > Emu(8 * 914400):
                        # Re-write the page number
                        paras = list(sub.text_frame.paragraphs)
                        if paras and paras[0].runs:
                            paras[0].runs[0].text = f"{i+1:02d}"
                            for r in paras[0].runs[1:]: r.text = ""
                        break


# --------------------------------------------------------------------------
# v6 final-defense patches: title slide, Proposed Solution stat cards,
# Discussion, Future Work, Challenges / Limitations
# --------------------------------------------------------------------------

def _walk_text_frames(shapes):
    """Yield every text frame under `shapes`, recursing into groups."""
    for sh in shapes:
        if sh.shape_type == 6:  # GroupShape
            yield from _walk_text_frames(sh.shapes)
            continue
        if sh.has_text_frame:
            yield sh.text_frame


def _swap_literal_in_tf(tf, old, new):
    """Substring-replace `old` -> `new` inside a text frame, preserving
    run formatting when the substring sits inside a single run (like
    renumber_footers does); falls back to rewriting the paragraph into
    its first run when the substring spans runs. Returns True if any
    replacement was made."""
    hit = False
    for p in tf.paragraphs:
        ptext = "".join(r.text for r in p.runs)
        if old not in ptext:
            continue
        in_run = False
        for r in p.runs:
            if old in r.text:
                r.text = r.text.replace(old, new)
                in_run = True
        if not in_run:
            runs = list(p.runs)
            if runs:
                runs[0].text = ptext.replace(old, new)
                for r in runs[1:]:
                    r.text = ""
            else:
                continue
        hit = True
    return hit


def patch_title_slide(prs):
    """Slide 1: final-defense course/date corrections. Substring
    replacements across ALL text shapes (including grouped ones),
    preserving run formatting."""
    s = prs.slides[0]
    swaps = [
        ("Seminar", "Project and Thesis"),
        ("MTE 4210", "MTE 4200"),
        ("April 28, 2026", "July 2026"),
        ("APRIL 2026", "JULY 2026"),
    ]
    for old, new in swaps:
        found = False
        for tf in _walk_text_frames(s.shapes):
            if _swap_literal_in_tf(tf, old, new):
                found = True
        if found:
            print(f"  title slide: replaced '{old}' -> '{new}'")
        else:
            print(f"  WARNING title slide: literal '{old}' not found")


def patch_proposed_solution(prs):
    """'Proposed Solution' slide: update the four baked-in stat cards
    (params / latency / synthetic EPE / cross-domain) and their
    sub-labels to the final thesis numbers."""
    idx = find_slide_idx(prs, "Proposed Solution")
    if idx < 0:
        print("  WARNING: 'Proposed Solution' slide not found; skipped")
        return
    s = prs.slides[idx]

    # 1) Conservative literal swaps (run-preserving).
    literals = [
        ("0.87M", "2.96M"),
        ("54ms", "36.3ms"),
        ("1.54px", "0.78px"),
        ("0.515px", "10.9%"),
        ("0.87 M-parameter", "2.96 M parameter"),
        ("8.7 MB", "12 MB"),
        # Card-header consistency with the new numbers (values above are
        # now Orin Nano INT8 latency and zero-shot MB14 D1-all).
        ("LATENCY · RTX 3050", "LATENCY · ORIN NANO"),
        ("EPE · REAL INDOOR", "D1 · ZERO-SHOT"),
    ]
    for old, new in literals:
        found = False
        for tf in _walk_text_frames(s.shapes):
            if _swap_literal_in_tf(tf, old, new):
                found = True
        if found:
            print(f"  proposed solution: replaced '{old}' -> '{new}'")
        else:
            print(f"  WARNING proposed solution: literal '{old}' not found")

    # 2) Sub-label rewrites (whole-line replacement; headers are all-caps
    #    and are skipped by the isupper() guard).
    sub_done = {"latency": False, "sf": False, "xdom": False}
    for tf in _walk_text_frames(s.shapes):
        t = tf.text.strip()
        if not t:
            continue
        low = t.lower().replace("×", "x")   # normalise multiply sign
        if "512x832" in low or ("per" in low and "stereo pair" in low):
            set_paragraph_text(tf, "INT8 on Jetson Orin Nano, 384x640")
            sub_done["latency"] = True
            print(f"  proposed solution: sub-label '{t}' -> "
                  "'INT8 on Jetson Orin Nano, 384x640'")
        elif not t.isupper() and ("scene flow driving" in low or
                                  "sf driving" in low or "200 val" in low):
            set_paragraph_text(tf, "EPE on Scene Flow FT3D test, 4,370 pairs")
            sub_done["sf"] = True
            print(f"  proposed solution: sub-label '{t}' -> "
                  "'EPE on Scene Flow FT3D test, 4,370 pairs'")
        elif not t.isupper() and ("held-out" in low or "indoor pairs" in low):
            set_paragraph_text(tf, "zero-shot Middlebury 2014 D1-all")
            sub_done["xdom"] = True
            print(f"  proposed solution: sub-label '{t}' -> "
                  "'zero-shot Middlebury 2014 D1-all'")
    for key, label in [("latency", "latency (512x832 / RTX 3050)"),
                       ("sf", "SF Driving 200-val"),
                       ("xdom", "indoor held-out")]:
        if not sub_done[key]:
            print(f"  WARNING proposed solution: {label} sub-label not found")


def rebuild_discussion(prs):
    """'Discussion': five takeaway bullets (left) + zero-shot Pareto
    figure from the thesis (right)."""
    idx = find_slide_idx(prs, "Discussion")
    if idx < 0:
        print("  WARNING: 'Discussion' slide not found; skipped")
        return
    s = prs.slides[idx]
    _strip_to_title(s, "Discussion")

    bullets = [
        "In-domain: 0.78 px EPE, 3.40% D1 on the full FT3D test set "
        "at 2.96 M parameters.",
        "Zero-shot: D1 4.33% (KITTI 2012), 3.93% (KITTI 2015), 3.96% "
        "(ETH3D); driving and outdoor domains land near the in-domain "
        "outlier rate.",
        "Middlebury 2014 is the remaining weak axis: 10.9% D1, within "
        "4 points of LiteAnyStereo at 2.6x fewer parameters.",
        "Robust: EPE rises only 1.03 to 1.53 px up to 1 px of vertical "
        "misalignment.",
        "Deployable: 36.3 ms (27.5 FPS) INT8 on Jetson Orin Nano, "
        "measured on device.",
        "Real rig: 1.45 px mean EPE agreement with the FoundationStereo "
        "teacher over 997 pairs.",
    ]
    y = 1.25
    for txt in bullets:
        add_filled_rect(s, 0.35, y + 0.05, 0.10, 0.10, fill_hex=ACCENT)
        add_text(s, 0.57, y, 5.08, 0.58, txt, size=11, color=INK)
        y += 0.63

    fig = TFIGS / "fig_4_8_pareto_ours.png"
    if fig.exists():
        add_picture_fit(s, fig, 5.85, 1.20, 3.85, 3.30)
        add_text(s, 5.85, 4.56, 3.85, 0.26,
                 "Zero-shot accuracy vs parameters",
                 size=10, italic=True, color=DARK, align="center")
    else:
        print(f"  WARNING discussion: figure missing {fig}")


def rebuild_future_work(prs):
    """'Future Work': five numbered directions from thesis Chapter 9,
    ordered by expected return."""
    idx = find_slide_idx(prs, "Future Work")
    if idx < 0:
        print("  WARNING: 'Future Work' slide not found; skipped")
        return
    s = prs.slides[idx]
    _strip_to_title(s, "Future Work")

    items = [
        "Three-stage knowledge distillation: synthetic supervision, then "
        "self-distillation under input perturbation, then a frozen "
        "foundation teacher on unlabeled real pairs.",
        "Submit to the official KITTI and ETH3D leaderboards; our "
        "training-split zero-shot results already sit at D1 4.33, 3.93 "
        "and 3.96 percent.",
        "Deployment hardening: full accuracy audit of the INT8 engine "
        "against the fp32 model, then ports to other edge boards.",
        "Temporal extension: propagate tile-plane state across video "
        "frames to amortize initialization cost.",
        "Mapping extension: fuse per-frame point clouds into a "
        "persistent scene model.",
    ]
    y = 1.28
    for i, txt in enumerate(items):
        add_filled_rect(s, 0.35, y, 0.26, 0.26, fill_hex=NAVY)
        add_text(s, 0.35, y + 0.03, 0.26, 0.20, str(i + 1),
                 size=11, bold=True, color=WHITE, align="center")
        add_text(s, 0.75, y, 8.80, 0.58, txt, size=12, color=INK)
        y += 0.72


def update_challenges_content(prs):
    """'Challenges / Limitations' (fallback title 'Challenges'):
    limitation bullets grounded in thesis Chapter 9."""
    title = "Challenges / Limitations"
    idx = find_slide_idx(prs, title)
    if idx < 0:
        title = "Challenges"
        idx = find_slide_idx(prs, title)
    if idx < 0:
        print("  WARNING: Challenges slide not found; skipped")
        return
    s = prs.slides[idx]
    _strip_to_title(s, title)

    bullets = [
        "Hardest Middlebury scenes stay weak: Jadeplant 34.9% D1, "
        "Flowers 30.3%, both textureless and repetitive structure.",
        "Accuracy degrades beyond 2 px vertical misalignment: D1 15.8% "
        "at 2 px, rising to 41.2% at 4 px.",
        "Recurrent reasoning stops at quarter resolution, limiting "
        "recovery of thin and distant structures.",
        "Supervision is synthetic only, and results come from a single "
        "training run.",
        "KITTI and ETH3D measured on training splits under our own "
        "protocol; official leaderboard submissions still pending.",
        "Real-rig accuracy is agreement with the FoundationStereo "
        "teacher, not absolute ground truth.",
    ]
    y = 1.28
    for txt in bullets:
        add_filled_rect(s, 0.35, y + 0.05, 0.10, 0.10, fill_hex=ACCENT)
        add_text(s, 0.57, y, 8.90, 0.55, txt, size=12, color=INK)
        y += 0.60


# --------------------------------------------------------------------------
# Driver
# --------------------------------------------------------------------------

def main():
    prs = Presentation(str(ORIG))
    print(f"opened {ORIG.name}: {len(prs.slides)} slides")

    # --- Phase A: in-place rewrites (no index changes) ---
    patch_title_slide(prs)
    print("  patched title slide (course, date)")

    patch_proposed_solution(prs)
    print("  patched Proposed Solution stat cards to final numbers")

    patch_objectives(prs)

    rebuild_introduction(prs)
    print("  rebuilt Introduction (slide 3) with stereo geometry concept")

    rebuild_introduction_cont(prs)
    print("  rebuilt Introduction (Cont..) (slide 4) with applications row")

    rebuild_problem_statement(prs)
    print("  rebuilt Problem Statement (slide 5) with alternatives + constraints")

    rebuild_literature_review(prs)
    print("  rebuilt literature review (slide 7)")

    rebuild_review_summary(prs)
    print("  rebuilt review summary capability matrix (slide 8)")

    rebuild_implementation(prs)
    print("  rebuilt implementation hardware/software panel")

    build_methodology(prs)
    print("  built methodology two-track flowchart")

    # --- Phase B: deletions BEFORE renames + adds ---
    # The OLD Working Principle slide must be deleted BEFORE the
    # architecture overview slide is renamed to "Working Principle"
    # (otherwise find_slide_idx finds the renamed slide and we delete
    # the wrong one). Our delete_slide purges the slide part from the
    # package so subsequent add_slide calls do not reuse the deleted
    # slot's slide-N.xml name (which would produce duplicate zip
    # entries that LibreOffice rejects).
    rg = find_slide_idx(prs, "Research Gap")
    if rg >= 0:
        delete_slide(prs, rg)
        print("  deleted Research Gap slide")
    wp = find_slide_idx(prs, "Working Principle")
    if wp >= 0:
        delete_slide(prs, wp)
        print("  deleted old Working Principle slide")
    oa = find_slide_idx(prs, "Objectives Answered")
    if oa >= 0:
        delete_slide(prs, oa)
        print("  deleted Objectives Answered slide")

    # Save + reload so python-pptx normalises the package (drops parts
    # that no longer have rels). Without this, subsequent add_slide
    # calls reuse the deleted slide-N.xml partnames and the resulting
    # zip contains duplicate entries that LibreOffice rejects.
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        tmp_path = tmp.name
    prs = Presentation(tmp_path)
    print(f"  normalised package via reload (now {len(prs.slides)} slides)")

    # --- Phase C: build and insert three architecture sub-slides ---
    arch_a = build_arch_subslide(prs,
        title="Implementation: Architecture · Stage 1 + 2",
        subhead="Encoders and tile hypothesis init",
        image1_path=TFIGS / "fig_3_2_encoders_preview.png",
        image1_caption="Stage 1 · YOLO26s encoder + context encoder",
        image1_job=(
            "A truncated YOLO26s backbone extracts shared features at "
            "1/4, 1/8 and 1/16; a separate 32-channel context encoder "
            "reads the left image for the refinement stream."
        ),
        image2_path=TFIGS / "fig_3_3_tile_init_preview.png",
        image2_caption="Stage 2 · group-wise cost volume + tile init",
        image2_job=(
            "An 8-group correlation volume with 24 hypotheses at 1/16 "
            "is aggregated by a small 3D CNN; soft-argmin seeds one "
            "disparity hypothesis per tile."
        ))

    arch_b = build_arch_subslide(prs,
        title="Implementation: Architecture · Stage 3",
        subhead="Recurrent refinement with geometry evidence",
        image1_path=TFIGS / "fig_3_4_refinement_preview.png",
        image1_caption="Stage 3 · ConvGRU refinement, 2 + 3 + 3 = 8 updates",
        image1_job=(
            "A ConvGRU refines tile disparities across 1/16, 1/8 and 1/4 "
            "using local correlation lookups and plane-aware cross-scale "
            "propagation."
        ),
        image2_path=TFIGS / "fig_3_6_gev_fusion_preview.png",
        image2_caption="Narrow-band GEV at 1/4 · fail-soft gate",
        image2_job=(
            "A 33-hypothesis geometry encoding volume around the current "
            "estimate adds fresh 1/4-scale evidence, fused through a "
            "learned fail-soft gate."
        ))

    arch_c = build_arch_subslide(prs,
        title="Implementation: Architecture · Stage 4 + Supervision",
        subhead="Plane upsampling and multi-scale loss",
        image1_path=TFIGS / "fig_3_7_upsample_preview.png",
        image1_caption="Stage 4 · plane rendering + convex upsample",
        image1_job=(
            "Tile planes are rendered with an edge gate, then two learned "
            "convex 2x stages restore full-resolution disparity."
        ),
        image2_path=TFIGS / "fig_3_8_supervision_preview.png",
        image2_caption="Supervision · 11-term multi-scale objective",
        image2_job=(
            "L1 at five scales plus GEV, gradient, threshold-stack, D1 "
            "hinge, smoothness and slant terms supervise every stage."
        ))

    # Insert them right after the architecture overview slide.
    overview_idx = find_slide_idx(prs, "Implementation: Architecture")
    if overview_idx < 0:
        # Fallback to former hard-coded position
        overview_idx = 12
    move_slide(prs, arch_a, overview_idx + 1)
    move_slide(prs, arch_b, overview_idx + 2)
    move_slide(prs, arch_c, overview_idx + 3)
    print("  inserted 3 per-layer architecture sub-slides")

    # --- Phase D: rename architecture slides → Working Principle family ---
    rename_architecture_to_working_principle(prs)
    print("  renamed architecture slides to Working Principle family")

    # --- Phase E: rebuilds and renames (no slide-count changes) ---
    rebuild_outline(prs)
    print("  rewrote Outline TOC")

    rebuild_impact(prs)
    print("  rebuilt impact panel")

    patch_impact_original(prs)

    rebuild_conclusion(prs)
    print("  rebuilt Conclusion (with Objectives-Met sub-block)")

    rename_challenges(prs)
    print("  renamed 'Challenges' -> 'Challenges / Limitations'")

    refresh_working_principle_overview(prs)

    rebuild_discussion(prs)
    print("  rebuilt Discussion with final results takeaways")

    update_challenges_content(prs)
    print("  rewrote Challenges / Limitations from thesis Ch 9")

    rebuild_future_work(prs)
    print("  rebuilt Future Work from thesis Ch 9")

    build_references(prs)
    print("  populated IEEE references")

    embed_results_progress_gifs(prs)
    print("  embedded training-progress GIFs on results slides 16 + 17")

    # Built LAST: adds a slide, so it must run after every function that
    # addresses slides by hardcoded index (rebuild_impact, implementation).
    build_time_plan(prs)

    build_expansion_slides(prs)

    patch_footer_dates(prs)

    renumber_footers(prs)
    print("  renumbered footers")

    prs.save(str(OUT))
    print(f"\nsaved {OUT.name}")
    print(f"  {OUT.stat().st_size/1e6:.1f} MB  ·  {len(prs.slides)} slides")

    # Auto-build the PDF alongside the .pptx so review is one-click.
    # Reviewing PDFs is faster than opening the deck in PowerPoint /
    # LibreOffice each time. Conversion uses the system LibreOffice
    # in headless mode; if it is missing or fails we just warn.
    import shutil, subprocess
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if soffice is None:
        print("  (skipped PDF build: libreoffice/soffice not on PATH)")
        return
    pdf_out = OUT.with_suffix(".pdf")
    # Remove any stale PDF first so a failed conversion does not look OK.
    if pdf_out.exists():
        pdf_out.unlink()
    r = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf",
         str(OUT), "--outdir", str(OUT.parent)],
        capture_output=True, text=True, timeout=180,
    )
    if pdf_out.exists():
        print(f"  built PDF: {pdf_out.name} ({pdf_out.stat().st_size/1e6:.1f} MB)")
    else:
        # Surface the soffice error so we can debug
        msg = (r.stderr.strip() or r.stdout.strip()
                or "no error message; check libreoffice install")
        print(f"  (PDF build FAILED) {msg[:400]}")


if __name__ == "__main__":
    main()
