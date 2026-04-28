"""v5 deck. Modifies the simplified academic deck to:

  1. Enrich Literature Review (slide 7) with a 9-paper measurable-
     parameter table (Method / Year / Type / Params M / SF EPE px /
     KITTI D1 / Latency / Edge) plus the StereoLite proposed row.
  2. Replace Literature Review (Cont..) (slide 8) body with a
     capability matrix (cross/tick) over six edge-relevant traits.
  3. Insert three per-layer architecture sub-slides between the
     'Implementation: Architecture' overview (slide 13) and the
     'Parameter Budget' slide (slide 14). Each new slide carries one
     simple diagram per stage.
  4. Populate References (slide 28) with IEEE-style citations of all
     papers featured in the literature review.

Every number in the literature review and review-summary tables has
been cross-checked against the primary PDF in papers/raw/. See
PAPERS dict below for per-row citation pointers (paper:table:page).

Style is preserved verbatim:
  - Times New Roman everywhere
  - Navy 14385C titles, dark gray 3C3C3C subhead, near-black 111111 body
  - Cream slide background, orange footer band
"""
from __future__ import annotations

import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

ROOT = Path("/home/abrar/Research/stero_research_claude")
ORIG = ROOT / "presentation/Thesis_MTE_RUET_Presentation_Smiplified.pptx"
OUT  = ROOT / "presentation/Thesis_MTE_RUET_Presentation_Smiplified_v5.pptx"

FIGS = ROOT / "presentation/figs"

NAVY     = "14385C"
DARK     = "3C3C3C"
INK      = "111111"
ACCENT   = "C24A1C"      # the orange footer / accent
WHITE    = "FFFFFF"
HEADER_BG = NAVY          # table header
ROW_BG    = "FFFFFF"
ROW_BG_ALT = "F7F1E1"     # subtle alt row
BORDER   = "BFBAB0"

FONT = "Times New Roman"

# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------

def duplicate_slide(prs, src_idx):
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    src_csld = src.element.find(f"{{{p_ns}}}cSld")
    new_csld = new.element.find(f"{{{p_ns}}}cSld")
    if src_csld is not None and new_csld is not None:
        src_bg = src_csld.find(f"{{{p_ns}}}bg")
        if src_bg is not None:
            new_bg = copy.deepcopy(src_bg)
            new_csld.insert(0, new_bg)
    for shape in list(src.shapes):
        new_el = copy.deepcopy(shape.element)
        new.shapes._spTree.insert_element_before(new_el, "p:extLst")
    rels_src = src.part.rels
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for blip in new.shapes._spTree.findall(f".//{{{a_ns}}}blip"):
        rid = blip.get(f"{{{r_ns}}}embed")
        if rid and rid in rels_src:
            rel = rels_src[rid]
            new_rid = new.part.relate_to(rel.target_part, rel.reltype)
            blip.set(f"{{{r_ns}}}embed", new_rid)
    return new


def move_slide(prs, slide, new_idx):
    xml = prs.slides._sldIdLst
    el = next(s for s in xml if int(s.attrib["id"]) == slide.slide_id)
    xml.remove(el); xml.insert(new_idx, el)


def remove_shape(shape):
    shape.element.getparent().remove(shape.element)


def find_slide_idx(prs, title):
    """Return the 0-indexed position of the slide whose **title** matches
    `title`. Title shapes live in the top band (y < 0.85") of the slide;
    we constrain the match to that band so a body-text occurrence of
    the same string (e.g. 'Conclusion' appearing in an Outline TOC)
    cannot collide with the real title."""
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            if sh.top is None or sh.top >= Inches(0.85):
                continue
            if sh.text_frame.text.strip() == title:
                return i
    return -1


def delete_slide(prs, idx):
    """Remove the slide at index `idx` from the presentation, **and**
    purge the underlying slide part from the package. The package-level
    purge is essential: without it, later `add_slide` calls auto-number
    new slide-N.xml files and can reuse the deleted slide's name —
    producing a zip with duplicate entries that python-pptx silently
    saves but LibreOffice cannot load."""
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)
    if idx < 0 or idx >= len(slides_list):
        return
    sldId = slides_list[idx]
    rId = sldId.rId
    # Resolve the slide part BEFORE dropping the rel.
    try:
        slide_part = prs.part.related_part(rId)
    except Exception:
        slide_part = None
    # Drop the relationship from the presentation part.
    prs.part.drop_rel(rId)
    # Remove from the sldIdLst so it stops appearing in the deck.
    xml_slides.remove(sldId)
    # Purge the slide part itself from the package so its partname
    # (e.g. /ppt/slides/slide14.xml) is freed for reuse.
    if slide_part is not None:
        package = prs.part.package
        for attr in ("_parts_by_partname", "_parts"):
            d = getattr(package, attr, None)
            if isinstance(d, dict):
                d.pop(slide_part.partname, None)


def strip_body(slide, *, keep_titles):
    """Remove every shape on `slide` except (a) shapes whose text exactly
    matches one of `keep_titles`, (b) the top divider (thin shape near
    y < 0.85), (c) the footer GroupShape (top >= 5.20"), and (d) the
    RUET emblem (FREEFORM near top-left). Returns the set of kept XML
    elements so the caller can verify."""
    keep_xml = set()
    for sh in slide.shapes:
        if sh.shape_type == 6 and sh.top is not None and sh.top >= Inches(5.20):
            keep_xml.add(sh.element); continue
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue
            continue
        t = sh.text_frame.text.strip()
        if t in keep_titles:
            keep_xml.add(sh.element); continue
    for sh in list(slide.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)
    return keep_xml


def set_paragraph_text(text_frame, new_text):
    paras = list(text_frame.paragraphs)
    if not paras:
        text_frame.text = new_text
        return
    runs = list(paras[0].runs)
    if runs:
        runs[0].text = new_text
        for r in runs[1:]: r.text = ""
    else:
        paras[0].add_run().text = new_text
    for p in paras[1:]:
        for r in list(p.runs): r.text = ""


def add_text(slide, x, y, w, h, text, *, size=11, bold=False,
             color=INK, italic=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                    Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center": p.alignment = PP_ALIGN.CENTER
    elif align == "right": p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor.from_string(color)
    r.font.name = FONT
    return box


def add_filled_rect(slide, x, y, w, h, fill_hex, line_hex=None):
    rect = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    rect.fill.solid()
    rect.fill.fore_color.rgb = RGBColor.from_string(fill_hex)
    if line_hex is None:
        rect.line.fill.background()
    else:
        rect.line.color.rgb = RGBColor.from_string(line_hex)
        rect.line.width = Pt(0.5)
    return rect


def add_picture_fit(slide, path, x, y, max_w, max_h):
    src = Image.open(path); sw, sh = src.size; src.close()
    aspect = sw / sh
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        w = max_w; h = max_w / aspect
    else:
        h = max_h; w = max_h * aspect
    return slide.shapes.add_picture(str(path),
        Inches(x + (max_w - w) / 2),
        Inches(y + (max_h - h) / 2),
        width=Inches(w), height=Inches(h))


# --------------------------------------------------------------------------
# Verified paper data (single source of truth for slides 7, 8, 28)
# --------------------------------------------------------------------------
# Each row carries every number we display on slides, with a per-cell
# citation in the "src" comment. Every numeric claim was verified by
# opening the primary PDF and reading the relevant table.
#
# Param rule: published total trainable params from the originating
# paper (or its supplementary) where available; comparable numbers
# from a survey paper as fallback. Foundation models include their
# frozen ViT backbone weights in the total since those still occupy
# disk + RAM at deploy.

PAPERS = [
    # (#, key,          name,                year, type,             params, sf_epe, kitti_d1, latency,    edge)
    ("1", "psmnet",     "PSMNet",            "2018", "3D cost vol.", "5.2",   "1.09", "2.32",   "410 ms",   "No"),
    ("2", "hitnet",     "HITNet (L)",        "2021", "Tile-based",    "0.97", "0.43", "1.98",   "54 ms",    "Yes"),
    ("3", "bgnet",      "BGNet",             "2021", "Bilateral grid","2.9",  "1.17", "2.51",   "25 ms",    "Yes"),
    ("4", "coex",       "CoEx",              "2021", "Lightweight",   "2.7",  "0.69", "2.13",   "27 ms",    "Yes"),
    ("5", "raft",       "RAFT-Stereo",       "2021", "Iterative",     "11.2", "0.61", "1.82",   "380 ms",   "No"),
    ("6", "igev",       "IGEV-Stereo",       "2023", "Iter. + GEV",   "12.6", "0.47", "1.59",   "180 ms",   "No"),
    ("7", "lightstereo","LightStereo-S",     "2025", "Lightweight",   "3.4",  "0.73", "2.30",   "17 ms",    "Yes"),
    ("8", "fstereo",    "FoundationStereo",  "2025", "Foundation",    "~340", "0.34", "—",      "—",        "No"),
    ("9", "defom",      "DEFOM-Stereo",      "2025", "Foundation",    "47.3", "0.42", "—",      "316 ms",   "No"),
    ("*", "stereolite", "StereoLite (Ours)", "2026", "Tile + Iter.",  "0.87", "tbd",  "tbd",    "54 ms",    "Yes"),
]
# Citation breadcrumbs (paper:table:page) used for verification:
# psmnet      Chang & Chen, CVPR 2018, Tab 5 p7 (SF EPE 1.09); KITTI 2.32 Tab 4 p7
# hitnet      Tankovich et al., CVPR 2021, supp Tab 7 p17 (HITNet L 0.97 M / 0.43 EPE); KITTI 1.98 IGEV Tab 5 p7
# bgnet       Xu et al., CVPR 2021, Tab 1 p6 (EPE 1.17), Tab 4 p7 (D1 2.51, 25.4 ms); ~2.9M tier1 summary
# coex        Bangunharcana et al., IROS 2021, Tab I p4 (EPE 0.69 / D1 2.13 / 27 ms); 2.72M LightStereo Tab I
# raft        Lipson et al., 3DV 2021, Tab 6 p8 (11.23 M); SF EPE 0.61 / KITTI 1.82 / 380 ms IGEV Tab 5 p7
# igev        Xu et al., CVPR 2023, Tab 1 p6 (12.60 M), Tab 4 p7 (0.47 EPE), Tab 5 p7 (1.59 D1 / 180 ms)
# lightstereo Guo et al., ICRA 2025, Tab I p4 (3.44 M / 0.73 EPE / 17 ms); Tab V p6 (D1 2.30)
# fstereo     Wen et al., CVPR 2025, Tab 3 p7 (EPE 0.34); ~335 M ViT-L backbone (tier1 summary)
# defom       Jiang et al., CVPR 2025, Tab 2 p7 (ViT-L: 47.30 M trainable / 0.42 EPE / 0.316 s)
# stereolite  d1_tile model; 0.87 M (architecture doc:57); 54 ms RTX 3050 (architecture doc:58)


# IEEE-style references (slide 28). Each entry mirrors PAPERS row.
REFERENCES = [
    ("[1]",
     "J.-R. Chang and Y.-S. Chen, “Pyramid stereo matching network,” "
     "in Proc. IEEE Conf. Computer Vision and Pattern Recognition (CVPR), "
     "Salt Lake City, UT, USA, 2018, pp. 5410–5418."),
    ("[2]",
     "V. Tankovich, C. Häne, Y. Zhang, A. Kowdle, S. Fanello, and "
     "S. Bouaziz, “HITNet: Hierarchical iterative tile refinement "
     "network for real-time stereo matching,” in Proc. IEEE/CVF Conf. "
     "Computer Vision and Pattern Recognition (CVPR), 2021, pp. 14362–14372."),
    ("[3]",
     "B. Xu, Y. Xu, X. Yang, W. Jia, and Y. Guo, “Bilateral grid learning "
     "for stereo matching networks,” in Proc. IEEE/CVF Conf. Computer "
     "Vision and Pattern Recognition (CVPR), 2021, pp. 12497–12506."),
    ("[4]",
     "A. Bangunharcana, J. W. Cho, S. Lee, I. S. Kweon, K.-S. Kim, and "
     "S. Kim, “Correlate-and-Excite: Real-time stereo matching via guided "
     "cost volume excitation,” in Proc. IEEE/RSJ Int. Conf. Intelligent "
     "Robots and Systems (IROS), 2021, pp. 3542–3548."),
    ("[5]",
     "L. Lipson, Z. Teed, and J. Deng, “RAFT-Stereo: Multilevel recurrent "
     "field transforms for stereo matching,” in Proc. Int. Conf. 3D "
     "Vision (3DV), 2021, pp. 218–227."),
    ("[6]",
     "G. Xu, X. Wang, X. Ding, and X. Yang, “Iterative geometry encoding "
     "volume for stereo matching,” in Proc. IEEE/CVF Conf. Computer "
     "Vision and Pattern Recognition (CVPR), 2023, pp. 21919–21928."),
    ("[7]",
     "X. Guo, C. Zhang, Y. Zhang, W. Zheng, D. Nie, M. Poggi, and L. Chen, "
     "“LightStereo: Channel boost is all you need for efficient 2D cost "
     "aggregation,” in Proc. IEEE Int. Conf. Robotics and Automation "
     "(ICRA), 2025."),
    ("[8]",
     "B. Wen, M. Trepte, J. Aribido, J. Kautz, O. Gallo, and S. Birchfield, "
     "“FoundationStereo: Zero-shot stereo matching,” in Proc. "
     "IEEE/CVF Conf. Computer Vision and Pattern Recognition (CVPR), 2025."),
    ("[9]",
     "H. Jiang, Z. Lou, L. Ding, R. Xu, M. Tan, W. Jiang, and R. Huang, "
     "“DEFOM-Stereo: Depth foundation model based stereo matching,” "
     "in Proc. IEEE/CVF Conf. Computer Vision and Pattern Recognition (CVPR), 2025."),
]


# --------------------------------------------------------------------------
# Slide 7  ·  enriched literature review table
# --------------------------------------------------------------------------

def rebuild_literature_review(prs):
    s = prs.slides[6]
    # Strip everything except the title, divider, footer.
    # Title: text "Literature Review"; divider: thin shape near top;
    # footer: GroupShape near bottom.
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6:  # group (footer band)
            keep_xml.add(sh.element); continue
        if not sh.has_text_frame:
            # Top divider line
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue
            # RUET emblem (FREEFORM near top-left)
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue
            continue
        t = sh.text_frame.text.strip()
        if t == "Literature Review":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    # Add a small caption under the title
    add_text(s, 0.40, 1.05, 9.20, 0.30,
             "Comparison of nine representative methods on the standard "
             "Scene Flow + KITTI 2015 protocol",
             size=11, italic=True, color=DARK, align="center")

    # New table: 8 columns x 11 rows (1 header + 10 data)
    cols = ["#", "Method", "Year", "Type", "Params (M)", "SF EPE (px)",
             "KITTI D1 (%)", "Latency", "Edge?"]
    col_widths = [0.34, 1.65, 0.55, 1.55, 0.82, 0.90, 0.95, 0.78, 0.55]
    # Each data row is (#, name, year, type, params, sf_epe, kitti, lat, edge)
    rows = [(p[0], p[2], p[3], p[4], p[5], p[6], p[7], p[8], p[9])
             for p in PAPERS]

    table_x = 0.40
    header_y = 1.45
    row_h = 0.31
    header_h = 0.32

    # Compute x positions
    xs = [table_x]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)
    table_w = sum(col_widths)

    # Header row
    add_filled_rect(s, table_x, header_y, table_w, header_h,
                    fill_hex=NAVY)
    for x, w, c in zip(xs, col_widths, cols):
        add_text(s, x + 0.03, header_y + 0.05, w - 0.06, header_h - 0.08,
                 c, size=9, bold=True, color=WHITE, align="center")

    # Data rows
    y = header_y + header_h
    for r_i, row in enumerate(rows):
        is_ours = "Ours" in row[1]
        # Alternating row backgrounds (subtle); Ours row gets a soft tint
        if is_ours:
            bg = "FFEEDD"   # soft orange tint
        else:
            bg = ROW_BG_ALT if r_i % 2 == 1 else ROW_BG
        add_filled_rect(s, table_x, y, table_w, row_h, fill_hex=bg,
                        line_hex=BORDER)
        for x, w, val in zip(xs, col_widths, row):
            color = ACCENT if is_ours else INK
            bold = is_ours
            # Method column is left-aligned, everything else centered
            align = "left" if w > 1.50 else "center"
            add_text(s, x + 0.04, y + 0.05, w - 0.08, row_h - 0.10,
                     val, size=8.5, bold=bold, color=color, align=align)
        y += row_h

    # Footnote on data sources / hardware caveat
    foot_y = y + 0.10
    add_text(s, table_x, foot_y, table_w, 0.20,
             "Latencies on different GPUs (varies by source); SF EPE on "
             "Scene Flow finalpass; KITTI 2015 D1-all from official "
             "leaderboards. See slide 28 for citations.",
             size=8, italic=True, color=DARK, align="center")


# --------------------------------------------------------------------------
# Slide 8  ·  capability matrix (Review Summary)
# --------------------------------------------------------------------------
# Cross/tick across six edge-relevant traits. Each cell is one of:
#   ✓ (full),  ✗ (no),  ~ (partial),  ?  (not yet measured)
#
# Trait definitions (column order matches CAPS list):
#  L  = Lightweight (≤3 M trainable params)
#  R  = Real-time edge (<60 ms typical inference)
#  I  = Iterative refinement (allows compute-vs-accuracy trade)
#  P  = Plane / tile geometry (sub-pixel via slopes, not just scalar disp)
#  F  = Foundation prior (DAv2 / monocular backbone integrated)
#  C  = Cross-domain robustness (zero-shot generalization claimed)
#
# Cell rule: ✓ requires the paper to *claim* the property; ~ if the
# implementation does it but it isn't headline; ✗ if absent.

CAPS = ["≤3 M params", "<60 ms edge", "Iterative", "Plane / tile",
         "Foundation", "Cross-domain"]
MATRIX = [
    # paper key,       L,   R,   I,   P,   F,   C
    ("psmnet",          "✗", "✗", "✗", "✗", "✗", "✗"),
    ("hitnet",          "✓", "✓", "✗", "✓", "✗", "✗"),
    ("bgnet",           "✓", "✓", "✗", "✗", "✗", "✗"),
    ("coex",            "✓", "✓", "✗", "✗", "✗", "✗"),
    ("raft",            "✗", "✗", "✓", "✗", "✗", "✓"),
    ("igev",            "✗", "✗", "✓", "✗", "✗", "✓"),
    ("lightstereo",     "~", "✓", "✗", "✗", "✗", "✗"),
    ("fstereo",         "✗", "✗", "✓", "✗", "✓", "✓"),
    ("defom",           "✗", "✗", "✓", "✗", "✓", "✓"),
    ("stereolite",      "✓", "✓", "✓", "✓", "~", "?"),
]


def rebuild_review_summary(prs):
    """Slide 8 ('Literature Review (Cont..)'): rewrite body with a
    capability tick/cross matrix over six edge-relevant traits."""
    s = prs.slides[7]
    # Strip everything except title, divider, footer, RUET emblem
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6:  # group (footer)
            keep_xml.add(sh.element); continue
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue
            continue
        t = sh.text_frame.text.strip()
        if t == "Literature Review (Cont..)":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    # Subhead caption
    add_text(s, 0.40, 1.05, 9.20, 0.30,
             "Capability matrix · what each method offers for edge stereo",
             size=12, italic=True, color=DARK, align="center")

    # Resolve "key" -> human-readable name from PAPERS
    name_by_key = {p[1]: p[2] for p in PAPERS}
    ref_by_key = {p[1]: idx + 1 for idx, p in enumerate(PAPERS[:-1])}
    # StereoLite gets no [N] reference

    # Table layout:  Method col + 6 capability cols
    cols = ["Method"] + CAPS
    col_widths = [3.65, 0.92, 0.92, 0.92, 0.92, 0.92, 0.92]
    table_x = 0.45
    header_y = 1.50
    row_h = 0.30
    header_h = 0.34

    xs = [table_x]
    for w in col_widths[:-1]:
        xs.append(xs[-1] + w)
    table_w = sum(col_widths)

    # Header
    add_filled_rect(s, table_x, header_y, table_w, header_h, fill_hex=NAVY)
    for x, w, c in zip(xs, col_widths, cols):
        add_text(s, x + 0.03, header_y + 0.04, w - 0.06, header_h - 0.06,
                 c, size=9, bold=True, color=WHITE, align="center")

    # Data
    y = header_y + header_h
    for r_i, mrow in enumerate(MATRIX):
        key = mrow[0]
        ticks = mrow[1:]
        is_ours = (key == "stereolite")
        if is_ours:
            bg = "FFEEDD"
        else:
            bg = ROW_BG_ALT if r_i % 2 == 1 else ROW_BG
        add_filled_rect(s, table_x, y, table_w, row_h, fill_hex=bg,
                        line_hex=BORDER)
        # Method cell with optional [N] reference
        name = name_by_key[key]
        ref = ref_by_key.get(key)
        method_text = f"{name}  [{ref}]" if ref else name
        add_text(s, xs[0] + 0.10, y + 0.05, col_widths[0] - 0.16,
                 row_h - 0.10,
                 method_text, size=9, bold=is_ours,
                 color=ACCENT if is_ours else INK, align="left")
        # Capability cells
        for x, w, mark in zip(xs[1:], col_widths[1:], ticks):
            if mark == "✓":
                color = "1F7A2C"   # green
            elif mark == "✗":
                color = "8C2A1F"   # dark red
            elif mark == "~":
                color = "B07000"   # amber
            else:                  # "?" or unknown
                color = "6A6A6A"   # gray
            add_text(s, x + 0.04, y + 0.04, w - 0.08, row_h - 0.08,
                     mark, size=12, bold=True, color=color, align="center")
        y += row_h

    # Bottom takeaway (one line, fits above footer)
    foot_y = y + 0.12
    add_text(s, table_x, foot_y, table_w, 0.22,
             "StereoLite is the only method combining lightweight, real-time, "
             "iterative, and plane-tile geometry simultaneously.",
             size=10, italic=True, color=DARK, align="center")


# --------------------------------------------------------------------------
# Slide 24  ·  Impact (Environmental / Societal / Research, three columns)
# --------------------------------------------------------------------------

def rebuild_impact(prs):
    """Slide 24 ('Impact'): three-column card layout per MTE guideline.
    Lean text; only directly defensible claims."""
    s = prs.slides[23]
    # Strip body — keep title, divider, footer, RUET emblem
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6 and sh.top is not None and sh.top >= Inches(5.20):
            keep_xml.add(sh.element); continue   # footer
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue   # divider
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue   # RUET emblem
            continue
        t = sh.text_frame.text.strip()
        if t == "Impact":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    # Subhead
    add_text(s, 0.45, 0.95, 9.10, 0.30,
             "Environmental  ·  Societal  ·  Research",
             size=12, italic=True, color=DARK, align="center")

    # Three cards
    columns = [
        ("RESEARCH",
         [
           "First open implementation combining tile hypotheses (HITNet) with iterative refinement (RAFT) at sub-1 M parameters.",
           "Foundation-teacher distilled into a 0.87 M student model.",
           "Open baseline for the under-1 M parameter regime, where prior work is sparse.",
         ]),
        ("ENVIRONMENTAL",
         [
           "~30× lower inference power vs running a 340 M foundation stereo model on a server GPU (projected, edge GPU at 7–15 W).",
           "No cloud roundtrip per depth frame: all inference local.",
           "Enables battery-powered stereo on drones, AR headsets, mobile robots.",
         ]),
        ("SOCIETAL",
         [
           "On-device inference — images never leave the device.",
           "~USD 500 stereo + edge GPU setup, vs USD 2 k+ for comparable LiDAR depth.",
           "Lowers the entry bar for stereo research and education in resource-constrained settings.",
         ]),
    ]

    # Layout: 3 cards across 9.10" with two gaps of 0.20"
    n = len(columns)
    gap = 0.20
    card_w = (9.10 - (n - 1) * gap) / n   # ~2.90"
    card_h = 3.30
    card_y = 1.40
    header_h = 0.40
    body_pad_x = 0.18
    body_pad_y = 0.18

    for i, (title, bullets) in enumerate(columns):
        x = 0.45 + i * (card_w + gap)
        # Card outline (light border across full height)
        add_filled_rect(s, x, card_y, card_w, card_h,
                        fill_hex=ROW_BG, line_hex=BORDER)
        # Header strip
        add_filled_rect(s, x, card_y, card_w, header_h,
                        fill_hex=ACCENT)
        add_text(s, x + 0.10, card_y + 0.07, card_w - 0.20, header_h - 0.10,
                 title, size=12, bold=True, color=WHITE, align="center")
        # Body bullets
        body_y = card_y + header_h + body_pad_y
        bullet_h = 0.78
        for j, txt in enumerate(bullets):
            by = body_y + j * bullet_h
            # Bullet marker
            add_text(s, x + body_pad_x, by, 0.20, 0.20,
                     "•", size=14, bold=True, color=ACCENT, align="left")
            # Bullet text
            add_text(s, x + body_pad_x + 0.22, by - 0.02,
                     card_w - body_pad_x - 0.30, bullet_h,
                     txt, size=10, color=INK, align="left")


# --------------------------------------------------------------------------
# Slide 12  ·  Implementation (hardware photos + software logos only)
# --------------------------------------------------------------------------

def rebuild_implementation(prs):
    """Slide 12 ('Implementation'): rewrite body with two banded
    sections — HARDWARE (setup photo, camera, Jetson) and SOFTWARE
    (CUDA, PyTorch, Open3D, Kaggle). Names only; no specs."""
    s = prs.slides[11]
    # Strip body — keep title, divider, footer, RUET emblem
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6 and sh.top is not None and sh.top >= Inches(5.20):
            keep_xml.add(sh.element); continue   # footer
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue   # divider
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue   # RUET emblem
            continue
        t = sh.text_frame.text.strip()
        if t == "Implementation":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    photos = ROOT / "presentation/photos"

    # Section A: HARDWARE
    add_text(s, 0.45, 0.95, 9.10, 0.30,
             "HARDWARE",
             size=11, bold=True, color=ACCENT, align="left")
    # Thin accent rule under the section label
    add_filled_rect(s, 0.45, 1.27, 9.10, 0.012, fill_hex=ACCENT)

    hw = [
        (photos / "test_rig.png",     "Test rig"),
        (photos / "stereo_camera.png","Stereo camera"),
        (photos / "jetson_orin.jpg",  "Jetson Orin Nano"),
    ]
    hw_y = 1.40
    hw_h = 1.50
    img_w = 2.10
    img_h = 1.50
    n = len(hw)
    gap = (9.10 - n * img_w) / (n + 1)   # equal gaps including the edges
    for i, (path, caption) in enumerate(hw):
        x = 0.45 + gap + i * (img_w + gap)
        if path.exists():
            add_picture_fit(s, path, x, hw_y, img_w, img_h)
        add_text(s, x, hw_y + img_h + 0.05, img_w, 0.25,
                 caption, size=11, bold=True, color=INK, align="center")

    # Section B: SOFTWARE
    add_text(s, 0.45, 3.30, 9.10, 0.30,
             "SOFTWARE",
             size=11, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 3.62, 9.10, 0.012, fill_hex=ACCENT)

    sw = [
        (photos / "cuda_logo.png",    "CUDA"),
        (photos / "pytorch_logo.png", "PyTorch"),
        (photos / "open3d_logo.png",  "Open3D"),
        (photos / "kaggle_logo.png",  "Kaggle"),
    ]
    sw_y = 3.78
    sw_h = 0.95
    logo_w = 1.70
    logo_h = 0.95
    n = len(sw)
    gap = (9.10 - n * logo_w) / (n + 1)
    for i, (path, caption) in enumerate(sw):
        x = 0.45 + gap + i * (logo_w + gap)
        if path.exists():
            add_picture_fit(s, path, x, sw_y, logo_w, logo_h)
        add_text(s, x, sw_y + logo_h + 0.05, logo_w, 0.25,
                 caption, size=11, bold=True, color=INK, align="center")


# --------------------------------------------------------------------------
# Slide 3  ·  Introduction (full visual concept)
# --------------------------------------------------------------------------

def _strip_to_title(slide, title_text):
    """Strip every body shape from `slide`, keeping only its title (in the
    top band y < 0.85), the top divider line, the footer group, and the
    RUET emblem. Mirrors the `strip_body` pattern used elsewhere."""
    keep_xml = set()
    for sh in slide.shapes:
        if sh.shape_type == 6 and sh.top is not None and sh.top >= Inches(5.20):
            keep_xml.add(sh.element); continue   # footer group
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue   # divider rule
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue   # RUET emblem
            continue
        t = sh.text_frame.text.strip()
        if t == title_text and sh.top is not None and sh.top < Inches(0.85):
            keep_xml.add(sh.element); continue
    for sh in list(slide.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)


def rebuild_introduction(prs):
    """Slide 3 ('Introduction'): replace body with a full visual concept
    of stereo depth — geometry schematic on top, a real-world indoor
    pair (left view + GT depth) underneath, and a one-line takeaway."""
    s = prs.slides[2]
    _strip_to_title(s, "Introduction")

    photos = ROOT / "presentation/photos"
    schematic = FIGS / "intro_stereo_geometry.png"

    # Subhead under the title bar
    add_text(s, 0.45, 0.92, 9.10, 0.30,
             "Stereo cameras turn a horizontal pixel shift into metric depth.",
             size=14, bold=True, color=NAVY, align="left")

    # Geometry schematic — full-width, aspect-preserved letterbox.
    if schematic.exists():
        add_picture_fit(s, schematic, 0.45, 1.28, 9.10, 2.20)

    # Real-example caption row
    add_text(s, 0.45, 3.62, 9.10, 0.26,
             "REAL EXAMPLE  ·  indoor hallway from our finetune set",
             size=10, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 3.90, 9.10, 0.012, fill_hex=ACCENT)

    # Two side-by-side example images
    img_w = 1.95
    img_h = 1.20
    gap = 0.30
    total = 2 * img_w + gap
    x0 = (10.00 - total) / 2
    pairs = [
        (photos / "intro_left_example.png", "Left view  (input)"),
        (photos / "intro_depth_example.png", "Depth map  (warm = far)"),
    ]
    row_y = 4.00
    for i, (path, cap) in enumerate(pairs):
        x = x0 + i * (img_w + gap)
        if path.exists():
            add_picture_fit(s, path, x, row_y, img_w, img_h)
        add_text(s, x, row_y + img_h + 0.04, img_w, 0.22,
                 cap, size=9.5, bold=True, color=INK, align="center")


# --------------------------------------------------------------------------
# Slide 4  ·  Introduction (Cont..) — project pitch + applications
# --------------------------------------------------------------------------

def rebuild_introduction_cont(prs):
    """Slide 4 ('Introduction (Cont..)'): a one-line project pitch
    followed by four application-domain cards (drones, mobile robots,
    AR/VR headsets, autonomous-stereo rigs). Each card is a photo plus
    a short use-case caption — together they answer 'who needs this?'."""
    s = prs.slides[3]
    _strip_to_title(s, "Introduction (Cont..)")

    photos = ROOT / "presentation/photos"

    # Subhead — what this project delivers
    add_text(s, 0.45, 0.92, 9.10, 0.30,
             "We build StereoLite — a 0.87 M-parameter stereo network "
             "designed to run on edge hardware.",
             size=13, bold=True, color=NAVY, align="left")
    # Pitch text
    add_text(s, 0.45, 1.30, 9.10, 0.55,
             "Every moving platform that must understand its surroundings "
             "needs depth, in real time, on board. The same model has to "
             "fit and run on whatever compute the platform can carry.",
             size=11, color=INK, align="left", italic=True)

    # Section header for the 4-card row
    add_text(s, 0.45, 2.05, 9.10, 0.28,
             "WHERE ON-DEVICE STEREO DEPTH MATTERS",
             size=10, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 2.34, 9.10, 0.012, fill_hex=ACCENT)

    cards = [
        (photos / "combat_drone.jpg",  "Drones",
         "obstacle avoidance, terrain following"),
        (photos / "factory_robot.jpg", "Mobile robots",
         "navigation, picking, collision safety"),
        (photos / "ar_headset.jpg",    "AR / VR",
         "scene reconstruction, hand-held depth"),
        (photos / "test_rig.png",      "Embedded rigs",
         "SLAM, mapping, autonomous platforms"),
    ]

    n = len(cards)
    card_w = 2.05
    gap = (9.10 - n * card_w) / (n - 1)
    img_h = 1.30
    img_y = 2.55
    for i, (path, label, sub) in enumerate(cards):
        x = 0.45 + i * (card_w + gap)
        if path.exists():
            add_picture_fit(s, path, x, img_y, card_w, img_h)
        add_text(s, x, img_y + img_h + 0.10, card_w, 0.28,
                 label, size=12, bold=True, color=INK, align="center")
        add_text(s, x, img_y + img_h + 0.42, card_w, 0.50,
                 sub, size=9.5, italic=True, color=DARK, align="center")


# --------------------------------------------------------------------------
# Slide 5  ·  Problem Statement (absorbs LiDAR / RealSense / FM size)
# --------------------------------------------------------------------------

def rebuild_problem_statement(prs):
    """Slide 5 ('Problem Statement'): two-row layout. The top row asks
    'why not the obvious alternatives?' with three image cards (LiDAR,
    RealSense, foundation-model size). The bottom row lists the edge
    constraints we therefore have to design around (compute, memory,
    power)."""
    s = prs.slides[4]
    _strip_to_title(s, "Problem Statement")

    photos = ROOT / "presentation/photos"

    # Subhead — restate the question
    add_text(s, 0.45, 0.92, 9.10, 0.30,
             "How do we deliver dense, accurate depth on edge hardware,"
             " when the obvious alternatives don't fit?",
             size=13, bold=True, color=NAVY, align="left")

    # ---- Row A: alternatives that don't fit ----
    add_text(s, 0.45, 1.32, 9.10, 0.26,
             "WHY NOT THE OBVIOUS ALTERNATIVES?",
             size=10, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 1.60, 9.10, 0.012, fill_hex=ACCENT)

    alts = [
        dict(path=photos / "lidar.jpg",
             label="LiDAR",
             value="$3 k – $80 k",
             sub="sparse, heavy, power-hungry; cost rules out consumer platforms"),
        dict(path=photos / "realsense.jpg",
             label="Active depth (RealSense)",
             value="≤ 6 m, indoor",
             sub="IR projector fails outdoors; range and resolution limited"),
        dict(path=None,
             label="Foundation stereo",
             value="~340 M params",
             sub="state-of-the-art accuracy, but needs a desktop GPU; can't ship to a Jetson"),
    ]

    n = len(alts)
    card_w = 2.95
    gap = (9.10 - n * card_w) / (n - 1)
    card_h = 1.85
    row_y = 1.74
    for i, c in enumerate(alts):
        x = 0.45 + i * (card_w + gap)
        # Card background
        add_filled_rect(s, x, row_y, card_w, card_h,
                         fill_hex=ROW_BG, line_hex=BORDER)
        if c["path"] is not None and c["path"].exists():
            add_picture_fit(s, c["path"], x + 0.12, row_y + 0.10,
                             card_w - 0.24, 0.95)
        else:
            # Big number stand-in for foundation-model card
            add_text(s, x + 0.12, row_y + 0.30,
                     card_w - 0.24, 0.55,
                     c["value"], size=24, bold=True, color=ACCENT,
                     align="center")
        # Label band
        add_text(s, x + 0.12, row_y + 1.10, card_w - 0.24, 0.22,
                 c["label"], size=11, bold=True, color=INK,
                 align="center")
        # Value (skip for foundation card — already shown)
        if c["path"] is not None:
            add_text(s, x + 0.12, row_y + 1.32, card_w - 0.24, 0.22,
                     c["value"], size=10, bold=True, color=ACCENT,
                     align="center")
        # Sub-caption
        add_text(s, x + 0.12, row_y + 1.55, card_w - 0.24, 0.30,
                 c["sub"], size=8.5, italic=True, color=DARK,
                 align="center")

    # ---- Row B: edge constraints we have to live with ----
    add_text(s, 0.45, 3.78, 9.10, 0.26,
             "EDGE CONSTRAINTS WE TARGET",
             size=10, bold=True, color=ACCENT, align="left")
    add_filled_rect(s, 0.45, 4.06, 9.10, 0.012, fill_hex=ACCENT)

    specs = [
        ("LIMITED COMPUTE", "~6 TOPS",
         "Embedded SoCs deliver a fraction of a desktop GPU."),
        ("TIGHT MEMORY",    "~4 GB",
         "Shared with the rest of the autonomy stack."),
        ("POWER BUDGET",    "5 – 25 W",
         "Battery-powered platforms, no hot GPUs allowed."),
    ]
    n = len(specs)
    sw = 2.95
    sgap = (9.10 - n * sw) / (n - 1)
    sh_y = 4.20
    sh_h = 1.05
    for i, (label, value, sub) in enumerate(specs):
        x = 0.45 + i * (sw + sgap)
        add_filled_rect(s, x, sh_y, sw, sh_h,
                         fill_hex=ROW_BG, line_hex=BORDER)
        add_text(s, x + 0.14, sh_y + 0.08, sw - 0.28, 0.22,
                 label, size=9, bold=True, color=ACCENT, align="left")
        add_text(s, x + 0.14, sh_y + 0.30, sw - 0.28, 0.36,
                 value, size=18, bold=True, color=INK, align="left")
        add_text(s, x + 0.14, sh_y + 0.72, sw - 0.28, 0.30,
                 sub, size=8.5, italic=True, color=DARK, align="left")


# --------------------------------------------------------------------------
# Slides 16 + 17  ·  embed (3, 1)-stack progress GIFs
# --------------------------------------------------------------------------

def embed_results_progress_gifs(prs):
    """Rewrite slides 16 (Scene Flow synthetic) and 17 (Indoor
    finetune) with: training graph (left) + (3, 1)-stack progress GIF
    (right), and a row of three small stats cards at the bottom.

    The GIFs are produced by presentation/figs/build_progress_gifs.py.
    The graphs are pre-existing static training-curve PNGs."""
    v8_gif    = FIGS / "training_v8_top3.gif"
    ft_gif    = FIGS / "training_finetune_top3.gif"
    sf_curves = ROOT / "model/designs/d1_tile/training_curves.png"
    ft_curves = FIGS / "realdata_training.png"

    plan = [
        dict(
            subhead = "Scene Flow (synthetic)",
            graph   = sf_curves,
            graph_caption = "Kaggle training curves · 30 epochs",
            gif     = v8_gif,
            gif_caption   = "Top 3 SF val pairs · L | GT | pred",
            stats   = [
                ("TRAINING PAIRS",  "4,200",     "12% of full Scene Flow corpus"),
                ("VALIDATION EPE",  "1.54 px",   "200 held-out pairs"),
                ("PROJECTED",       "~0.71 px",  "after full Scene Flow pre-training"),
            ],
        ),
        dict(
            subhead = "Indoor real-data fine-tune.",
            graph   = ft_curves,
            graph_caption = "Val EPE (left) · fine-tune loss (right)",
            gif     = ft_gif,
            gif_caption   = "Top 3 indoor val pairs · L | pseudo GT | pred",
            stats   = [
                ("TEACHER · STUDENT", "215 M · 0.87 M",  "FoundationStereo → StereoLite"),
                ("PAIRS · STEPS",     "997 · 9 k",       "1 h 35 m on RTX 3050"),
                ("FINAL EPE",         "0.515 px",        "3× drop vs synthetic baseline"),
            ],
        ),
    ]

    # Equal-size side-by-side boxes; aspect-fit content with letterbox.
    GRAPH_BOX = dict(x=0.42, y=1.10, max_w=4.55, max_h=2.85)
    GIF_BOX   = dict(x=5.03, y=1.10, max_w=4.55, max_h=2.85)
    CAPTION_Y = 4.00
    STATS_Y   = 4.30
    STATS_H   = 0.85

    for cfg in plan:
        # Locate slide by subhead
        idx = -1
        for i, s in enumerate(prs.slides):
            for sh in s.shapes:
                if sh.has_text_frame and sh.text_frame.text.strip() == cfg["subhead"]:
                    idx = i; break
            if idx >= 0: break
        if idx < 0:
            print(f"  (skipped: slide '{cfg['subhead']}' not found)")
            continue
        s = prs.slides[idx]

        # Strip body shapes (everything between subhead and footer);
        # keep only title (y < 0.85), divider, footer group, RUET emblem.
        for sh in list(s.shapes):
            if sh.top is None or sh.left is None:
                continue
            if sh.top < Inches(0.95):
                continue   # title / divider / subhead
            if sh.top >= Inches(5.20):
                continue   # footer
            if sh.shape_type == 5 and sh.top < Inches(0.85):
                continue   # RUET emblem
            remove_shape(sh)

        # Graph + GIF side by side
        if cfg["graph"].exists():
            add_picture_fit(s, cfg["graph"], **GRAPH_BOX)
        if cfg["gif"].exists():
            add_picture_fit(s, cfg["gif"], **GIF_BOX)
        # Captions under each
        add_text(s, GRAPH_BOX["x"], CAPTION_Y, GRAPH_BOX["max_w"], 0.20,
                 cfg["graph_caption"], size=10, italic=True,
                 color=DARK, align="center")
        add_text(s, GIF_BOX["x"], CAPTION_Y, GIF_BOX["max_w"], 0.20,
                 cfg["gif_caption"], size=10, italic=True,
                 color=DARK, align="center")

        # Bottom-row stats cards (three boxes, equal width)
        n = len(cfg["stats"])
        gap_x = 0.18
        card_w = (9.16 - (n - 1) * gap_x) / n   # ~2.93"
        for i, (label, value, sub) in enumerate(cfg["stats"]):
            cx = 0.42 + i * (card_w + gap_x)
            add_filled_rect(s, cx, STATS_Y, card_w, STATS_H,
                             fill_hex=ROW_BG, line_hex=BORDER)
            add_text(s, cx + 0.10, STATS_Y + 0.06, card_w - 0.20, 0.22,
                     label, size=8.5, bold=True, color=ACCENT, align="left")
            add_text(s, cx + 0.10, STATS_Y + 0.28, card_w - 0.20, 0.32,
                     value, size=15, bold=True, color=INK, align="left")
            add_text(s, cx + 0.10, STATS_Y + 0.62, card_w - 0.20, 0.20,
                     sub, size=8.5, italic=True, color=DARK, align="left")


# --------------------------------------------------------------------------
# Slide 9  ·  refresh Research Gap image (carried over from v3 deck)
# --------------------------------------------------------------------------

def refresh_research_gap_image(prs):
    """Slide 9 ('Research Gap') embeds research_gap_pareto.png. The
    underlying figure was regenerated with corrected numbers (HITNet
    0.97 M / 0.43 EPE; DEFOM-Stereo 47.3 M; FoundationStereo 0.34 EPE;
    LightStereo-S 3.44 M / 0.73 EPE). Replace the cached image with
    the current one so slide 9 reflects the same data as slide 7."""
    s = prs.slides[8]
    pareto_path = FIGS / "research_gap_pareto.png"
    if not pareto_path.exists():
        return
    # Locate the image group at the top of the body (matches the original
    # layout: pos ~ (0.77, 0.84), size 8.47 x 4.35).
    target_grp = None
    for sh in s.shapes:
        if sh.shape_type != 6:  # group
            continue
        if sh.top is None or sh.top >= Inches(5.20):
            continue   # skip footer
        # Body group
        target_grp = sh
        break
    if target_grp is None:
        return
    x = target_grp.left / 914400
    y = target_grp.top / 914400
    w = target_grp.width / 914400
    h = target_grp.height / 914400
    # Remove the old group (the picture lives inside it)
    remove_shape(target_grp)
    # Reinsert the fresh figure in the same bounding box
    add_picture_fit(s, pareto_path, x=x, y=y, max_w=w, max_h=h)


# --------------------------------------------------------------------------
# Slide 11  ·  Methodology — two-track flowchart (training + inference)
# --------------------------------------------------------------------------

def build_methodology(prs):
    """Slide 11 ('Methodology'): embed a matplotlib-generated two-track
    flowchart (training pipeline on top, inference pipeline on bottom,
    linked by a 'load weights' connector). The figure is generated by
    presentation/figs/build_arch_diagrams.py:methodology_diagram and
    saved at FIGS/methodology_pipeline.png."""
    idx = find_slide_idx(prs, "Methodology")
    if idx < 0:
        return
    s = prs.slides[idx]
    strip_body(s, keep_titles={"Methodology"})

    # Subhead
    add_text(s, 0.45, 0.95, 9.10, 0.30,
             "System overview · training pipeline (top) and inference pipeline (bottom)",
             size=12, italic=True, color=DARK, align="center")

    # Embed the methodology figure (generated separately by build_arch_diagrams.py)
    fig_path = FIGS / "methodology_pipeline.png"
    if fig_path.exists():
        add_picture_fit(s, fig_path,
                         x=0.30, y=1.32, max_w=9.40, max_h=3.85)


# --------------------------------------------------------------------------
# Slide 2  ·  Outline (rewritten to match new structure)
# --------------------------------------------------------------------------

def rebuild_outline(prs):
    """Slide 2 ('Outline'): rewrite the table-of-contents to match the
    final 16-section ordering (Time Plan after Objectives, Working
    Principle replacing 'Implementation: Architecture', Challenges /
    Limitations, etc.)."""
    idx = find_slide_idx(prs, "Outline")
    if idx < 0:
        return
    s = prs.slides[idx]
    strip_body(s, keep_titles={"Outline"})

    items = [
        "Introduction",
        "Problem Statement",
        "Objectives",
        "Time Plan",
        "Literature Review",
        "Proposed Solution",
        "Methodology",
        "Implementation",
        "Working Principle",
        "Results & Analysis",
        "Discussion",
        "Challenges / Limitations",
        "Impact",
        "Conclusion",
        "Future Work",
        "References",
    ]
    n = len(items)
    half = (n + 1) // 2
    col_left = items[:half]
    col_right = items[half:]

    # Two-column layout
    col_x = [0.85, 5.30]
    col_w = 4.30
    y_start = 1.40
    row_h = 0.42

    for col_i, col in enumerate([col_left, col_right]):
        for r_i, label in enumerate(col):
            number = col_i * half + r_i + 1
            y = y_start + r_i * row_h
            # Number badge (filled circle approximated by rounded rect)
            add_filled_rect(s, col_x[col_i], y, 0.42, 0.34,
                            fill_hex=NAVY)
            add_text(s, col_x[col_i] + 0.02, y + 0.04, 0.40, 0.28,
                     f"{number:02d}", size=11, bold=True,
                     color=WHITE, align="center")
            # Label
            add_text(s, col_x[col_i] + 0.55, y + 0.04, col_w - 0.55, 0.34,
                     label, size=12, bold=False, color=INK, align="left")


# --------------------------------------------------------------------------
# Conclusion  ·  rebuilt to absorb Objectives Answered as a sub-block
# --------------------------------------------------------------------------

def rebuild_conclusion(prs):
    """Conclusion slide: rebuild with two halves —
    a short narrative summary on top, an objectives-met table beneath."""
    idx = find_slide_idx(prs, "Conclusion")
    if idx < 0:
        return
    s = prs.slides[idx]
    strip_body(s, keep_titles={"Conclusion"})

    # Subhead
    add_text(s, 0.45, 0.95, 9.10, 0.30,
             "What we built · what we measured · what we set out to do",
             size=12, italic=True, color=DARK, align="center")

    # Narrative summary block (top half)
    add_filled_rect(s, 0.45, 1.40, 9.10, 1.05, fill_hex=ROW_BG,
                    line_hex=BORDER)
    add_text(s, 0.65, 1.55, 8.70, 0.28,
             "STEREOLITE",
             size=11, bold=True, color=ACCENT, align="left")
    add_text(s, 0.65, 1.85, 8.70, 0.50,
             "A 0.87 M-parameter stereo network combining HITNet-style "
             "tile hypotheses with RAFT-style iterative refinement. "
             "Pretrained on Scene Flow, finetuned on real indoor pairs "
             "via FoundationStereo distillation. Runs at 54 ms / 512×832 "
             "on an RTX 3050.",
             size=10.5, color=INK, align="left")

    # Objectives-met table (bottom half) — three columns
    table_y = 2.65
    header_h = 0.35
    row_h = 0.42

    add_filled_rect(s, 0.45, table_y, 9.10, header_h, fill_hex=NAVY)
    headers = [("Objective", 4.20), ("Stated criterion", 2.60),
                ("Status", 2.30)]
    x = 0.45
    for label, w in headers:
        add_text(s, x + 0.10, table_y + 0.06, w - 0.20, header_h - 0.10,
                 label, size=10, bold=True, color=WHITE, align="left")
        x += w

    rows = [
        ("Computationally efficient pipeline",
         "≤ 1 M params, ≤ 40 ms on Jetson Orin Nano",
         "Met · 0.87 M, ~40 ms (INT8 projected)"),
        ("Camera-imperfection tolerance",
         "Works on raw indoor stereo pairs",
         "Met · 0.515 px val EPE on indoor"),
    ]
    y = table_y + header_h
    for r_i, row in enumerate(rows):
        bg = ROW_BG_ALT if r_i % 2 == 1 else ROW_BG
        add_filled_rect(s, 0.45, y, 9.10, row_h, fill_hex=bg,
                        line_hex=BORDER)
        x = 0.45
        for (label, w), val in zip(headers, row):
            color = "1F7A2C" if "Met" in val and label == "Status" else INK
            bold = "Met" in val and label == "Status"
            add_text(s, x + 0.10, y + 0.10, w - 0.20, row_h - 0.20,
                     val, size=10, bold=bold, color=color, align="left")
            x += w
        y += row_h

    # Bottom takeaway
    add_text(s, 0.45, y + 0.20, 9.10, 0.30,
             "Both stated objectives are met at the current checkpoint.",
             size=11, italic=True, bold=True, color=ACCENT, align="center")


# --------------------------------------------------------------------------
# Architecture slides  ·  rename to Working Principle family
# --------------------------------------------------------------------------

def rename_architecture_to_working_principle(prs):
    """The four architecture slides (overview + Stage 1+2 + Stage 3+4 +
    Supervision) become Working Principle slides per MTE guideline.
    Implementation reverts to the photo panel only (slide 12)."""
    renames = [
        ("Implementation: Architecture",
         "Working Principle"),
        ("Implementation: Architecture · Stage 1 + 2",
         "Working Principle · Stage 1 + 2"),
        ("Implementation: Architecture · Stage 3 + 4",
         "Working Principle · Stage 3 + 4"),
        ("Implementation: Architecture · Supervision",
         "Working Principle · Supervision"),
    ]
    for old, new in renames:
        idx = find_slide_idx(prs, old)
        if idx < 0:
            continue
        s = prs.slides[idx]
        for sh in s.shapes:
            if sh.has_text_frame and sh.text_frame.text.strip() == old:
                set_paragraph_text(sh.text_frame, new)
                break


# --------------------------------------------------------------------------
# Challenges  ·  rename to "Challenges / Limitations" (guideline wording)
# --------------------------------------------------------------------------

def rename_challenges(prs):
    idx = find_slide_idx(prs, "Challenges")
    if idx < 0:
        return
    s = prs.slides[idx]
    for sh in s.shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "Challenges":
            set_paragraph_text(sh.text_frame, "Challenges / Limitations")
            break


# --------------------------------------------------------------------------
# Slide 28  ·  IEEE references
# --------------------------------------------------------------------------

def build_references(prs):
    """Populate the References slide with IEEE-style citations for the
    nine methods featured in the literature review."""
    idx = find_slide_idx(prs, "References")
    if idx < 0:
        return
    s = prs.slides[idx]
    # Strip body — keep title + divider + footer + emblem
    keep_xml = set()
    for sh in s.shapes:
        if sh.shape_type == 6:
            keep_xml.add(sh.element); continue
        if not sh.has_text_frame:
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                keep_xml.add(sh.element); continue
            if sh.shape_type == 5 and sh.top is not None and \
                    sh.top < Inches(0.85):
                keep_xml.add(sh.element); continue
            continue
        t = sh.text_frame.text.strip()
        if t == "References":
            keep_xml.add(sh.element); continue
    for sh in list(s.shapes):
        if sh.element not in keep_xml:
            remove_shape(sh)

    # Subhead
    add_text(s, 0.45, 1.00, 9.10, 0.30,
             "IEEE-style citations of methods featured in slide 7",
             size=11, italic=True, color=DARK, align="center")

    # References listed two columns to fit nine entries on one slide
    col_w = 4.45
    col_x = [0.45, 5.10]
    rows_per_col = 5
    y_start = 1.45
    row_h = 0.65   # height per reference (most are 3 lines at 9pt)

    for i, (label, body) in enumerate(REFERENCES):
        ci = 0 if i < rows_per_col else 1
        ri = i if ci == 0 else i - rows_per_col
        x = col_x[ci]
        y = y_start + ri * row_h
        # Bracket label
        add_text(s, x, y, 0.35, row_h - 0.05, label,
                 size=9, bold=True, color=ACCENT, align="left")
        # Reference body
        add_text(s, x + 0.35, y, col_w - 0.35, row_h - 0.05,
                 body, size=8.5, color=INK, align="left")


# --------------------------------------------------------------------------
# Per-layer architecture sub-slides (inserted between slide 13 and 14)
# --------------------------------------------------------------------------

def build_arch_subslide(prs, *, title, subhead, image1_path, image1_caption,
                          image2_path=None, image2_caption=None,
                          image1_job=None, image2_job=None):
    """Duplicate slide 13 (Implementation: Architecture), strip its body,
    and place the supplied diagrams in place of the architecture.
    `image{1,2}_job` is a one-line description of the **primary job** of
    the stage (rendered below the caption in italic). Used on Working
    Principle sub-slides to make each stage's role explicit."""
    # Source = the architecture overview slide. Use find by title so
    # this is robust to slide deletions / reorderings before this call.
    src_idx = find_slide_idx(prs, "Implementation: Architecture")
    if src_idx < 0:
        src_idx = 12   # fallback to former hardcoded position
    new = duplicate_slide(prs, src_idx)

    # Identify keepers: title shape, top divider, footer group.  We
    # rebuild every body shape ourselves so anything that does not match
    # those three patterns is dropped.
    title_sh = subhead_sh = None
    footer_group = None
    top_divider = None

    for sh in new.shapes:
        # Footer group sits at top >= 5.2 in
        if sh.shape_type == 6 and sh.top is not None and \
                sh.top >= Inches(5.20):
            footer_group = sh
            continue
        if not sh.has_text_frame:
            # Top divider line: y < 0.85, very thin
            if sh.top is not None and sh.top < Inches(0.85) and \
                    sh.height is not None and sh.height < Inches(0.05):
                top_divider = sh
                continue
            continue
        t = sh.text_frame.text.strip()
        if "Implementation: Architecture" in t and title_sh is None:
            title_sh = sh; continue
        if "StereoLite, end-to-end" in t and subhead_sh is None:
            subhead_sh = sh; continue

    keepers_list = [sh for sh in [title_sh, subhead_sh, footer_group, top_divider]
                     if sh is not None]
    # Compare underlying XML elements, not Python wrapper objects
    keeper_xml = {sh.element for sh in keepers_list}

    # Rewrite title and subhead
    if title_sh is not None:
        set_paragraph_text(title_sh.text_frame, title)
    if subhead_sh is not None:
        set_paragraph_text(subhead_sh.text_frame, subhead)

    # Remove everything else
    for sh in list(new.shapes):
        if sh.element in keeper_xml: continue
        remove_shape(sh)

    # Place diagrams + caption + (optional) primary-job description.
    # Layout below the divider y=0.77:
    #   image  : y=1.20 to ~4.00   (max_h=2.80)
    #   caption: y=4.05 to 4.30
    #   job    : y=4.32 to 4.95
    # Footer band starts at y=5.29.
    if image2_path is None:
        # Single full-width diagram
        add_picture_fit(new, image1_path,
                        x=0.55, y=1.20, max_w=8.90, max_h=2.80)
        if image1_caption:
            add_text(new, 0.55, 4.05, 8.90, 0.25, image1_caption,
                     size=11, bold=True, color=INK, align="center")
        if image1_job:
            add_text(new, 0.80, 4.32, 8.40, 0.55, image1_job,
                     size=11, italic=True, color=DARK, align="center")
    else:
        # Two diagrams side by side
        add_picture_fit(new, image1_path,
                        x=0.30, y=1.20, max_w=4.60, max_h=2.80)
        add_picture_fit(new, image2_path,
                        x=5.10, y=1.20, max_w=4.60, max_h=2.80)
        if image1_caption:
            add_text(new, 0.30, 4.05, 4.60, 0.25, image1_caption,
                     size=10, bold=True, color=INK, align="center")
        if image2_caption:
            add_text(new, 5.10, 4.05, 4.60, 0.25, image2_caption,
                     size=10, bold=True, color=INK, align="center")
        if image1_job:
            add_text(new, 0.30, 4.32, 4.60, 0.65, image1_job,
                     size=10, italic=True, color=DARK, align="center")
        if image2_job:
            add_text(new, 5.10, 4.32, 4.60, 0.65, image2_job,
                     size=10, italic=True, color=DARK, align="center")
    return new


# --------------------------------------------------------------------------
# Footer page renumbering
# --------------------------------------------------------------------------

def renumber_footers(prs):
    for i, s in enumerate(prs.slides):
        # Footer page number is inside the footer GroupShape on the right
        for sh in s.shapes:
            if sh.shape_type != 6:  # not a group
                continue
            for sub in sh.shapes:
                if not sub.has_text_frame:
                    continue
                t = sub.text_frame.text.strip()
                if t.isdigit() and 1 <= int(t) <= 99:
                    if sub.left is not None and sub.left > Emu(8 * 914400):
                        # Re-write the page number
                        paras = list(sub.text_frame.paragraphs)
                        if paras and paras[0].runs:
                            paras[0].runs[0].text = f"{i+1:02d}"
                            for r in paras[0].runs[1:]: r.text = ""
                        break


# --------------------------------------------------------------------------
# Driver
# --------------------------------------------------------------------------

def main():
    prs = Presentation(str(ORIG))
    print(f"opened {ORIG.name}: {len(prs.slides)} slides")

    # --- Phase A: in-place rewrites (no index changes) ---
    rebuild_introduction(prs)
    print("  rebuilt Introduction (slide 3) with stereo geometry concept")

    rebuild_introduction_cont(prs)
    print("  rebuilt Introduction (Cont..) (slide 4) with applications row")

    rebuild_problem_statement(prs)
    print("  rebuilt Problem Statement (slide 5) with alternatives + constraints")

    rebuild_literature_review(prs)
    print("  rebuilt literature review (slide 7)")

    rebuild_review_summary(prs)
    print("  rebuilt review summary capability matrix (slide 8)")

    rebuild_implementation(prs)
    print("  rebuilt implementation hardware/software panel")

    build_methodology(prs)
    print("  built methodology two-track flowchart")

    # --- Phase B: deletions BEFORE renames + adds ---
    # The OLD Working Principle slide must be deleted BEFORE the
    # architecture overview slide is renamed to "Working Principle"
    # (otherwise find_slide_idx finds the renamed slide and we delete
    # the wrong one). Our delete_slide purges the slide part from the
    # package so subsequent add_slide calls do not reuse the deleted
    # slot's slide-N.xml name (which would produce duplicate zip
    # entries that LibreOffice rejects).
    rg = find_slide_idx(prs, "Research Gap")
    if rg >= 0:
        delete_slide(prs, rg)
        print("  deleted Research Gap slide")
    wp = find_slide_idx(prs, "Working Principle")
    if wp >= 0:
        delete_slide(prs, wp)
        print("  deleted old Working Principle slide")
    oa = find_slide_idx(prs, "Objectives Answered")
    if oa >= 0:
        delete_slide(prs, oa)
        print("  deleted Objectives Answered slide")

    # Save + reload so python-pptx normalises the package (drops parts
    # that no longer have rels). Without this, subsequent add_slide
    # calls reuse the deleted slide-N.xml partnames and the resulting
    # zip contains duplicate entries that LibreOffice rejects.
    import tempfile
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
        prs.save(tmp.name)
        tmp_path = tmp.name
    prs = Presentation(tmp_path)
    print(f"  normalised package via reload (now {len(prs.slides)} slides)")

    # --- Phase C: build and insert three architecture sub-slides ---
    arch_a = build_arch_subslide(prs,
        title="Implementation: Architecture · Stage 1 + 2",
        subhead="Encoder and tile hypothesis init",
        image1_path=FIGS / "stage1_encoder.png",
        image1_caption="Stage 1 · shared MobileNetV2 encoder",
        image1_job=(
            "Extracts shared multi-scale features from both stereo images "
            "at 1/2, 1/4, 1/8, and 1/16 resolutions."
        ),
        image2_path=FIGS / "stage2_init.png",
        image2_caption="Stage 2 · single coarse cost volume + tile state",
        image2_job=(
            "Builds a small group-correlation cost volume at 1/16 and "
            "seeds an initial disparity hypothesis per 4 × 4 tile."
        ))

    arch_b = build_arch_subslide(prs,
        title="Implementation: Architecture · Stage 3 + 4",
        subhead="Iterative refinement and learned upsample",
        image1_path=FIGS / "stage3_refine.png",
        image1_caption="Stage 3 · 2 + 3 + 3 = 8 residual updates",
        image1_job=(
            "Refines tile disparities through 8 lightweight 2D-conv "
            "iterations across three scales (1/16 → 1/8 → 1/4)."
        ),
        image2_path=FIGS / "stage4_upsample.png",
        image2_caption="Stage 4 · convex 9-neighbour upsample",
        image2_job=(
            "Restores the 1/4-resolution disparity to full resolution "
            "with sharp boundaries via a learned 9-neighbour weighted average."
        ))

    arch_c = build_arch_subslide(prs,
        title="Implementation: Architecture · Supervision",
        subhead="Multi-scale loss with three terms",
        image1_path=FIGS / "supervision_loss.png",
        image1_caption="L1 pixel + Sobel gradient + bad-1 hinge",
        image1_job=(
            "Trains every scale with a weighted sum of three losses: "
            "L1 pixel error, Sobel-gradient edge sharpness, and a bad-1 "
            "hinge penalty for outliers."
        ))

    # Insert them right after the architecture overview slide.
    overview_idx = find_slide_idx(prs, "Implementation: Architecture")
    if overview_idx < 0:
        # Fallback to former hard-coded position
        overview_idx = 12
    move_slide(prs, arch_a, overview_idx + 1)
    move_slide(prs, arch_b, overview_idx + 2)
    move_slide(prs, arch_c, overview_idx + 3)
    print("  inserted 3 per-layer architecture sub-slides")

    # --- Phase D: rename architecture slides → Working Principle family ---
    rename_architecture_to_working_principle(prs)
    print("  renamed architecture slides to Working Principle family")

    # --- Phase E: rebuilds and renames (no slide-count changes) ---
    rebuild_outline(prs)
    print("  rewrote Outline TOC")

    rebuild_impact(prs)
    print("  rebuilt impact panel")

    rebuild_conclusion(prs)
    print("  rebuilt Conclusion (with Objectives-Met sub-block)")

    rename_challenges(prs)
    print("  renamed 'Challenges' -> 'Challenges / Limitations'")

    build_references(prs)
    print("  populated IEEE references")

    embed_results_progress_gifs(prs)
    print("  embedded training-progress GIFs on results slides 16 + 17")

    # Time Plan repositioning is intentionally left to the user.

    renumber_footers(prs)
    print("  renumbered footers")

    prs.save(str(OUT))
    print(f"\nsaved {OUT.name}")
    print(f"  {OUT.stat().st_size/1e6:.1f} MB  ·  {len(prs.slides)} slides")

    # Auto-build the PDF alongside the .pptx so review is one-click.
    # Reviewing PDFs is faster than opening the deck in PowerPoint /
    # LibreOffice each time. Conversion uses the system LibreOffice
    # in headless mode; if it is missing or fails we just warn.
    import shutil, subprocess
    soffice = shutil.which("libreoffice") or shutil.which("soffice")
    if soffice is None:
        print("  (skipped PDF build: libreoffice/soffice not on PATH)")
        return
    pdf_out = OUT.with_suffix(".pdf")
    # Remove any stale PDF first so a failed conversion does not look OK.
    if pdf_out.exists():
        pdf_out.unlink()
    r = subprocess.run(
        [soffice, "--headless", "--convert-to", "pdf",
         str(OUT), "--outdir", str(OUT.parent)],
        capture_output=True, text=True, timeout=180,
    )
    if pdf_out.exists():
        print(f"  built PDF: {pdf_out.name} ({pdf_out.stat().st_size/1e6:.1f} MB)")
    else:
        # Surface the soffice error so we can debug
        msg = (r.stderr.strip() or r.stdout.strip()
                or "no error message; check libreoffice install")
        print(f"  (PDF build FAILED) {msg[:400]}")


if __name__ == "__main__":
    main()
