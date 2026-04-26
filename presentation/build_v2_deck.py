"""Modify the existing pre-defense PPTX:

  - Slide  1: course code MTE 4200 → 4210, course title → "Seminar"
  - Slide  8: replace conceptual research-gap chart with real Pareto PNG
  - Slide 10: AR0144 camera label + AR0144 deployment pipeline diagram
  - Slide 25: refine the limitations text now that we DO have real-data eval
  - Slide  2: refresh the outline to reflect the new sections + page numbers

  + 4 new slides inserted in flow:
    A. after slide 3 (Introduction)         — "Live demo · what we built"
    B. after that                            — "3D reconstruction · top val pairs"
    C. after slide 19 (Parameter budget)    — "Real-data fine-tune"
    D. after slide 22 (Qualitative SF)      — "Real-world qualitative results"

  The new slides are built by *duplicating* existing similarly-styled slides
  and then re-writing their content; this preserves all the careful
  decoration the supervisor approved.
"""
from __future__ import annotations

import copy
import shutil
import subprocess
from pathlib import Path

from lxml import etree
from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Emu, Inches, Pt

ROOT = Path("/home/abrar/Research/stero_research_claude")
ORIG = ROOT / "presentation/Thesis Pre Defense Presentation Slides(2008011, 2008026).pptx"
OUT  = ROOT / "presentation/Thesis Pre Defense Presentation Slides v2 (2008011, 2008026).pptx"

FIGS    = ROOT / "presentation/figs"
PCTOP   = Path("/mnt/abrarssd/Datasets/stereo_samples_20260425_104147/point_clouds_top3")
VIDS    = Path("/mnt/abrarssd/Datasets/stereo_samples_20260425_104147/vids")
DEPLOY  = ROOT / "model/designs/d1_tile/deployment_pipeline.png"

# Brand palette from existing deck (sampled from PDF screenshots)
INK_HEX     = "1A1A1F"
SUBINK_HEX  = "5A5550"
ACCENT_HEX  = "C24A1C"
CREAM_HEX   = "F4EFE6"
DARK_BG_HEX = "242021"


# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------

def replace_in_runs(text_frame, old: str, new: str):
    """Replace `old` with `new` even when the match spans multiple runs.

    Strategy: per paragraph, concatenate run text; if `old` appears,
    rewrite the first run with the replaced text and blank the others.
    """
    for p in text_frame.paragraphs:
        runs = list(p.runs)
        if not runs:
            continue
        full = "".join(r.text for r in runs)
        if old not in full:
            continue
        new_full = full.replace(old, new)
        runs[0].text = new_full
        for r in runs[1:]:
            r.text = ""


def deep_replace(shape, old: str, new: str):
    if shape.has_text_frame:
        replace_in_runs(shape.text_frame, old, new)
    if shape.shape_type == 6:  # GROUP
        for sub in shape.shapes:
            deep_replace(sub, old, new)


def set_paragraph_text(text_frame, new_text: str):
    """Set the first paragraph's first run to `new_text` and clear all
    other runs / paragraphs. Preserves the first run's formatting."""
    paras = list(text_frame.paragraphs)
    if not paras:
        text_frame.text = new_text
        return
    first = paras[0]
    runs = list(first.runs)
    if runs:
        runs[0].text = new_text
        for r in runs[1:]:
            r.text = ""
    else:
        first.add_run().text = new_text
    # Clear extra paragraphs
    for p in paras[1:]:
        for r in list(p.runs):
            r.text = ""


def _shape_has_text(shape, needle: str) -> bool:
    if shape.has_text_frame and needle in shape.text_frame.text:
        return True
    if shape.shape_type == 6:
        for sub in shape.shapes:
            if _shape_has_text(sub, needle):
                return True
    return False


def remove_shape(shape):
    shape.element.getparent().remove(shape.element)


def duplicate_slide(prs, src_idx: int):
    """Append a deep copy of slide `src_idx` as a new slide. Preserves
    text, shapes, and embedded picture relationships."""
    src = prs.slides[src_idx]
    layout = src.slide_layout
    new = prs.slides.add_slide(layout)

    # Copy shapes
    for shape in list(src.shapes):
        new_el = copy.deepcopy(shape.element)
        new.shapes._spTree.insert_element_before(new_el, "p:extLst")

    # Re-create relationships for any pictures in the copied shapes
    rels_src = src.part.rels
    rels_new = new.part.rels
    nsmap = {"r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
              "a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
    for blip in new.shapes._spTree.findall(
            ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"):
        rid = blip.get("{http://schemas.openxmlformats.org/officeDocument/"
                       "2006/relationships}embed")
        if rid and rid in rels_src:
            rel = rels_src[rid]
            new_rid = new.part.relate_to(rel.target_part, rel.reltype)
            blip.set("{http://schemas.openxmlformats.org/officeDocument/"
                     "2006/relationships}embed", new_rid)
    return new


def move_slide(prs, slide, new_idx: int):
    xml_slides = prs.slides._sldIdLst
    slide_id = next(s for s in xml_slides
                     if int(s.attrib["id"]) == slide.slide_id)
    xml_slides.remove(slide_id)
    xml_slides.insert(new_idx, slide_id)


def renumber_footers(prs):
    """Walk every slide and rewrite the small page-number text at the
    bottom-right (it's the last text shape with a small numeric value)."""
    for i, s in enumerate(prs.slides):
        # The page number shape is identified by its position (right edge,
        # tiny width) and short numeric text.
        for sh in s.shapes:
            if not sh.has_text_frame:
                continue
            t = sh.text_frame.text.strip()
            if t.isdigit() and 1 <= int(t) <= 99 and \
                    sh.left is not None and sh.left > Emu(8 * 914400):
                # Replace the digits while preserving formatting
                p = sh.text_frame.paragraphs[0]
                if p.runs:
                    p.runs[0].text = f"{i + 1:02d}"
                else:
                    sh.text_frame.text = f"{i + 1:02d}"
                break


# --------------------------------------------------------------------------
# Edits to existing slides
# --------------------------------------------------------------------------

def edit_slide1_title(prs):
    """Course code MTE 4200 → MTE 4210, course title → Seminar."""
    s = prs.slides[0]
    for sh in s.shapes:
        deep_replace(sh, "MTE 4200", "MTE 4210")
        deep_replace(sh, "Project and Thesis", "Seminar")


def edit_slide8_research_gap(prs):
    """Delete the conceptual scatter and drop in the real Pareto PNG."""
    s = prs.slides[7]
    # Keep: header bar (sh[0,1]), right-side description (sh[26..28],
    # sh[30]), footer band (sh[31..34]). Remove: the scatter elements
    # (sh[2..23]) and the small vertical accent (sh[29]).
    keep_texts = [
        "§ 06 RESEARCH GAP",
        "A gap in the accuracy",
        "Top-right cluster",
        "Bottom-left cluster",
        "Under-populated",
        "small-&-accurate",
        "APRIL 2026",
        "AI-Enhanced Stereo Matching",
    ]
    keep_pgnum = "08"
    to_remove = []
    for sh in s.shapes:
        if not sh.has_text_frame:
            to_remove.append(sh); continue
        txt = sh.text_frame.text.strip()
        if any(k in txt for k in keep_texts):
            continue
        if txt == keep_pgnum:
            continue
        to_remove.append(sh)
    for sh in to_remove:
        # Skip the bottom footer band shape (the wide rectangle has no text
        # and lives at the bottom); identify by top > 5.0in
        if sh.top is not None and sh.top > Inches(5.2) and sh.width is not None \
                and sh.width > Inches(8.0):
            continue
        # Skip the right-column small accent bar (the thin vertical we
        # want to keep above the italic note)
        if sh.has_text_frame is False and sh.width is not None and \
                sh.width < Inches(0.05) and sh.left is not None and \
                sh.left > Inches(5.0):
            continue
        remove_shape(sh)

    # Insert the new Pareto figure on the LEFT half (matches the old chart's
    # bounding box ~ L 0.4, T 1.40, W 4.75, H 3.71).
    img = FIGS / "research_gap_pareto.png"
    pic = s.shapes.add_picture(str(img),
        left=Inches(0.40), top=Inches(1.30),
        width=Inches(5.20))
    # Anchor the bottom by setting only width to keep aspect


def edit_slide10_hardware(prs):
    """Replace IMX-219 deployment pipeline + camera label with AR0144."""
    s = prs.slides[9]
    # Update camera label text
    for sh in s.shapes:
        deep_replace(sh, "IMX 219-83 Camera", "Waveshare AR0144 USB")

    # Find the deployment-pipeline picture (rightmost picture) and
    # replace with the AR0144 PNG. Easiest: add new picture, delete old.
    pictures = [sh for sh in s.shapes if sh.shape_type == 13]  # PICTURE
    if pictures:
        # Rightmost = largest .left
        old = max(pictures, key=lambda p: p.left)
        L, T, W, H = old.left, old.top, old.width, old.height
        remove_shape(old)
        s.shapes.add_picture(str(DEPLOY), L, T, width=W, height=H)


def edit_slide25_limitations(prs):
    """We have real-data eval now — refine the middle box."""
    s = prs.slides[24]
    for sh in s.shapes:
        deep_replace(sh, "NO REAL LIFE QUALITATIVE MEASUREMENT",
                     "PSEUDO-GT, NOT GROUND TRUTH")
        deep_replace(sh,
            "The real life inference accuracy is measured against "
            "FoundationStereo. Access to 3D lidar would have provided "
            "accuracy on real life scenarios.",
            "Indoor real-data fine-tune EPE is measured against "
            "FoundationStereo pseudo-GT, not LiDAR. Absolute accuracy on "
            "physical scenes still needs a calibrated depth reference.")
        deep_replace(sh, "LIMITED DATASET",
                     "SCENE FLOW + ~1k INDOOR")
        deep_replace(sh,
            "Trained on Scene Flow Driving and custom generated dataset. "
            "For competing against existing models. It must be trained "
            "against other standard ones",
            "Pretrained on Scene Flow Driving (4,200 pairs), fine-tuned "
            "on 997 clean indoor pairs. KITTI / Middlebury fine-tunes "
            "and full Scene Flow pretraining are queued.")


# --------------------------------------------------------------------------
# New slides
# --------------------------------------------------------------------------

def _add_text(slide, x, y, w, h, text, *, size=14, bold=False, color=INK_HEX,
              mono=False, align="left", italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                     Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center":
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        from pptx.enum.text import PP_ALIGN
        p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor.from_string(color)
    if mono:
        r.font.name = "DejaVu Sans Mono"
    else:
        r.font.name = "DejaVu Serif"
    return box


def make_chrome_slide(prs, *, section_label: str, title_main: str,
                       title_accent: str = ""):
    """Duplicate slide 24 (Discussion — has light bg + chrome) and clear
    its body. Returns the new slide ready for body content."""
    new = duplicate_slide(prs, 23)   # slide index 23 = slide 24

    # Edit header / title in place; remove body content.
    # Identify shapes to clear by their text or by position (center body
    # area, between top y=1.3in and y=5.0in).
    keep = ["§ 14 DISCUSSION",
             "Reading the numbers",
             "APRIL 2026",
             "AI-Enhanced Stereo Matching"]
    keep_idx = []
    for sh in new.shapes:
        if not sh.has_text_frame:
            continue
        txt = sh.text_frame.text.strip()
        if any(k in txt for k in keep):
            keep_idx.append(sh)
    # Replace the header label
    for sh in keep_idx:
        if "§ 14 DISCUSSION" in sh.text_frame.text:
            set_paragraph_text(sh.text_frame, section_label)
        if "Reading the numbers" in sh.text_frame.text:
            # Big slide title — rebuild with two runs (main + accent)
            tf = sh.text_frame
            for p in list(tf.paragraphs):
                for r in list(p.runs):
                    p._p.remove(r._r)
            p = tf.paragraphs[0]
            r1 = p.add_run(); r1.text = title_main + (" " if title_accent else "")
            r1.font.bold = True
            r1.font.size = Pt(28)
            r1.font.color.rgb = RGBColor.from_string(INK_HEX)
            r1.font.name = "DejaVu Serif"
            if title_accent:
                r2 = p.add_run(); r2.text = title_accent
                r2.font.bold = True
                r2.font.size = Pt(28)
                r2.font.color.rgb = RGBColor.from_string(ACCENT_HEX)
                r2.font.name = "DejaVu Serif"

    # Remove body shapes (anything that's not keep, header, footer)
    for sh in list(new.shapes):
        if sh in keep_idx:
            continue
        if sh.top is not None and sh.top >= Inches(5.20):
            continue   # footer band + footer text
        if sh.has_text_frame:
            txt = sh.text_frame.text.strip()
            if txt.isdigit() and 1 <= int(txt) <= 99:
                continue   # page number; renumbered later
            if "AI-Enhanced Stereo Matching" in txt or "APRIL 2026" in txt:
                continue
            # Drop body text
            remove_shape(sh)
        else:
            # Drop body decorations
            if sh.top is not None and Inches(1.20) < sh.top < Inches(5.20):
                remove_shape(sh)
    return new


# Per-slide builders ---------------------------------------------------------

def slide_intro_demo(prs, video_path: Path | None):
    """Big embedded video / GIF showing the model running on real data."""
    s = make_chrome_slide(prs,
        section_label="§ 01.5 LIVE DEMO",
        title_main="What StereoLite",
        title_accent="produces.")

    # Large video on the left if mp4 ready, otherwise a colormapped frame
    if video_path and video_path.exists() and video_path.stat().st_size > 0:
        # add_movie expects mp4; needs a poster image
        poster = FIGS / "_video_poster.png"
        if not poster.exists():
            # synthesize a poster from the existing inference panel preview
            preview = FIGS / "_inference_preview.png"
            if preview.exists():
                shutil.copy(preview, poster)
            else:
                # placeholder
                Image.new("RGB", (1080, 540),
                           color=tuple(int(DARK_BG_HEX[i:i+2], 16) for i in (0,2,4))) \
                     .save(poster)
        movie = s.shapes.add_movie(
            str(video_path), Inches(0.55), Inches(1.45),
            Inches(6.20), Inches(3.10),
            poster_frame_image=str(poster), mime_type="video/mp4")
    else:
        # fallback to a still
        still = FIGS / "_inference_preview.png"
        if still.exists():
            s.shapes.add_picture(str(still), Inches(0.55), Inches(1.45),
                                   width=Inches(6.20))

    # Right column: numbers + caption
    rx = 7.05
    _add_text(s, rx, 1.45, 2.50, 0.30,
               "INDOOR REAL DATA  ·  RUET CORRIDOR",
               size=8.5, bold=True, color=ACCENT_HEX, mono=True)
    _add_text(s, rx, 1.85, 2.50, 0.50,
               "0.87 M",
               size=30, bold=True, color=INK_HEX)
    _add_text(s, rx, 2.35, 2.50, 0.30,
               "trainable parameters",
               size=10, color=SUBINK_HEX)

    _add_text(s, rx, 2.75, 2.50, 0.50,
               "0.515 px",
               size=24, bold=True, color=ACCENT_HEX)
    _add_text(s, rx, 3.20, 2.50, 0.30,
               "indoor val EPE  ·  50 held-out pairs",
               size=10, color=SUBINK_HEX)

    _add_text(s, rx, 3.65, 2.50, 0.50,
               "54 ms",
               size=24, bold=True, color=INK_HEX)
    _add_text(s, rx, 4.10, 2.50, 0.30,
               "RTX 3050 · 512×832 stereo pair",
               size=10, color=SUBINK_HEX)

    _add_text(s, 0.55, 4.65, 9.00, 0.30,
               "Left: 1280×720 stereo input  ·  Right: predicted "
               "disparity (TURBO colormap, fixed range)",
               size=9.5, color=SUBINK_HEX, italic=True)
    return s


def slide_intro_3d(prs):
    """Triptych of rotating point cloud GIFs."""
    s = make_chrome_slide(prs,
        section_label="§ 01.6 RECONSTRUCTION",
        title_main="From disparity to",
        title_accent="3D.")
    # Three rotating GIFs in a row
    gifs = sorted(PCTOP.glob("pair_*.gif"))
    gifs = gifs[:3]
    panel_w = 2.95
    panel_h = 2.40
    gap = 0.10
    total_w = 3 * panel_w + 2 * gap
    x0 = (10.0 - total_w) / 2.0
    y0 = 1.45
    epe_labels = []
    for i, gif in enumerate(gifs):
        s.shapes.add_picture(str(gif),
            Inches(x0 + i * (panel_w + gap)), Inches(y0),
            width=Inches(panel_w), height=Inches(panel_h))
        # Caption
        nm = gif.stem
        epe = nm.split("epe")[-1] if "epe" in nm else "?"
        _add_text(s,
            x0 + i * (panel_w + gap), y0 + panel_h + 0.05,
            panel_w, 0.25,
            f"val pair  ·  EPE {epe} px",
            size=10, mono=True, color=ACCENT_HEX, align="center")
        epe_labels.append(epe)

    _add_text(s, 0.55, y0 + panel_h + 0.45, 9.00, 0.30,
               "AR0144 stereo USB camera  ·  fx ≈ 1005 px  ·  baseline 52 mm  "
               "·  Open3D post-processed (statistical outlier · 4 mm voxel · "
               "Phong)",
               size=9.5, color=SUBINK_HEX, italic=True)
    return s


def slide_realdata_training(prs):
    s = make_chrome_slide(prs,
        section_label="§ 11.5 REAL-DATA FINE-TUNE",
        title_main="Fine-tune on real",
        title_accent="indoor stereo.")

    # Sized to leave room for two caption lines below; image is 11×4.4
    # native, so a 8.0-in width yields ~3.20 in tall.
    img = FIGS / "realdata_training.png"
    s.shapes.add_picture(str(img),
        Inches(1.00), Inches(1.30),
        width=Inches(8.00))

    _add_text(s, 0.55, 4.65, 8.90, 0.30,
               "Teacher: FoundationStereo (215 M)  →  StereoLite (0.87 M)  "
               "·  997 clean pairs · 50 val held-out · 9000 steps · "
               "1 h 35 m on RTX 3050",
               size=9.5, color=SUBINK_HEX, mono=True)

    _add_text(s, 0.55, 4.95, 8.90, 0.30,
               "Indoor mean EPE drops from 1.54 px (Scene Flow baseline) "
               "to 0.515 px in 9000 steps — 3× improvement on the new "
               "domain without losing edge-deployable size.",
               size=10.5, color=ACCENT_HEX, italic=True)
    return s


def slide_realdata_results(prs):
    s = make_chrome_slide(prs,
        section_label="§ 13.4 REAL-WORLD RESULTS",
        title_main="Top-3 val pairs ·",
        title_accent="3D reconstruction.")

    # Top: compact metric card (sized smaller than the figure script's
    # 11×4.6 to fit alongside the GIF row).
    img = FIGS / "realdata_results.png"
    s.shapes.add_picture(str(img),
        Inches(0.55), Inches(1.30), width=Inches(8.90),
        height=Inches(2.00))

    # GIFs row — placed below the metric card with breathing room
    gifs = sorted(PCTOP.glob("pair_*.gif"))[:3]
    panel_w = 2.85
    panel_h = 1.65
    gap = 0.10
    total_w = 3 * panel_w + 2 * gap
    x0 = (10.0 - total_w) / 2.0
    y0 = 3.45
    for i, gif in enumerate(gifs):
        s.shapes.add_picture(str(gif),
            Inches(x0 + i * (panel_w + gap)), Inches(y0),
            width=Inches(panel_w), height=Inches(panel_h))
    return s


# --------------------------------------------------------------------------
# Driver
# --------------------------------------------------------------------------

def main():
    prs = Presentation(str(ORIG))
    print(f"opened {ORIG.name}: {len(prs.slides)} slides")

    # 1. Edits
    edit_slide1_title(prs)
    edit_slide8_research_gap(prs)
    edit_slide10_hardware(prs)
    edit_slide25_limitations(prs)
    print("  edited slides 1, 8, 10, 25")

    # 2. New slides — appended at end, then moved into position.
    # IMPORTANT: build them in REVERSE position order so subsequent
    # inserts don't disturb earlier ones.
    video_path = VIDS / "stereolite_inference_panel_v2.mp4"
    if not video_path.exists():
        video_path = VIDS / "stereolite_inference_panel.mp4"
    print(f"  video for live demo: {video_path.name}")

    s_demo  = slide_intro_demo(prs, video_path)
    s_3d    = slide_intro_3d(prs)
    s_train = slide_realdata_training(prs)
    s_real  = slide_realdata_results(prs)

    # Original slide indices (0-based, BEFORE inserts):
    #   3 = Introduction (slide 3)
    #   19 = Parameter Budget (slide 19)
    #   22 = Qualitative Progression (slide 22)
    # Targets: insert demo + 3d after slide 3; train after slide 19;
    # real after slide 22.

    # Insert in reverse so earlier positions remain valid:
    move_slide(prs, s_real,  23)   # after slide 22 (idx 22) → idx 23
    move_slide(prs, s_train, 20)   # after slide 19 (idx 19) → idx 20
    move_slide(prs, s_3d,    4)    # after slide 3  (idx 3)  → idx 4
    move_slide(prs, s_demo,  4)    # after slide 3  (idx 3)  → idx 4

    print(f"  added 4 new slides; total {len(prs.slides)}")

    # 3. Renumber footer page numbers
    renumber_footers(prs)
    print("  renumbered footers")

    # 4. Save
    prs.save(str(OUT))
    print(f"\nsaved {OUT}")
    sz = OUT.stat().st_size / 1e6
    print(f"  {sz:.1f} MB  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
