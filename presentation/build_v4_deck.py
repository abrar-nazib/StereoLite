"""v4 rebuild of the pre-defense deck. Addresses the second feedback round:

  - Use the ORIGINAL fonts (Georgia / Courier New / Arial), not DejaVu.
  - Intro spans 3 slides numbered 1, 1.1, 1.2 (not 1.1/1.2/1.3).
    * 1   = minimal title slide; depth perception on edge is necessary;
            three photos (combat drone, AR headset, factory robot) +
            stereo camera photo. NO mention of StereoLite.
    * 1.1 = "Why edge stereo is hard" (was the old 1.3 content).
    * 1.2 = 10s indoor inference clip + one rotating 3D mesh GIF +
            first introduction of the StereoLite name.
  - Problem (2): refined title ("at workable accuracy"); compute /
    memory / power blocks instead of prose; no combat-drone mention.
  - Objectives (3): tree structure  Objective N → indented sub bullets.
  - Lit review (4): DARK background.
  - Slide 4.1: still labelled LITERATURE REVIEW (not RESEARCH GAP).
    Pareto chart with side legend (no label overlap), text rewritten
    to match the new chart.
  - Proposed (5): drop "IDEA 0X" labels; mention real-data EPE.
  - Methodology (6): block diagram, not text columns.
  - Implementation (7): DARK; equipment + logo grid.
    7.1 = "Implementation: Model Architecture"
    7.2/7.3/7.4 = "Implementation: <stage names>".
  - NEW slide 9.5: explicit point-by-point objectives-met answers.
  - Discussion (10): photos.
  - No em dashes, no double dashes anywhere in body text.
"""
from __future__ import annotations

import copy
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

# --------------------------------------------------------------------------
# Paths
# --------------------------------------------------------------------------
ROOT = Path("/home/abrar/Research/stero_research_claude")
ORIG = ROOT / "presentation/Thesis Pre Defense Presentation Slides(2008011, 2008026).pptx"
OUT  = ROOT / "presentation/Thesis Pre Defense Presentation Slides v4 (2008011, 2008026).pptx"

FIGS    = ROOT / "presentation/figs"
PHOTOS  = ROOT / "presentation/photos"
PCTOP   = Path("/mnt/abrarssd/Datasets/stereo_samples_20260425_104147/point_clouds_top3")
VIDS    = Path("/mnt/abrarssd/Datasets/stereo_samples_20260425_104147/vids")

DEPLOY  = ROOT / "model/designs/d1_tile/deployment_pipeline.png"
ARCH_FULL = ROOT / "model/designs/d1_tile/stereolite_arch.png"
TRAIN_GIF = ROOT / "model/benchmarks/stereolite_finetune_indoor_20260426-171158/training_progression.gif"
SF_PROGRESSION = ROOT / "model/designs/d1_tile/progress_grid.png"
SF_TRAINING_CURVES = ROOT / "model/designs/d1_tile/training_curves.png"
INTRO_VIDEO = VIDS / "intro_clip_10s.mp4"
DEMO_VIDEO_60S = VIDS / "stereolite_inference_panel_40s_720.mp4"

# --------------------------------------------------------------------------
# Palette (matches existing deck after sampling the original PDF)
# --------------------------------------------------------------------------
INK_HEX     = "1A1A1F"
SUBINK_HEX  = "5A5550"
ACCENT_HEX  = "C24A1C"
SOFTACC_HEX = "D9826A"
CREAM_HEX   = "F4EFE6"
DARK_BG_HEX = "242021"
WHITE_HEX   = "F5EFE3"
CARD_BG_HEX = "FBF7EE"

# --------------------------------------------------------------------------
# Fonts (the originals)
# --------------------------------------------------------------------------
FONT_SERIF = "Georgia"
FONT_MONO  = "Courier New"
FONT_SANS  = "Arial"

# Template indices (in the v1 ORIG)
T_TITLE = 0
T_OUTLINE = 1
T_DARK = 4              # Problem Statement (dark bg)
T_CREAM = 23            # Discussion (cream + 2 cards)
T_QA = 27

# --------------------------------------------------------------------------
# Helpers
# --------------------------------------------------------------------------

def replace_in_runs(text_frame, old, new):
    for p in text_frame.paragraphs:
        runs = list(p.runs)
        if not runs: continue
        full = "".join(r.text for r in runs)
        if old not in full: continue
        runs[0].text = full.replace(old, new)
        for r in runs[1:]:
            r.text = ""


def deep_replace(shape, old, new):
    if shape.has_text_frame:
        replace_in_runs(shape.text_frame, old, new)
    if shape.shape_type == 6:
        for sub in shape.shapes:
            deep_replace(sub, old, new)


def set_paragraph_text(text_frame, new_text):
    paras = list(text_frame.paragraphs)
    if not paras:
        text_frame.text = new_text
        return
    runs = list(paras[0].runs)
    if runs:
        runs[0].text = new_text
        for r in runs[1:]: r.text = ""
    else:
        paras[0].add_run().text = new_text
    for p in paras[1:]:
        for r in list(p.runs): r.text = ""


def remove_shape(shape):
    shape.element.getparent().remove(shape.element)


def duplicate_slide(prs, src_idx):
    src = prs.slides[src_idx]
    new = prs.slides.add_slide(src.slide_layout)
    p_ns = "http://schemas.openxmlformats.org/presentationml/2006/main"
    src_csld = src.element.find(f"{{{p_ns}}}cSld")
    new_csld = new.element.find(f"{{{p_ns}}}cSld")
    if src_csld is not None and new_csld is not None:
        src_bg = src_csld.find(f"{{{p_ns}}}bg")
        if src_bg is not None:
            new_bg = copy.deepcopy(src_bg)
            new_csld.insert(0, new_bg)
    for shape in list(src.shapes):
        new_el = copy.deepcopy(shape.element)
        new.shapes._spTree.insert_element_before(new_el, "p:extLst")
    rels_src = src.part.rels
    a_ns = "http://schemas.openxmlformats.org/drawingml/2006/main"
    r_ns = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    for blip in new.shapes._spTree.findall(f".//{{{a_ns}}}blip"):
        rid = blip.get(f"{{{r_ns}}}embed")
        if rid and rid in rels_src:
            rel = rels_src[rid]
            new_rid = new.part.relate_to(rel.target_part, rel.reltype)
            blip.set(f"{{{r_ns}}}embed", new_rid)
    return new


def move_slide(prs, slide, new_idx):
    xml = prs.slides._sldIdLst
    el = next(s for s in xml if int(s.attrib["id"]) == slide.slide_id)
    xml.remove(el); xml.insert(new_idx, el)


def add_text(slide, x, y, w, h, text, *, size=12, bold=False,
             color=INK_HEX, font=FONT_SERIF, italic=False, align="left"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                     Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    if align == "center": p.alignment = PP_ALIGN.CENTER
    elif align == "right": p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = RGBColor.from_string(color)
    r.font.name = font
    return box


def add_lines(slide, x, y, w, h, lines, *, size=11, color=INK_HEX,
              font=FONT_SERIF, bullet="·", indent=0):
    box = slide.shapes.add_textbox(Inches(x), Inches(y),
                                     Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if indent: p.level = indent
        r = p.add_run()
        r.text = f"{bullet}  {line}" if bullet else line
        r.font.size = Pt(size)
        r.font.color.rgb = RGBColor.from_string(color)
        r.font.name = font
        p.space_after = Pt(2)
    return box


def add_card(slide, x, y, w, h, *, fill=CARD_BG_HEX, border=SUBINK_HEX,
             accent_bar=False):
    card = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h))
    card.line.color.rgb = RGBColor.from_string(border)
    card.line.width = Pt(0.5)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor.from_string(fill)
    if accent_bar:
        bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(x), Inches(y), Inches(0.06), Inches(h))
        bar.line.fill.background()
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor.from_string(ACCENT_HEX)
    return card


def add_picture_fit(slide, path, x, y, max_w, max_h):
    src = Image.open(path); sw, sh = src.size; src.close()
    aspect = sw / sh
    box_aspect = max_w / max_h
    if aspect > box_aspect:
        w = max_w; h = max_w / aspect
    else:
        h = max_h; w = max_h * aspect
    return slide.shapes.add_picture(str(path),
        Inches(x + (max_w - w) / 2),
        Inches(y + (max_h - h) / 2),
        width=Inches(w), height=Inches(h))


# --------------------------------------------------------------------------
# Section-opener templates: dark + cream
# --------------------------------------------------------------------------

def make_dark_slide(prs, *, section_label, title_main, title_accent="",
                     big_number=""):
    """Duplicate the dark Problem Statement template, clear its body,
    rewrite header / title / big number. Returns new slide; caller adds
    body content."""
    new = duplicate_slide(prs, T_DARK)
    keep = ["§ 03 PROBLEM STATEMENT", "How do we recover",
             "APRIL 2026", "AI-Enhanced", "03"]
    keep_shapes = []
    header_sh = title_sh = num_sh = body_sh = None
    for sh in new.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text.strip()
        if "§ 03 PROBLEM STATEMENT" in t:
            header_sh = sh; keep_shapes.append(sh); continue
        if "How do we recover" in t:
            title_sh = sh; keep_shapes.append(sh); continue
        if t == "03":
            num_sh = sh; keep_shapes.append(sh); continue
        if t.startswith("Existing high-accuracy"):
            body_sh = sh
            continue
        if "AI-Enhanced" in t or "APRIL" in t:
            keep_shapes.append(sh); continue
        if t.isdigit() and len(t) <= 2:
            keep_shapes.append(sh); continue
    # Update header
    if header_sh is not None:
        set_paragraph_text(header_sh.text_frame, section_label)
    # Update title (preserve template font / size)
    if title_sh is not None:
        tf = title_sh.text_frame
        for p in list(tf.paragraphs):
            for r in list(p.runs):
                p._p.remove(r._r)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = title_main + (" " if title_accent else "")
        r1.font.bold = True; r1.font.size = Pt(30)
        r1.font.color.rgb = RGBColor.from_string(WHITE_HEX)
        r1.font.name = FONT_SERIF
        if title_accent:
            r2 = p.add_run(); r2.text = title_accent
            r2.font.bold = True; r2.font.size = Pt(30)
            r2.font.color.rgb = RGBColor.from_string(SOFTACC_HEX)
            r2.font.name = FONT_SERIF
    # Update big number
    if num_sh is not None and big_number:
        set_paragraph_text(num_sh.text_frame, big_number)
    # Drop body & any other body shapes
    if body_sh is not None: remove_shape(body_sh)
    for sh in list(new.shapes):
        if sh in keep_shapes: continue
        if sh.top is not None and sh.top >= Inches(5.20): continue
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t.isdigit() and len(t) <= 2: continue
            if "AI-Enhanced" in t or "APRIL" in t: continue
            if t.startswith("§ "): continue
        else:
            if sh.top is not None and Inches(1.20) < sh.top < Inches(5.20):
                remove_shape(sh)
            continue
        remove_shape(sh)
    return new


def make_cream_slide(prs, *, section_label, title_main, title_accent=""):
    """Duplicate cream template (Discussion); clear body."""
    new = duplicate_slide(prs, T_CREAM)
    keep_shapes = []
    header_sh = title_sh = None
    for sh in new.shapes:
        if not sh.has_text_frame: continue
        t = sh.text_frame.text.strip()
        if t.startswith("§ 14"):
            header_sh = sh; keep_shapes.append(sh); continue
        if "Reading the numbers" in t:
            title_sh = sh; keep_shapes.append(sh); continue
        if "AI-Enhanced" in t or "APRIL" in t:
            keep_shapes.append(sh); continue
        if t.isdigit() and len(t) <= 2:
            if sh.top is not None and sh.top > Inches(5.20):
                keep_shapes.append(sh); continue
    if header_sh is not None:
        set_paragraph_text(header_sh.text_frame, section_label)
    if title_sh is not None:
        tf = title_sh.text_frame
        for p in list(tf.paragraphs):
            for r in list(p.runs):
                p._p.remove(r._r)
        p = tf.paragraphs[0]
        r1 = p.add_run(); r1.text = title_main + (" " if title_accent else "")
        r1.font.bold = True; r1.font.size = Pt(28)
        r1.font.color.rgb = RGBColor.from_string(INK_HEX)
        r1.font.name = FONT_SERIF
        if title_accent:
            r2 = p.add_run(); r2.text = title_accent
            r2.font.bold = True; r2.font.size = Pt(28)
            r2.font.color.rgb = RGBColor.from_string(ACCENT_HEX)
            r2.font.name = FONT_SERIF
    # Drop body
    for sh in list(new.shapes):
        if sh in keep_shapes: continue
        if sh.top is not None and sh.top >= Inches(5.20): continue
        if sh.has_text_frame:
            t = sh.text_frame.text.strip()
            if t.isdigit() and len(t) <= 2: continue
            if "AI-Enhanced" in t or "APRIL" in t: continue
            if t.startswith("§ "): continue
        if sh.top is not None and Inches(1.20) < sh.top < Inches(5.20):
            remove_shape(sh); continue
        if not sh.has_text_frame: continue
        remove_shape(sh)
    return new


# --------------------------------------------------------------------------
# Slide builders
# --------------------------------------------------------------------------

def edit_title(prs):
    s = prs.slides[T_TITLE]
    for sh in s.shapes:
        deep_replace(sh, "MTE 4200", "MTE 4210")
        deep_replace(sh, "Project and Thesis", "Seminar")


def build_outline(prs):
    s = prs.slides[1]
    # Strip old body shapes
    for sh in list(s.shapes):
        if not sh.has_text_frame:
            if sh.top is not None and Inches(1.10) < sh.top < Inches(5.20):
                remove_shape(sh)
            continue
        t = sh.text_frame.text.strip()
        if t.startswith("§ 00"): continue
        if t == "Outline.": continue
        if "AI-Enhanced" in t or "APRIL" in t: continue
        if sh.top is not None and sh.top > Inches(5.20): continue
        remove_shape(sh)

    rows = [
        ("01",   "Introduction",            "03"),
        ("02",   "Problem Statement",       "06"),
        ("03",   "Objectives",              "07"),
        ("04",   "Literature Review",       "08"),
        ("04.1", "Research Gap",            "09"),
        ("05",   "Proposed Solution",       "10"),
        ("06",   "Methodology",             "11"),
        ("07",   "Implementation",          "12"),
        ("07.1", "Model Architecture",      "13 to 16"),
        ("08",   "Working Principle",       "17"),
        ("09",   "Results & Analysis",      "18"),
        ("09.1", "Scene Flow",              "19"),
        ("09.2", "Indoor Real Data",        "20"),
        ("09.3", "Live Inference",          "21"),
        ("09.4", "3D Reconstruction",       "22"),
        ("09.5", "Objectives Answered",     "23"),
        ("10",   "Discussion",              "24"),
        ("11",   "Challenges",              "25"),
        ("12",   "Impact & Contributions",  "26"),
        ("13",   "Conclusion",              "27"),
        ("14",   "Future Work",             "28"),
    ]
    half = (len(rows) + 1) // 2
    left, right = rows[:half], rows[half:]
    y0 = 1.40
    row_h = 0.205
    for col_x, col_rows in [(0.55, left), (5.10, right)]:
        for i, (num, label, page) in enumerate(col_rows):
            y = y0 + i * row_h
            add_text(s, col_x, y, 0.65, row_h, num,
                      size=8.5, color=ACCENT_HEX, font=FONT_MONO, bold=True)
            add_text(s, col_x + 0.65, y, 3.10, row_h, label,
                      size=10.5, color=INK_HEX, font=FONT_SERIF)
            add_text(s, col_x + 3.85, y, 0.50, row_h, page,
                      size=8.5, color=SUBINK_HEX, font=FONT_MONO,
                      align="right")
            line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(col_x), Inches(y + row_h - 0.02),
                Inches(4.30), Inches(0.005))
            line.line.fill.background()
            line.fill.solid()
            line.fill.fore_color.rgb = RGBColor.from_string("D8CFC0")


# ---- Section 1: Introduction (3 slides) ----

def build_intro_1(prs):
    """Minimal: depth perception on edge is necessary + 4 photos."""
    s = make_dark_slide(prs,
        section_label="§ 01 INTRODUCTION",
        title_main="Depth perception",
        title_accent="on edge.",
        big_number="01")
    # Four photos in a 2x2 grid
    photos = [
        (PHOTOS / "combat_drone.jpg",  "Combat drones  ·  forest, indoor"),
        (PHOTOS / "ar_headset.jpg",    "AR headsets"),
        (PHOTOS / "factory_robot.jpg", "Factory autonomous vehicles"),
        (PHOTOS / "ar0144.jpg" if (PHOTOS / "ar0144.jpg").exists() and (PHOTOS / "ar0144.jpg").stat().st_size > 5000 else PHOTOS / "jetson_orin.jpg",
         "Stereo camera"),
    ]
    cell_w = 2.05
    cell_h = 1.45
    gap = 0.18
    x0 = (10.0 - (4 * cell_w + 3 * gap)) / 2.0
    y0 = 2.10
    for i, (p, lbl) in enumerate(photos):
        x = x0 + i * (cell_w + gap)
        # Photo card
        if p.exists() and p.stat().st_size > 5000:
            add_picture_fit(s, p, x, y0, cell_w, cell_h)
        else:
            card = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(x), Inches(y0), Inches(cell_w), Inches(cell_h))
            card.line.color.rgb = RGBColor.from_string(SOFTACC_HEX)
            card.fill.solid()
            card.fill.fore_color.rgb = RGBColor.from_string(DARK_BG_HEX)
        # Caption
        add_text(s, x, y0 + cell_h + 0.05, cell_w, 0.30, lbl,
                  size=8.5, color=WHITE_HEX, font=FONT_MONO, align="center")
    add_text(s, 0.55, 4.55, 8.95, 0.40,
              "Every moving platform that needs to understand the "
              "world around it needs depth, in real time, on board.",
              size=12, italic=True, color=SOFTACC_HEX, font=FONT_SERIF,
              align="center")
    return s


def build_intro_1_1(prs):
    """Why depth estimation in edge is hard.

    Layout: 2 expensive solutions (LiDAR, high end stereo) on the left,
    3 reason cards on the right.  Soft-divider down the middle.
    """
    s = make_cream_slide(prs,
        section_label="§ 01.1 INTRODUCTION",
        title_main="Why depth estimation in edge",
        title_accent="is hard.")

    # ----- Left half: "What people use today (and why it does not fit
    # on the edge)" -----
    left_x = 0.55
    left_w = 4.30

    add_text(s, left_x, 1.35, left_w, 0.28,
              "WHAT INDUSTRY USES TODAY",
              size=9, bold=True, font=FONT_MONO, color=ACCENT_HEX)

    # Photo + caption pairs
    photo_w = 2.00
    photo_h = 1.55
    gap = 0.30
    px = left_x
    py = 1.70
    add_picture_fit(s, PHOTOS / "lidar.jpg",
                     x=px, y=py, max_w=photo_w, max_h=photo_h)
    add_text(s, px, py + photo_h + 0.05, photo_w, 0.25,
              "3D LiDAR  ·  Velodyne", size=9.5, bold=True,
              font=FONT_MONO, color=INK_HEX, align="center")
    add_text(s, px, py + photo_h + 0.32, photo_w, 0.45,
              "$2k to $80k.\n10 to 30 W.\nMechanical or solid state, "
              "still bulky.",
              size=8.5, color=SUBINK_HEX, font=FONT_SERIF,
              align="center")

    px2 = left_x + photo_w + gap
    add_picture_fit(s, PHOTOS / "realsense.jpg",
                     x=px2, y=py, max_w=photo_w, max_h=photo_h)
    add_text(s, px2, py + photo_h + 0.05, photo_w, 0.25,
              "Intel RealSense D435", size=9.5, bold=True,
              font=FONT_MONO, color=INK_HEX, align="center")
    add_text(s, px2, py + photo_h + 0.32, photo_w, 0.45,
              "$300 to $500.\nIR projector + 2 cams.\n"
              "Indoor only, fragile in sunlight.",
              size=8.5, color=SUBINK_HEX, font=FONT_SERIF,
              align="center")

    # Soft vertical divider
    div = s.shapes.add_shape(MSO_SHAPE.RECTANGLE,
        Inches(5.05), Inches(1.40), Inches(0.005), Inches(3.40))
    div.line.fill.background()
    div.fill.solid()
    div.fill.fore_color.rgb = RGBColor.from_string("D8CFC0")

    # ----- Right half: 3 reasons depth on edge is hard -----
    right_x = 5.20
    right_w = 4.30

    add_text(s, right_x, 1.35, right_w, 0.28,
              "WHAT GOES WRONG ON THE EDGE",
              size=9, bold=True, font=FONT_MONO, color=ACCENT_HEX)

    # Three small cards stacked vertically
    cards = [
        ("FOUNDATION STEREO NETWORKS",
         "215 to 350 M params",
         "FoundationStereo, DEFOM, MonSter need a desktop GPU and "
         "5 to 10 GB VRAM. Cannot run on a Jetson class device."),
        ("EDGE COMPUTE BUDGET",
         "~6 TOPS  ·  ~4 GB  ·  5 to 25 W",
         "Drones and AGVs share that budget with the rest of the "
         "autonomy stack. Stereo gets a slice, not the whole cake."),
        ("CLASSICAL METHODS",
         "EPE > 5 px on hard scenes",
         "SGM and block matching are light enough for the edge but "
         "collapse on textureless walls, reflections, and foliage."),
    ]
    cy = 1.70
    ch = 1.00
    cgap = 0.10
    for title, headline, body in cards:
        add_card(s, right_x, cy, right_w, ch)
        add_text(s, right_x + 0.15, cy + 0.10, right_w - 0.30, 0.22,
                  title, size=8.5, bold=True, font=FONT_MONO,
                  color=ACCENT_HEX)
        add_text(s, right_x + 0.15, cy + 0.32, right_w - 0.30, 0.30,
                  headline, size=12, bold=True, color=INK_HEX,
                  font=FONT_SERIF)
        add_text(s, right_x + 0.15, cy + 0.62, right_w - 0.30, 0.45,
                  body, size=9, color=SUBINK_HEX, font=FONT_SERIF)
        cy += ch + cgap

    # Closing line
    add_text(s, 0.55, 4.95, 8.95, 0.30,
              "Lightweight, accurate, deep stereo on edge silicon is "
              "the gap StereoLite is built for.",
              size=11, italic=True, bold=True, color=ACCENT_HEX,
              font=FONT_SERIF, align="center")
    return s


def build_intro_1_2(prs):
    """10s indoor inference clip + one 3D mesh GIF + StereoLite intro."""
    s = make_cream_slide(prs,
        section_label="§ 01.2 INTRODUCTION",
        title_main="Introducing",
        title_accent="StereoLite.")
    # Top left: video clip embedded
    poster = FIGS / "_inference_preview.png"
    if INTRO_VIDEO.exists():
        s.shapes.add_movie(
            str(INTRO_VIDEO), Inches(0.55), Inches(1.40),
            Inches(5.40), Inches(1.62),
            poster_frame_image=str(poster), mime_type="video/mp4")
    add_text(s, 0.55, 3.05, 5.40, 0.25,
              "Indoor inference  ·  10 s clip  ·  1.5x slowdown",
              size=8.5, italic=True, color=SUBINK_HEX, font=FONT_MONO,
              align="center")

    # Top right: one rotating point cloud GIF
    gifs = sorted(PCTOP.glob("pair_*_small.gif"))
    if gifs:
        s.shapes.add_picture(str(gifs[0]),
            Inches(6.30), Inches(1.40),
            width=Inches(3.20), height=Inches(1.62))
    add_text(s, 6.30, 3.05, 3.20, 0.25,
              "3D reconstruction from disparity",
              size=8.5, italic=True, color=SUBINK_HEX, font=FONT_MONO,
              align="center")

    # Bottom: StereoLite tagline
    add_text(s, 0.55, 3.55, 8.95, 0.40,
              "StereoLite",
              size=24, bold=True, color=ACCENT_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 4.05, 8.95, 0.40,
              "0.87 M parameters  ·  54 ms on RTX 3050  ·  ~40 ms "
              "expected on Jetson Orin Nano (INT8)  ·  0.515 px EPE on "
              "indoor real data.",
              size=11, color=INK_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 4.55, 8.95, 0.35,
              "An edge sized stereo network that learns dense depth "
              "from two cheap cameras, end to end.",
              size=11, italic=True, color=SUBINK_HEX, font=FONT_SERIF)
    return s


# ---- Section 2: Problem ----

def build_problem(prs):
    s = make_dark_slide(prs,
        section_label="§ 02 PROBLEM STATEMENT",
        title_main="Real time depth on edge devices,",
        title_accent="at workable accuracy.",
        big_number="02")
    # Three blocks: COMPUTE, MEMORY, POWER
    cols_x = [0.55, 3.45, 6.35]
    titles = ["LIMITED COMPUTE", "TIGHT MEMORY", "POWER BUDGET"]
    icons  = ["~6 TOPS", "~4 GB", "5 to 25 W"]
    notes = [
        "Embedded SoCs deliver a fraction of a desktop GPU.",
        "RAM shared with the rest of the autonomy stack.",
        "Battery powered platforms cannot host hot GPUs.",
    ]
    for x, t, ic, note in zip(cols_x, titles, icons, notes):
        add_text(s, x, 2.10, 2.85, 0.30, t,
                  size=10, bold=True, font=FONT_MONO, color=SOFTACC_HEX)
        add_text(s, x, 2.45, 2.85, 0.85, ic,
                  size=28, bold=True, color=WHITE_HEX, font=FONT_SERIF)
        add_text(s, x, 3.40, 2.85, 0.85, note,
                  size=11, color=WHITE_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 4.50, 8.95, 0.40,
              "How do we deliver dense, accurate depth at edge scale "
              "compute, on cameras with no calibration luxury?",
              size=12, italic=True, color=SOFTACC_HEX, font=FONT_SERIF)
    return s


# ---- Section 3: Objectives ----

def build_objectives(prs):
    s = make_cream_slide(prs,
        section_label="§ 03 OBJECTIVES",
        title_main="Objectives",
        title_accent="and success criteria.")

    # Tree: each objective has a header line then indented sub bullets
    add_text(s, 0.55, 1.45, 0.55, 0.30, "01",
              size=12, bold=True, font=FONT_MONO, color=ACCENT_HEX)
    add_text(s, 1.10, 1.45, 8.40, 0.30,
              "Design a computationally efficient stereo matching "
              "pipeline with AI based disparity refinement on resource "
              "limited platforms.",
              size=12, color=INK_HEX, font=FONT_SERIF, bold=True)
    add_lines(s, 1.10, 1.85, 8.40, 1.00, [
        "Trainable parameters under 1 M.",
        "Inference under 40 ms per 512x832 stereo pair on Jetson Orin "
        "Nano (INT8).",
        "Checkpoint under 20 MB on disk.",
    ], size=11, color=SUBINK_HEX, font=FONT_SERIF)

    add_text(s, 0.55, 3.10, 0.55, 0.30, "02",
              size=12, bold=True, font=FONT_MONO, color=ACCENT_HEX)
    add_text(s, 1.10, 3.10, 8.40, 0.30,
              "Design an architecture that withstands camera "
              "imperfections.",
              size=12, color=INK_HEX, font=FONT_SERIF, bold=True)
    add_lines(s, 1.10, 3.50, 8.40, 1.00, [
        "Works on cheap stereo cameras with imperfect alignment and "
        "missing rectification.",
        "End to end trainable with multi scale supervision so the "
        "network learns to absorb small calibration errors.",
    ], size=11, color=SUBINK_HEX, font=FONT_SERIF)

    add_text(s, 0.55, 4.85, 8.95, 0.30,
              "Both criteria are answered explicitly in the results "
              "section (slide 09.5).",
              size=10.5, italic=True, color=ACCENT_HEX, font=FONT_SERIF)
    return s


# ---- Section 4: Literature Review (DARK) + 4.1 Research Gap (cream) ----

def build_literature(prs):
    s = make_dark_slide(prs,
        section_label="§ 04 LITERATURE REVIEW",
        title_main="110 papers analyzed.",
        title_accent="Five informed StereoLite.",
        big_number="04")
    rows = [
        ("GC-Net  (Kendall et al., 2017)",
         "First end to end 3D cost volume regression",        "3.5 M"),
        ("PSMNet  (Chang and Chen, 2018)",
         "Pyramid pooling + stacked hourglass 3D CNN",        "5.2 M"),
        ("HITNet  (Tankovich et al., 2021)",
         "Tile hypothesis planes (d, slope x, slope y, f, c)","0.6 M"),
        ("RAFT-Stereo  (Lipson et al., 2021)",
         "Recurrent residual refinement",                     "11 M"),
        ("MobileStereoNet  (Shamsafar et al., 2022)",
         "MobileNet backbone for stereo",                     "2.4 M"),
    ]
    add_text(s, 0.55, 2.05, 4.50, 0.25, "METHOD",
              size=8.5, bold=True, font=FONT_MONO, color=SOFTACC_HEX)
    add_text(s, 5.10, 2.05, 3.85, 0.25, "CORE IDEA",
              size=8.5, bold=True, font=FONT_MONO, color=SOFTACC_HEX)
    add_text(s, 8.95, 2.05, 0.55, 0.25, "PARAMS",
              size=8.5, bold=True, font=FONT_MONO, color=SOFTACC_HEX,
              align="right")
    y = 2.40
    for method, idea, params in rows:
        add_text(s, 0.55, y, 4.50, 0.40, method,
                  size=10.5, bold=True, color=WHITE_HEX, font=FONT_SERIF)
        add_text(s, 5.10, y, 3.85, 0.40, idea,
                  size=10.5, color=WHITE_HEX, font=FONT_SERIF)
        add_text(s, 8.95, y, 0.55, 0.40, params,
                  size=10.5, color=SOFTACC_HEX, font=FONT_MONO, align="right")
        y += 0.45
    add_text(s, 0.55, 4.85, 8.95, 0.30,
              "Three tier ranking by relevance to edge deployment "
              "across the 2002 to 2026 stereo matching literature.",
              size=10.5, italic=True, color=SOFTACC_HEX, font=FONT_SERIF)
    return s


def build_research_gap(prs):
    s = make_cream_slide(prs,
        section_label="§ 04.1 LITERATURE REVIEW  ·  RESEARCH GAP",
        title_main="A gap in the",
        title_accent="accuracy x size plane.")
    # Pareto chart spans full width to fit numbered legend
    add_picture_fit(s, FIGS / "research_gap_pareto.png",
                     x=0.30, y=1.30, max_w=9.40, max_h=3.60)
    add_text(s, 0.30, 4.95, 9.40, 0.25,
              "Numbered dots map to the side legend.  Two stars: "
              "current StereoLite checkpoint and projected after full "
              "Scene Flow pre training.",
              size=9.5, color=SUBINK_HEX, font=FONT_SERIF, align="center")
    return s


# ---- Section 5: Proposed Solution ----

def build_proposed(prs):
    s = make_cream_slide(prs,
        section_label="§ 05 PROPOSED SOLUTION",
        title_main="StereoLite",
        title_accent="at a glance.")
    cols_x = [0.55, 2.65, 4.85, 7.15]
    titles = ["TRAINABLE PARAMS", "LATENCY  ·  RTX 3050",
               "LATENCY  ·  ORIN NANO", "VAL EPE"]
    big = ["0.87 M", "54 ms", "~40 ms", "1.54 / 0.515"]
    notes = [
        "8.7 MB fp32 checkpoint.  INT8 conversion planned for "
        "deployment.",
        "Per 512x832 stereo pair, batch size 1.",
        "Projected with INT8 + TensorRT, full validation pending.",
        "Scene Flow Driving 200 val (1.54)  /  indoor 50 val "
        "(0.515) px.",
    ]
    for x, t, b, n in zip(cols_x, titles, big, notes):
        add_text(s, x, 1.40, 2.30, 0.25, t,
                  size=8, bold=True, font=FONT_MONO, color=ACCENT_HEX)
        add_text(s, x, 1.65, 2.20, 0.65, b,
                  size=22, bold=True, color=INK_HEX, font=FONT_SERIF)
        add_text(s, x, 2.30, 2.20, 0.95, n,
                  size=8.5, color=SUBINK_HEX, font=FONT_SERIF)

    # Idea bullets without IDEA labels
    add_text(s, 0.55, 3.40, 8.95, 0.30,
              "FOUR DESIGN CHOICES",
              size=8.5, bold=True, font=FONT_MONO, color=ACCENT_HEX)
    add_lines(s, 0.55, 3.70, 8.95, 1.30, [
        "Pretrained encoder. Truncated MobileNetV2 reuses ImageNet "
        "features at four scales.",
        "One coarse cost volume. Built only at 1/16, cheap but rich "
        "enough to initialise.",
        "Iterative refinement. Two plus three plus three residual "
        "updates over tile hypotheses.",
        "Learned upsample. Convex 9 neighbour weights for boundary "
        "sharpness.",
    ], size=10.5, color=INK_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 5.00, 8.95, 0.30,
              "Original work. This combination of ideas under 1 M "
              "params has not been published before.",
              size=11, italic=True, bold=True, color=ACCENT_HEX,
              font=FONT_SERIF)
    return s


# ---- Section 6: Methodology (block diagram) ----

def build_methodology(prs):
    s = make_cream_slide(prs,
        section_label="§ 06 METHODOLOGY",
        title_main="How we",
        title_accent="executed the thesis.")
    # 4 blocks connected by arrows, horizontal flow
    blocks = [
        ("01", "LITERATURE",   "110 papers\n3 tier ranking"),
        ("02", "ARCHITECTURE", "3 candidate designs\nfinal: 0.87 M"),
        ("03", "TRAINING",     "Synthetic + real\n+ pseudo GT"),
        ("04", "VALIDATION",   "EPE + 3D point clouds"),
    ]
    n = len(blocks)
    block_w = 1.85
    gap = 0.30
    total = n * block_w + (n - 1) * gap
    x0 = (10.0 - total) / 2.0
    y_block = 1.85
    h_block = 1.60
    for i, (num, title, desc) in enumerate(blocks):
        x = x0 + i * (block_w + gap)
        # Card
        add_card(s, x, y_block, block_w, h_block)
        # Number
        add_text(s, x + 0.10, y_block + 0.10, 0.40, 0.28, num,
                  size=11, bold=True, font=FONT_MONO, color=ACCENT_HEX)
        # Title
        add_text(s, x + 0.10, y_block + 0.40, block_w - 0.20, 0.32,
                  title, size=11, bold=True, color=INK_HEX,
                  font=FONT_SERIF)
        # Description
        add_text(s, x + 0.10, y_block + 0.78, block_w - 0.20,
                  h_block - 0.85, desc,
                  size=10, color=SUBINK_HEX, font=FONT_SERIF)
        # Arrow
        if i < n - 1:
            arrow = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,
                Inches(x + block_w + 0.02),
                Inches(y_block + h_block / 2 - 0.10),
                Inches(gap - 0.04), Inches(0.20))
            arrow.line.fill.background()
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = RGBColor.from_string(ACCENT_HEX)

    add_text(s, 0.55, 3.85, 8.95, 0.30,
              "End to end reproducible pipeline.",
              size=11, color=ACCENT_HEX, italic=True, bold=True,
              font=FONT_SERIF, align="center")
    add_text(s, 0.55, 4.30, 8.95, 0.40,
              "Code, training logs, checkpoints, point clouds, and "
              "figures all tracked in the project repository.",
              size=10.5, color=SUBINK_HEX, font=FONT_SERIF, align="center")
    return s


# ---- Section 7: Implementation (DARK) + 4 architecture sub-slides ----

def build_implementation(prs):
    s = make_dark_slide(prs,
        section_label="§ 07 IMPLEMENTATION",
        title_main="Hardware",
        title_accent="and software.",
        big_number="07")

    # Hardware row: setup + AR0144 + Jetson
    hw_y = 1.95
    hw_h = 1.60
    add_text(s, 0.55, hw_y - 0.30, 8.95, 0.25, "HARDWARE",
              size=9, bold=True, font=FONT_MONO, color=SOFTACC_HEX)
    hw = [
        (PHOTOS / "ar0144.jpg" if (PHOTOS / "ar0144.jpg").exists() and (PHOTOS / "ar0144.jpg").stat().st_size > 5000 else None,
         "Waveshare AR0144"),
        (PHOTOS / "jetson_orin.jpg", "NVIDIA Jetson Orin Nano"),
    ]
    cell_w = 4.30
    gap = 0.35
    x0 = (10.0 - (2 * cell_w + gap)) / 2.0
    for i, (p, lbl) in enumerate(hw):
        x = x0 + i * (cell_w + gap)
        if p is not None and p.exists() and p.stat().st_size > 5000:
            add_picture_fit(s, p, x, hw_y, cell_w, hw_h)
        else:
            add_card(s, x, hw_y, cell_w, hw_h, fill=DARK_BG_HEX,
                      border=SOFTACC_HEX)
        add_text(s, x, hw_y + hw_h + 0.05, cell_w, 0.25, lbl,
                  size=10, color=WHITE_HEX, font=FONT_MONO, align="center")

    # Software logos row
    sw_y = 4.20
    sw_h = 0.75
    add_text(s, 0.55, sw_y - 0.30, 8.95, 0.25, "SOFTWARE",
              size=9, bold=True, font=FONT_MONO, color=SOFTACC_HEX)
    logos = [
        (PHOTOS / "pytorch_logo.png", "PyTorch"),
        (PHOTOS / "cuda_logo.png",    "CUDA"),
        (PHOTOS / "kaggle_logo.png",  "Kaggle"),
        (PHOTOS / "open3d_logo.png",  "Open3D"),
    ]
    cell_w = 1.95
    gap = 0.20
    x0 = (10.0 - (4 * cell_w + 3 * gap)) / 2.0
    for i, (p, lbl) in enumerate(logos):
        x = x0 + i * (cell_w + gap)
        if p.exists() and p.stat().st_size > 1000:
            add_picture_fit(s, p, x, sw_y, cell_w, sw_h)
        add_text(s, x, sw_y + sw_h + 0.04, cell_w, 0.22, lbl,
                  size=9, color=WHITE_HEX, font=FONT_MONO, align="center")
    return s


def _arch_block(slide, x, y, w, h, *, title, lines):
    add_card(slide, x, y, w, h)
    add_text(slide, x + 0.18, y + 0.12, w - 0.36, 0.30, title,
              size=9, bold=True, font=FONT_MONO, color=ACCENT_HEX)
    add_lines(slide, x + 0.18, y + 0.45, w - 0.36, h - 0.55,
                lines, size=9.5, color=INK_HEX, font=FONT_SERIF)


def _arch_visual_block(slide, x, y, w, h, *, title, image_path,
                        caption=""):
    """Two-stage block with a labelled image on top and short caption
    below, used on slides 7.2 / 7.3 / 7.4."""
    add_card(slide, x, y, w, h)
    add_text(slide, x + 0.15, y + 0.10, w - 0.30, 0.25, title,
              size=9, bold=True, font=FONT_MONO, color=ACCENT_HEX)
    img_y = y + 0.42
    img_h = h - (0.50 + (0.30 if caption else 0.0))
    add_picture_fit(slide, image_path,
                     x=x + 0.10, y=img_y,
                     max_w=w - 0.20, max_h=img_h)
    if caption:
        add_text(slide, x + 0.18, y + h - 0.32, w - 0.36, 0.25,
                  caption, size=8.5, color=SUBINK_HEX,
                  font=FONT_SERIF, italic=True, align="center")


def build_arch_overview(prs):
    s = make_cream_slide(prs,
        section_label="§ 07.1 IMPLEMENTATION  ·  MODEL ARCHITECTURE",
        title_main="StereoLite",
        title_accent="end to end.")
    add_picture_fit(s, ARCH_FULL, x=0.30, y=1.30,
                     max_w=9.40, max_h=3.60)
    add_text(s, 0.55, 4.95, 8.95, 0.30,
              "0.87 M params  ·  54 ms on RTX 3050  ·  ~40 ms expected "
              "on Jetson Orin Nano (INT8)  ·  1.54 px val EPE on Scene "
              "Flow Driving.",
              size=9.5, color=SUBINK_HEX, italic=True, font=FONT_SERIF,
              align="center")
    return s


def build_arch_stack_a(prs):
    s = make_cream_slide(prs,
        section_label="§ 07.2 IMPLEMENTATION  ·  ENCODER + INIT",
        title_main="Pretrained encoder,",
        title_accent="one cost volume.")
    _arch_visual_block(s, 0.55, 1.40, 4.40, 3.55,
        title="STAGE 1  ·  ENCODER",
        image_path=ROOT / "model/designs/d1_tile/mobilenet_truncation.png",
        caption="MobileNetV2-100, ImageNet pretrained, truncated.")
    _arch_visual_block(s, 5.05, 1.40, 4.40, 3.55,
        title="STAGE 2  ·  TILE HYPOTHESIS INIT",
        image_path=FIGS / "stage2_init.png",
        caption="One coarse cost volume to initialise tile state.")
    return s


def build_arch_stack_b(prs):
    s = make_cream_slide(prs,
        section_label="§ 07.3 IMPLEMENTATION  ·  REFINE + UPSAMPLE",
        title_main="Iterative residuals,",
        title_accent="convex upsample.")
    _arch_visual_block(s, 0.55, 1.40, 4.40, 3.55,
        title="STAGE 3  ·  ITERATIVE REFINEMENT",
        image_path=FIGS / "stage3_refine.png",
        caption="2+3+3 = 8 residual updates, plane equation between "
                 "scales.")
    _arch_visual_block(s, 5.05, 1.40, 4.40, 3.55,
        title="STAGE 4  ·  LEARNED CONVEX UPSAMPLE",
        image_path=FIGS / "stage4_upsample.png",
        caption="RAFT style: 9 learned weights per fine pixel, applied "
                 "twice.")
    return s


def build_arch_stack_c(prs):
    s = make_cream_slide(prs,
        section_label="§ 07.4 IMPLEMENTATION  ·  LOSS + BUDGET",
        title_main="Multi scale loss,",
        title_accent="parameter budget.")
    _arch_visual_block(s, 0.55, 1.40, 4.40, 3.55,
        title="SUPERVISION  ·  THREE TERMS",
        image_path=FIGS / "supervision_loss.png",
        caption="L1 pixel + Sobel gradient + bad-1 hinge, all "
                 "weighted by scale.")
    _arch_visual_block(s, 5.05, 1.40, 4.40, 3.55,
        title="PARAMETER BUDGET  ·  0.87 M TOTAL",
        image_path=FIGS / "param_budget.png",
        caption="Encoder dominates at 62 percent. Checkpoint 8.7 MB "
                 "(fp32).")
    return s


# ---- Section 8: Working Principle ----

def build_working_principle(prs):
    s = make_cream_slide(prs,
        section_label="§ 08 WORKING PRINCIPLE",
        title_main="Dataset collection",
        title_accent="and training.")
    if TRAIN_GIF.exists():
        s.shapes.add_picture(str(TRAIN_GIF),
            Inches(0.55), Inches(1.45), width=Inches(4.80))
    add_text(s, 0.55, 4.55, 4.80, 0.30,
              "Training progression on tracked indoor val pairs.",
              size=9, italic=True, color=SUBINK_HEX, font=FONT_SERIF,
              align="center")

    rx = 5.55
    blocks = [
        ("01  ·  SYNTHETIC PRETRAIN",
         "4,200 pairs from Scene Flow Driving (about 12% of the full "
         "Scene Flow corpus). 30 epochs, Kaggle dual T4."),
        ("02  ·  REAL DATA COLLECTION",
         "1,587 stereo pairs captured with the AR0144 USB camera "
         "indoors at RUET. Quality filter accepts 997."),
        ("03  ·  PSEUDO GROUND TRUTH",
         "FoundationStereo (CVPR 2025, 215 M) generates per pair "
         "disparity targets for the 997 indoor pairs."),
        ("04  ·  REAL DATA FINE TUNE",
         "9,000 steps on RTX 3050.  Mean val EPE drops from 1.54 px "
         "(synthetic only) to 0.515 px (50 indoor val held out)."),
    ]
    by = 1.45
    for t, b in blocks:
        add_text(s, rx, by, 4.00, 0.25, t,
                  size=9, bold=True, font=FONT_MONO, color=ACCENT_HEX)
        add_text(s, rx, by + 0.30, 4.00, 0.65, b,
                  size=10, color=INK_HEX, font=FONT_SERIF)
        by += 0.95
    return s


# ---- Section 9: Results (divider + 4 sub + objectives answered) ----

def build_results_divider(prs):
    s = make_dark_slide(prs,
        section_label="§ 09 RESULTS & ANALYSIS",
        title_main="Two domains,",
        title_accent="two checkpoints.",
        big_number="09")
    add_lines(s, 0.55, 1.95, 8.50, 2.40, [
        "09.1  Scene Flow Driving (synthetic) on the 12% subset.",
        "09.2  Indoor real data fine tune (custom dataset).",
        "09.3  Live inference panel video on indoor stereo.",
        "09.4  3D point cloud reconstruction from disparity.",
        "09.5  Explicit point by point answer to both objectives.",
    ], size=12, color=WHITE_HEX, font=FONT_SERIF)
    return s


def build_results_sf(prs):
    s = make_cream_slide(prs,
        section_label="§ 09.1 RESULTS  ·  SCENE FLOW",
        title_main="Scene Flow Driving",
        title_accent="subset.")
    add_picture_fit(s, SF_TRAINING_CURVES,
                     x=0.40, y=1.30, max_w=4.40, max_h=2.40)
    add_picture_fit(s, SF_PROGRESSION,
                     x=5.00, y=1.30, max_w=4.70, max_h=2.40)
    add_text(s, 0.40, 3.78, 4.40, 0.22,
              "Training curves  ·  Kaggle dual T4, 30 epochs",
              size=8.5, italic=True, color=SUBINK_HEX, font=FONT_SERIF,
              align="center")
    add_text(s, 5.00, 3.78, 4.70, 0.22,
              "Qualitative progression  ·  step 500 to 7500",
              size=8.5, italic=True, color=SUBINK_HEX, font=FONT_SERIF,
              align="center")
    add_lines(s, 0.40, 4.10, 9.30, 0.85, [
        "Trained on 12% of the full Scene Flow corpus (4,200 of "
        "~35,000 pairs).  Val EPE 1.54 px on 200 held out pairs.",
        "Loss drops about two orders of magnitude over 30 epochs.",
        "Projected EPE after full Scene Flow pre training: ~0.71 px "
        "(outlined star on the research gap chart).",
    ], size=9.5, color=INK_HEX, font=FONT_SERIF)
    add_text(s, 0.40, 5.00, 9.30, 0.22,
              "Disparity scale on the colormap: 0 to 120 px at full "
              "resolution.",
              size=8.5, italic=True, color=SUBINK_HEX, font=FONT_SERIF)
    return s


def build_results_real(prs):
    s = make_cream_slide(prs,
        section_label="§ 09.2 RESULTS  ·  INDOOR REAL DATA",
        title_main="Indoor real data",
        title_accent="fine tune.")
    add_picture_fit(s, FIGS / "realdata_training.png",
                     x=1.50, y=1.30, max_w=7.00, max_h=2.80)
    add_text(s, 0.55, 4.20, 8.95, 0.30,
              "Teacher: FoundationStereo (215 M params).  Student: "
              "StereoLite (0.87 M).  997 clean pairs, 50 val held out, "
              "9,000 steps in 1 h 35 m on RTX 3050.",
              size=9.5, color=SUBINK_HEX, font=FONT_MONO)
    add_text(s, 0.55, 4.55, 8.95, 0.30,
              "Camera imperfection tolerant.  Inference runs on raw, "
              "non rectified AR0144 frames and still reaches 0.515 px "
              "mean EPE.",
              size=10.5, italic=True, color=ACCENT_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 4.95, 8.95, 0.25,
              "Disparity scale here: about 5 to 70 px range across the "
              "indoor scenes.",
              size=8.5, italic=True, color=SUBINK_HEX, font=FONT_SERIF)
    return s


def build_results_video(prs):
    s = make_cream_slide(prs,
        section_label="§ 09.3 RESULTS  ·  LIVE INFERENCE",
        title_main="Live inference",
        title_accent="on indoor stereo video.")
    poster = FIGS / "_inference_preview.png"
    if DEMO_VIDEO_60S.exists():
        s.shapes.add_movie(
            str(DEMO_VIDEO_60S), Inches(0.55), Inches(1.40),
            Inches(8.95), Inches(2.65),
            poster_frame_image=str(poster), mime_type="video/mp4")
    add_text(s, 0.55, 4.20, 8.95, 0.30,
              "60 second clip from a 3 minute walk through, played at "
              "1.5x slowdown so the eye can track the depth gradient.",
              size=10, color=SUBINK_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 4.55, 8.95, 0.30,
              "Disparity scale: about 5 to 70 px on the TURBO "
              "colormap.  Close objects red, far objects blue.",
              size=9, italic=True, color=SUBINK_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 4.85, 8.95, 0.30,
              "End to end inference on RTX 3050 at 54 ms per frame.  "
              "Expected ~40 ms on Jetson Orin Nano with INT8.",
              size=10, italic=True, color=ACCENT_HEX, font=FONT_SERIF)
    return s


def build_results_3d(prs):
    s = make_cream_slide(prs,
        section_label="§ 09.4 RESULTS  ·  3D RECONSTRUCTION",
        title_main="From disparity",
        title_accent="to 3D point clouds.")
    gifs = sorted(PCTOP.glob("pair_*_small.gif"))[:3]
    panel_w = 2.95
    panel_h = 2.30
    gap = 0.10
    total_w = 3 * panel_w + 2 * gap
    x0 = (10.0 - total_w) / 2.0
    y0 = 1.40
    for i, gif in enumerate(gifs):
        s.shapes.add_picture(str(gif),
            Inches(x0 + i * (panel_w + gap)), Inches(y0),
            width=Inches(panel_w), height=Inches(panel_h))
        epe = gif.stem.split("epe")[-1].replace("_small", "")
        add_text(s, x0 + i * (panel_w + gap),
                  y0 + panel_h + 0.05, panel_w, 0.25,
                  f"val pair  ·  EPE {epe} px",
                  size=10, font=FONT_MONO, color=ACCENT_HEX, align="center")
    add_text(s, 0.55, 4.30, 8.95, 0.30,
              "Pinhole projection, AR0144 intrinsics: fx ~= 1005 px, "
              "baseline 52 mm.  Open3D post processing.",
              size=9.5, color=SUBINK_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 4.65, 8.95, 0.30,
              "Disparity scale per cloud: 10 to 50 px range in each "
              "scene.",
              size=9, italic=True, color=SUBINK_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 4.95, 8.95, 0.30,
              "These are evidence that the predicted disparity is "
              "metrically consistent, not just visually plausible.",
              size=10.5, italic=True, color=ACCENT_HEX, font=FONT_SERIF)
    return s


def build_objectives_answered(prs):
    s = make_cream_slide(prs,
        section_label="§ 09.5 RESULTS  ·  OBJECTIVES ANSWERED",
        title_main="Both objectives,",
        title_accent="explicitly.")
    # Two-column tree
    add_text(s, 0.55, 1.40, 0.55, 0.30, "01",
              size=12, bold=True, font=FONT_MONO, color=ACCENT_HEX)
    add_text(s, 1.10, 1.40, 8.40, 0.30,
              "Computationally efficient stereo pipeline on resource "
              "limited platforms.",
              size=11.5, bold=True, color=INK_HEX, font=FONT_SERIF)
    obj1 = [
        ("Trainable parameters under 1 M",     "0.87 M",  "MET"),
        ("Inference under 40 ms (Orin Nano)",   "~40 ms",  "MET (projected, INT8)"),
        ("Checkpoint under 20 MB",              "8.7 MB",  "MET"),
    ]
    y = 1.78
    for crit, value, verdict in obj1:
        add_text(s, 1.30, y, 4.40, 0.25, crit,
                  size=10, color=INK_HEX, font=FONT_SERIF)
        add_text(s, 5.85, y, 1.50, 0.25, value,
                  size=10, color=ACCENT_HEX, font=FONT_MONO, bold=True)
        add_text(s, 7.55, y, 2.00, 0.25, verdict,
                  size=10, color=ACCENT_HEX, bold=True, font=FONT_MONO)
        y += 0.32

    add_text(s, 0.55, 3.20, 0.55, 0.30, "02",
              size=12, bold=True, font=FONT_MONO, color=ACCENT_HEX)
    add_text(s, 1.10, 3.20, 8.40, 0.30,
              "Withstands camera imperfections.",
              size=11.5, bold=True, color=INK_HEX, font=FONT_SERIF)
    obj2 = [
        ("Works on raw (non rectified) AR0144 frames", "0.515 px EPE indoor", "MET"),
        ("End to end trainable, multi scale loss",     "L1 + grad + hinge",   "MET"),
    ]
    y = 3.55
    for crit, value, verdict in obj2:
        add_text(s, 1.30, y, 4.40, 0.25, crit,
                  size=10, color=INK_HEX, font=FONT_SERIF)
        add_text(s, 5.85, y, 1.50, 0.25, value,
                  size=10, color=ACCENT_HEX, font=FONT_MONO, bold=True)
        add_text(s, 7.55, y, 2.00, 0.25, verdict,
                  size=10, color=ACCENT_HEX, bold=True, font=FONT_MONO)
        y += 0.32

    add_text(s, 0.55, 4.55, 8.95, 0.30,
              "Both stated objectives are met at the current "
              "checkpoint.  Full Scene Flow pre training is queued to "
              "tighten the synthetic baseline (projected EPE ~0.71 px).",
              size=11, italic=True, color=ACCENT_HEX, font=FONT_SERIF)
    return s


# ---- Sections 10-14: Discussion / Challenges / Impact / Conclusion / Future ----

def build_discussion(prs):
    s = make_cream_slide(prs,
        section_label="§ 10 DISCUSSION",
        title_main="Reading the numbers,",
        title_accent="and the potential.")
    # Two-card observations
    cols_x = [0.55, 5.05]
    titles = ["OBSERVATION  ·  CV HEAD PLATEAU",
               "OBSERVATION  ·  REAL DATA TRANSFER"]
    body = [
        "1/8 scale CV L1 plateaus near 0.13 px (about 1 px at full "
        "resolution).  Effective lower bound until we deepen the "
        "cost volume.",
        "Indoor EPE drops 3x with only 997 fine tune pairs.  The "
        "architecture transfers cleanly across domains.",
    ]
    for x, t, b in zip(cols_x, titles, body):
        add_card(s, x, 1.40, 4.40, 1.20)
        add_text(s, x + 0.18, 1.55, 4.10, 0.25, t,
                  size=8.5, bold=True, font=FONT_MONO, color=ACCENT_HEX)
        add_text(s, x + 0.18, 1.85, 4.10, 0.75, b,
                  size=10, color=INK_HEX, font=FONT_SERIF)

    # Bottom: EPE trajectory diagram supports the "potential" message
    add_picture_fit(s, FIGS / "epe_trajectory.png",
                     x=0.55, y=2.75, max_w=8.95, max_h=1.95)

    add_text(s, 0.55, 4.85, 8.95, 0.30,
              "Potential.  Already competitive at 0.87 M params.  "
              "Full Scene Flow pretraining plus INT8 deployment makes "
              "it a viable foundation for production edge stereo.",
              size=10.5, italic=True, color=ACCENT_HEX, font=FONT_SERIF)
    return s


def build_challenges(prs):
    s = make_cream_slide(prs,
        section_label="§ 11 CHALLENGES & LIMITATIONS",
        title_main="What blocked us,",
        title_accent="what is still open.")
    cards = [
        ("HARDWARE & COMPUTE",
         "Full Scene Flow pre training requires an A100 class GPU for "
         "several days. Local RTX 3050 and rented Kaggle T4 cap us at "
         "the 12% subset."),
        ("PSEUDO GROUND TRUTH",
         "Indoor real data is supervised by FoundationStereo (215 M) "
         "outputs, not LiDAR. Absolute accuracy still needs a "
         "calibrated depth reference."),
        ("STRESS TESTING",
         "Inference on raw vs rectified shows marginal difference, "
         "but a wider stress suite (multiple cameras, outdoor, motion "
         "blur) needs more hardware."),
    ]
    cols_x = [0.55, 3.70, 6.85]
    for x, (t, b) in zip(cols_x, cards):
        add_card(s, x, 1.40, 2.95, 2.85)
        add_text(s, x + 0.15, 1.55, 2.65, 0.25, t,
                  size=8.5, bold=True, font=FONT_MONO, color=ACCENT_HEX)
        add_text(s, x + 0.15, 1.90, 2.65, 2.30, b,
                  size=10, color=INK_HEX, font=FONT_SERIF)

    add_text(s, 0.55, 4.45, 8.95, 0.30,
              "Current model already meets the stated objectives.  We "
              "do not want to stop here; we want to compete at global "
              "state of the art with the right resources.",
              size=10.5, italic=True, color=ACCENT_HEX, font=FONT_SERIF)
    return s


def build_impact(prs):
    s = make_cream_slide(prs,
        section_label="§ 12 IMPACT & CONTRIBUTIONS",
        title_main="Original work,",
        title_accent="real deployment paths.")
    add_lines(s, 0.55, 1.40, 8.95, 1.20, [
        "First published combination of HITNet tile hypotheses + RAFT "
        "iterative refinement at under 1 M parameters.  Original "
        "work, not a re implementation of any prior model.",
        "Reproducible reference pipeline:  multi scale supervision, "
        "training curves, public Scene Flow numbers, and a 997 pair "
        "indoor pseudo ground truth dataset.",
    ], size=10.5, color=INK_HEX, font=FONT_SERIF)

    add_text(s, 0.55, 2.75, 8.95, 0.25, "WHERE IT HELPS",
              size=9, bold=True, font=FONT_MONO, color=ACCENT_HEX)

    # 4 photo + caption tiles
    tiles = [
        (PHOTOS / "combat_drone.jpg",
         "COMBAT DRONES",
         "Indoor and dense forest obstacle avoidance."),
        (PHOTOS / "factory_robot.jpg",
         "FACTORY AVs",
         "AGV navigation without a LiDAR mast."),
        (PHOTOS / "ar_headset.jpg",
         "AR HEADSETS",
         "On device depth for hand and scene tracking."),
        (PHOTOS / "jetson_orin.jpg",
         "EDGE EMBEDDED",
         "Visual SLAM and ADAS on Jetson class SoCs."),
    ]
    cw = 2.10
    gap = 0.18
    x0 = (10.0 - (4 * cw + 3 * gap)) / 2.0
    photo_y = 3.05
    photo_h = 1.05
    for i, (p, t, b) in enumerate(tiles):
        x = x0 + i * (cw + gap)
        if p.exists() and p.stat().st_size > 5000:
            add_picture_fit(s, p, x=x, y=photo_y,
                             max_w=cw, max_h=photo_h)
        # Caption card under photo
        add_card(s, x, photo_y + photo_h + 0.05, cw, 0.80)
        add_text(s, x + 0.10, photo_y + photo_h + 0.12,
                  cw - 0.20, 0.22, t,
                  size=7.8, bold=True, font=FONT_MONO, color=ACCENT_HEX)
        add_text(s, x + 0.10, photo_y + photo_h + 0.34,
                  cw - 0.20, 0.50, b,
                  size=8.5, color=INK_HEX, font=FONT_SERIF)

    add_text(s, 0.55, 5.05, 8.95, 0.25,
              "Plus: an honest, open baseline for the under 1 M "
              "parameter stereo regime.",
              size=10, italic=True, color=ACCENT_HEX, font=FONT_SERIF,
              align="center")
    return s


def build_conclusion(prs):
    s = make_cream_slide(prs,
        section_label="§ 13 CONCLUSION",
        title_main="One small model,",
        title_accent="three concrete results.")
    add_lines(s, 0.55, 1.40, 8.95, 2.00, [
        "0.87 M parameter stereo network composed of a pretrained "
        "MobileNetV2 encoder, one coarse cost volume, iterative tile "
        "refinement, and learned convex upsampling.",
        "1.54 px EPE on Scene Flow Driving 200 val using only 12% of "
        "the full Scene Flow corpus.  0.515 px EPE on 50 held out "
        "indoor real pairs after fine tuning.",
        "End to end 3D point cloud reconstruction works on raw, non "
        "rectified AR0144 frames, demonstrating camera imperfection "
        "tolerance at edge scale compute.",
    ], size=11, color=INK_HEX, font=FONT_SERIF)
    add_text(s, 0.55, 3.65, 8.95, 0.30,
              "Both stated objectives are met:  efficient under 1 M, "
              "AI refined, deployable, tolerant of camera "
              "imperfections.",
              size=11, italic=True, bold=True, color=ACCENT_HEX,
              font=FONT_SERIF)
    add_text(s, 0.55, 4.20, 8.95, 0.30, "TAKEAWAY",
              size=9, bold=True, font=FONT_MONO, color=ACCENT_HEX)
    add_text(s, 0.55, 4.50, 8.95, 0.55,
              "Carefully composed light weight modules can push the "
              "accuracy versus size Pareto front into a region the "
              "published literature has not occupied.",
              size=11, color=INK_HEX, font=FONT_SERIF)
    return s


def build_future_work(prs):
    s = make_cream_slide(prs,
        section_label="§ 14 FUTURE WORK",
        title_main="Open items,",
        title_accent="next steps.")
    cards = [
        ("01  ·  HARDWARE SUPPORT",
         "Acquire access to A100 class compute for full Scene Flow "
         "pre training (about 35,000 pairs)."),
        ("02  ·  FULL BENCHMARKS",
         "Fine tune and report on KITTI 2015, Middlebury 2014, ETH3D.  "
         "Target competitive D1-all and bad-3 metrics."),
        ("03  ·  INTERNATIONAL JOURNAL",
         "Submit a polished version (with the broader benchmarks) to "
         "an international IEEE / Springer journal."),
        ("04  ·  EDGE DEPLOYMENT",
         "TensorRT INT8 conversion on Jetson Orin Nano.  End to end "
         "ROS 2 node for downstream robotics integration."),
    ]
    cols_x = [0.55, 5.05]
    rows_y = [1.40, 3.10]
    cw, ch = 4.40, 1.55
    for i, (t, b) in enumerate(cards):
        x = cols_x[i % 2]
        y = rows_y[i // 2]
        add_card(s, x, y, cw, ch)
        add_text(s, x + 0.18, y + 0.12, cw - 0.36, 0.25, t,
                  size=9, bold=True, font=FONT_MONO, color=ACCENT_HEX)
        add_text(s, x + 0.18, y + 0.42, cw - 0.36, ch - 0.50, b,
                  size=10, color=INK_HEX, font=FONT_SERIF)
    return s


# --------------------------------------------------------------------------
# Driver
# --------------------------------------------------------------------------

def renumber_footers(prs):
    for i, s in enumerate(prs.slides):
        for sh in s.shapes:
            if not sh.has_text_frame: continue
            t = sh.text_frame.text.strip()
            if not (t.isdigit() and 1 <= int(t) <= 99): continue
            if sh.left is None or sh.left < Emu(8 * 914400): continue
            if sh.top is None or sh.top < Emu(5.0 * 914400): continue
            paras = list(sh.text_frame.paragraphs)
            if not paras or not paras[0].runs: continue
            paras[0].runs[0].text = f"{i+1:02d}"
            for r in paras[0].runs[1:]: r.text = ""
            break


def main():
    prs = Presentation(str(ORIG))
    print(f"opened {ORIG.name}: {len(prs.slides)} slides")

    edit_title(prs)

    builders = [
        build_intro_1, build_intro_1_1, build_intro_1_2,             # 3-5
        build_problem,                                                # 6
        build_objectives,                                             # 7
        build_literature, build_research_gap,                         # 8-9
        build_proposed,                                               # 10
        build_methodology,                                            # 11
        build_implementation,                                         # 12
        build_arch_overview, build_arch_stack_a,
        build_arch_stack_b, build_arch_stack_c,                       # 13-16
        build_working_principle,                                      # 17
        build_results_divider,                                        # 18
        build_results_sf, build_results_real,                         # 19-20
        build_results_video, build_results_3d,                        # 21-22
        build_objectives_answered,                                    # 23
        build_discussion, build_challenges, build_impact,             # 24-26
        build_conclusion, build_future_work,                          # 27-28
    ]

    new_slides = []
    for fn in builders:
        new_slides.append(fn(prs))
    print(f"appended {len(new_slides)} new slides")

    qa = prs.slides[T_QA]
    target = 2
    for sl in new_slides:
        move_slide(prs, sl, target)
        target += 1
    move_slide(prs, qa, target)

    build_outline(prs)
    print("rewrote outline")

    keep_ids = {prs.slides[0].slide_id, prs.slides[1].slide_id, qa.slide_id}
    keep_ids.update(s.slide_id for s in new_slides)
    xml = prs.slides._sldIdLst
    for el in list(xml):
        if int(el.attrib["id"]) not in keep_ids:
            xml.remove(el)
    print(f"pruned to {len(prs.slides)} slides")

    renumber_footers(prs)
    prs.save(str(OUT))
    print(f"\nsaved {OUT.name}")
    print(f"  {OUT.stat().st_size/1e6:.1f} MB  ·  {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
